#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
EV ChargeOps v2 -Sprint 01: Pesquisa e Documentação
Gera versões expandidas do relatório PDF e apresentação PPT com anexos:
  Anexo A -Pesquisa de Mercado (soluções reais Brasil e internacional)
  Anexo B -Contexto Regulatório Ampliado (produção e consumo de energia)
  Anexo C -Arquitetura Detalhada da Solução (componentes e hardware)
  Anexo D -Decisões Técnicas (Q&A com prós, contras e recomendação)

Prazo Sprint 01: 21/06/2026
"""

import os
import sys

script_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, script_dir)

from ev_chargeops import (
    EVChargeOps, gerar_graficos,
    gerar_relatorio_pdf, gerar_apresentacao_ppt
)


# ============================================================================
# ANEXOS DO RELATÓRIO PDF (v2)
# ============================================================================

def _tabela_header(pdf, colunas, larguras):
    """Cabeçalho de tabela com fundo azul."""
    pdf.set_font('Helvetica', 'B', 8)
    pdf.set_fill_color(29, 53, 87)
    pdf.set_text_color(255, 255, 255)
    for texto, w in zip(colunas, larguras):
        pdf.cell(w, 7, texto, 1, 0, 'C', True)
    pdf.ln()
    pdf.set_text_color(40, 40, 40)


def _tabela_linha(pdf, valores, larguras, fill=False):
    """Linha de tabela com alternância de cor."""
    pdf.set_font('Helvetica', '', 7)
    if fill:
        pdf.set_fill_color(240, 245, 250)
    else:
        pdf.set_fill_color(255, 255, 255)
    for texto, w in zip(valores, larguras):
        pdf.cell(w, 6, texto, 1, 0, 'L', True)
    pdf.ln()


def adicionar_anexos_pdf(pdf):
    """Adiciona Anexos A-D ao relatório PDF (Sprint 01)."""

    # ================================================================
    # PÁGINA SEPARADORA
    # ================================================================
    pdf.add_page()
    pdf.ln(40)
    pdf.set_font('Helvetica', 'B', 28)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 15, 'ANEXOS', 0, 1, 'C')
    pdf.set_font('Helvetica', '', 14)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 10, 'Sprint 01 -Pesquisa e Documentação', 0, 1, 'C')
    pdf.cell(0, 8, 'Prazo: 21/06/2026', 0, 1, 'C')
    pdf.ln(5)
    pdf.set_draw_color(230, 57, 70)
    pdf.line(60, pdf.get_y(), 150, pdf.get_y())

    # ================================================================
    # ANEXO A -PESQUISA DE MERCADO
    # ================================================================
    pdf.add_page()
    pdf.titulo_secao('A', 'Pesquisa de Mercado -Soluções Existentes')

    pdf.subtitulo('A.1 Soluções no Brasil')
    pdf.corpo(
        'O mercado brasileiro de infraestrutura de recarga para veículos '
        'elétricos cresce a uma taxa de 40-50% ao ano, impulsionado pela '
        'chegada de novos modelos e pela regulação favorável da ANEEL. '
        'Diversas empresas já oferecem soluções comerciais:'
    )

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'Voltbras (Porto Alegre, 2017)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Maior plataforma de gestão de recarga do Brasil. Modelo SaaS B2B '
        'com integração OCPP a mais de 10 fabricantes de carregadores. '
        'Oferece gestão de hubs públicos e privados, pagamento via app e '
        'RFID, dashboard de analytics e relatórios. Clientes incluem '
        'shoppings, condomínios e frotas corporativas. Parceria com BMW, '
        'Volvo e BYD para integração de dados do veículo.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'WEG WEMOB (Jaraguá do Sul)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Fabricante brasileiro de carregadores EV com a linha WEMOB: '
        'modelos de 7,4 kW (AC monofásico), 22 kW (AC trifásico) e '
        '60-180 kW (DC fast charging). Protocolo OCPP 1.6J nativo, '
        'autenticação RFID, app de gestão. Vantagem competitiva: '
        'fabricação nacional com suporte técnico local e integração '
        'com inversores solares WEG.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'Zletric (São Paulo)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Plataforma SaaS focada exclusivamente em condomínios. Oferece '
        'rateio individualizado por unidade habitacional, integração com '
        'administradoras de condomínio e relatórios mensais automatizados. '
        'Modelo de negócio: assinatura mensal por carregador gerenciado.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'EDP Box (Grupo EDP) e Tupinambá Energia', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'EDP Box: programa de recarga residencial e condominial com '
        'carregador + instalação + app, integrado à conta de energia EDP. '
        'Tupinambá Energia: rede de mais de 500 eletropostos públicos no '
        'Brasil com app de localização, reserva e pagamento digital.')
    pdf.ln(3)

    pdf.subtitulo('A.2 Soluções Internacionais')

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'ChargePoint (EUA)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Maior rede aberta de recarga do mundo com mais de 200 mil pontos. '
        'Plataforma cloud completa: gestão de frota, analytics preditivo, '
        'pagamento integrado. Hardware próprio (CP6000, Express Plus) e '
        'software como serviço. IPO na NYSE em 2021. Referência em '
        'escalabilidade e gestão enterprise.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'Wallbox (Espanha)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Destaque para o Quasar 2: primeiro carregador residencial com '
        'V2G bidirecional -o veículo pode devolver energia à rede ou '
        'à residência. Pulsar Plus para uso doméstico. App com gestão '
        'de energia solar integrada. Presente em 107 países.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'EVBox (Holanda) e Virta (Finlândia)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'EVBox (grupo ENGIE): hardware (Elvi, Iqon, Troniq Modular) + '
        'plataforma Everon para operadores. OCPP nativo, foco enterprise. '
        'Virta: plataforma white-label API-first para operadores de recarga, '
        'OCPP 2.0.1, roaming europeu via Hubject/GIREVE.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'Tesla Supercharger (EUA/Global)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Maior rede proprietária do mundo com mais de 50 mil pontos. '
        'Conector NACS adotado como padrão SAE J3400 na América do Norte. '
        'Abertura da rede para outros fabricantes via Magic Dock (CCS). '
        'Potências de 72 kW a 250 kW (V3) e 350 kW (V4).')
    pdf.ln(3)

    pdf.subtitulo('A.3 Tabela Comparativa')

    colunas = ['Empresa', 'País', 'Foco', 'Protocolo', 'V2G', 'Destaque']
    larguras = [28, 18, 28, 22, 12, 82]

    _tabela_header(pdf, colunas, larguras)
    dados = [
        ['Voltbras', 'Brasil', 'B2B / Hubs', 'OCPP 1.6', 'Não', 'Maior plataforma BR, +10 fabricantes'],
        ['WEG WEMOB', 'Brasil', 'Hardware', 'OCPP 1.6', 'Não', 'Fabricação nacional, suporte local'],
        ['Zletric', 'Brasil', 'Condomínios', 'OCPP 1.6', 'Não', 'Rateio condominial especializado'],
        ['ChargePoint', 'EUA', 'Rede aberta', 'OCPP 2.0', 'Não', '200k+ pontos, plataforma cloud'],
        ['Wallbox', 'Espanha', 'Residencial', 'OCPP 1.6', 'Sim', 'Quasar V2G bidirecional'],
        ['EVBox', 'Holanda', 'Enterprise', 'OCPP 2.0', 'Não', 'Grupo ENGIE, plataforma Everon'],
        ['Virta', 'Finlândia', 'White-label', 'OCPP 2.0', 'Não', 'API-first, roaming europeu'],
        ['EV ChargeOps', 'Brasil', 'Condomínios', 'OCPP 1.6', 'Roadmap', 'IA integrada, Síndico Virtual'],
    ]
    for i, row in enumerate(dados):
        _tabela_linha(pdf, row, larguras, fill=(i % 2 == 0))

    pdf.ln(3)
    pdf.corpo(
        'O EV ChargeOps se diferencia das soluções existentes pela integração '
        'nativa de um Motor de IA com quatro dimensões (interpretação, '
        'preditividade, precificação e conversação), pelo agente Síndico '
        'Virtual que traduz dados técnicos em linguagem gerencial, e pelo '
        'foco específico no problema de rateio condominial justo.'
    )

    # ================================================================
    # ANEXO B -CONTEXTO REGULATÓRIO AMPLIADO
    # ================================================================
    pdf.add_page()
    pdf.titulo_secao('B', 'Contexto Regulatório Ampliado')

    pdf.subtitulo('B.1 Estrutura Regulatória do Setor Elétrico Brasileiro')
    pdf.corpo(
        'O setor elétrico brasileiro é regulado por um conjunto de '
        'instituições com funções distintas:'
    )
    pdf.item_lista(
        'ANEEL (Agência Nacional de Energia Elétrica): regulação e '
        'fiscalização da geração, transmissão, distribuição e '
        'comercialização de energia elétrica.')
    pdf.item_lista(
        'ONS (Operador Nacional do Sistema Elétrico): coordenação '
        'e controle da operação do Sistema Interligado Nacional (SIN).')
    pdf.item_lista(
        'CCEE (Câmara de Comercialização de Energia Elétrica): '
        'contabilização e liquidação do mercado de energia.')
    pdf.item_lista(
        'EPE (Empresa de Pesquisa Energética): planejamento '
        'energético de longo prazo, incluindo projeções de demanda EV.')

    pdf.subtitulo('B.2 Regulamentações Específicas para Recarga de EVs')

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'ANEEL Resolução Normativa 1.000/2021', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Consolida as regras de distribuição de energia elétrica. '
        'Artigos 311 a 318 tratam especificamente da recarga de veículos '
        'elétricos, classificando-a como atividade acessória não regulada, '
        'o que permite precificação livre pelo prestador do serviço. '
        'Esta resolução é a base legal que viabiliza o modelo de '
        'tarifação dinâmica do EV ChargeOps.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'Lei 14.300/2022 -Marco Legal da Geração Distribuída', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Regulamenta a micro e minigeração distribuída no Brasil. Permite '
        'que prosumidores (consumidor + produtor) gerem própria energia '
        'via painéis solares ou eólica e injetem excedente na rede '
        '(net metering). Relevância para o EV ChargeOps: condomínios '
        'com geração solar podem abastecer os carregadores com energia '
        'própria, reduzindo custos e pegada de carbono.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'ABNT NBR IEC 61851-1:2023 e INMETRO', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'ABNT NBR IEC 61851 define os requisitos de segurança para '
        'sistemas de recarga condutiva: modos de carga 1 a 4, proteções '
        'elétricas, comunicação piloto (Control Pilot). INMETRO Portaria '
        '111/2023 estabelece certificação compulsória para equipamentos '
        'de recarga comercializados no Brasil.')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'PNE 2050 -Plano Nacional de Energia', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Elaborado pela EPE, projeta 4 milhões de veículos elétricos no '
        'Brasil até 2030 e 33 milhões até 2050. A demanda adicional de '
        'energia é estimada em 24 TWh/ano até 2035, equivalente a '
        'aproximadamente 3% do consumo elétrico nacional atual.')
    pdf.ln(3)

    pdf.subtitulo('B.3 Produção e Consumo de Energia no Brasil')
    pdf.corpo(
        'A matriz elétrica brasileira é predominantemente renovável: '
        'hidráulica (63%), eólica (11%), solar (4%), biomassa (5%) e '
        'outras fontes (17%). A capacidade instalada total ultrapassa '
        '195 GW. A tarifa média residencial é de R$ 0,65/kWh com impostos.'
    )
    pdf.corpo(
        'O sistema de bandeiras tarifárias (verde, amarela, vermelha) '
        'reflete o custo de geração de energia. A tarifa branca, opcional '
        'para consumidores residenciais com medição horária, aplica '
        'valores diferenciados por período do dia -modelo análogo ao '
        'utilizado pelo EV ChargeOps para precificação de recargas.'
    )
    pdf.corpo(
        'O horário de ponta (18h-21h, definido pela distribuidora) concentra '
        'a maior demanda e o maior custo de geração. Carregar veículos '
        'fora deste horário reduz o impacto na rede e o custo para o '
        'consumidor -princípio central da precificação dinâmica do sistema.'
    )

    pdf.subtitulo('B.4 Regulação Internacional de Referência')
    pdf.item_lista(
        'União Europeia -AFIR (2023): obriga postos de recarga rápida '
        'a cada 60 km em rodovias TEN-T, com mínimo de 150 kW até 2025 '
        'e 350 kW até 2030. Pagamento por cartão obrigatório.')
    pdf.item_lista(
        'EUA -NEVI Program: investimento de US$ 7,5 bilhões para '
        'instalação de 500 mil carregadores até 2030. Exige OCPP, '
        'uptime mínimo de 97% e conectores CCS.')
    pdf.item_lista(
        'Reino Unido -Smart Charge Points Regulations 2021: exige '
        'que carregadores residenciais suportem demand response e '
        'V2G. Carga padrão fora do pico (off-peak default).')
    pdf.item_lista(
        'China -Padrões GB/T: maior mercado EV do mundo (>8 milhões '
        'vendidos em 2023). Padronização própria de conectores e '
        'protocolo de comunicação.')

    pdf.subtitulo('B.5 Perspectivas V2G (Vehicle-to-Grid)')
    pdf.corpo(
        'O V2G (Vehicle-to-Grid) permite utilizar a bateria do veículo '
        'elétrico como armazenamento distribuído, devolvendo energia à '
        'rede nos horários de ponta. A ANEEL estuda a regulamentação '
        'via Consulta Pública 020/2024. Aplicações potenciais: peak '
        'shaving, reserva de emergência e arbitragem tarifária. '
        'O EV ChargeOps inclui V2G no roadmap de desenvolvimento, '
        'preparando a arquitetura para suporte futuro ao protocolo '
        'ISO 15118 (comunicação veículo-carregador bidirecional).'
    )

    # ================================================================
    # ANEXO C -ARQUITETURA DETALHADA DA SOLUÇÃO
    # ================================================================
    pdf.add_page()
    pdf.titulo_secao('C', 'Arquitetura Detalhada da Solução')

    pdf.subtitulo('C.1 Visão Geral -Três Camadas')
    pdf.corpo(
        'O EV ChargeOps adota uma arquitetura híbrida de três camadas '
        'conforme especificado no playbook do Challenge. Cada camada tem '
        'responsabilidades bem definidas e interfaces claras, permitindo '
        'evolução independente dos componentes.'
    )

    pdf.subtitulo('C.2 Camada Física -Hardware GoodWe HCA G2')
    pdf.corpo(
        'O hardware de referência é a linha GoodWe HCA G2, projetada '
        'para instalação em condomínios e estacionamentos comerciais. '
        'Especificações técnicas por modelo:'
    )

    colunas_hw = ['Modelo', 'Potência', 'Fase', 'Conector', 'Proteções', 'IP']
    larguras_hw = [35, 20, 22, 25, 55, 13]
    _tabela_header(pdf, colunas_hw, larguras_hw)
    hw_dados = [
        ['GW7K-HCA-20', '7 kW', 'Monofásico', 'AC Tipo 2', 'RCD A, OVP, UVP, OCP, OTP', 'IP65'],
        ['GW11K-HCA-20', '11 kW', 'Trifásico', 'AC Tipo 2', 'RCD A, OVP, UVP, OCP, OTP', 'IP65'],
        ['GW22K-HCA-20', '22 kW', 'Trifásico', 'AC Tipo 2', 'RCD A, OVP, UVP, OCP, OTP', 'IP65'],
    ]
    for i, row in enumerate(hw_dados):
        _tabela_linha(pdf, row, larguras_hw, fill=(i % 2 == 0))

    pdf.ln(3)
    pdf.corpo(
        'Características comuns a todos os modelos: autenticação RFID '
        '(ISO 14443A, até 10 cartões), conectividade RS-485 + LAN + '
        'Wi-Fi + Bluetooth, display LED com indicadores de status, '
        'e firmware atualizável remotamente via OCPP.')
    pdf.corpo(
        'Infraestrutura elétrica recomendada: quadro de distribuição '
        'dedicado com disjuntores individuais por carregador, cabeamento '
        'mínimo de 6 mm² (7 kW) a 10 mm² (22 kW), e aterramento '
        'conforme padrão TN-S da NBR 5410.')

    pdf.subtitulo('C.3 Camada de Conectividade -Protocolos')

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'OCPP 1.6J (JSON over WebSocket)', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Protocolo aberto da Open Charge Alliance para comunicação '
        'bidirecional entre o sistema central (EV ChargeOps) e os '
        'pontos de carga (carregadores GoodWe). Operações principais: '
        'Authorize (validar RFID), StartTransaction / StopTransaction '
        '(controle de sessão), MeterValues (leituras de energia), '
        'StatusNotification (estado do carregador), Heartbeat '
        '(keep-alive) e FirmwareUpdate (atualização remota).')
    pdf.ln(2)

    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(69, 123, 157)
    pdf.cell(0, 7, 'MODBUS RTU/TCP', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 5,
        'Protocolo industrial para leitura de registradores de energia '
        'em tempo real: tensão (V), corrente (A), potência ativa (W), '
        'fator de potência, frequência (Hz) e energia acumulada (kWh). '
        'Polling configurável (padrão: a cada 5 segundos).')
    pdf.ln(3)

    pdf.subtitulo('C.4 Camada Digital -Módulos de Software')
    pdf.corpo(
        'A camada digital é composta por cinco módulos interdependentes, '
        'todos implementados em Python:'
    )
    pdf.item_lista(
        'Gerenciador de Sessões: controle de estado (disponível, em uso, '
        'finalizada), vinculação unidade-carregador via RFID, log de eventos.')
    pdf.item_lista(
        'Motor de Faturamento: cálculo de consumo individual (kWh x tarifa), '
        'geração de faturas mensais, rateio condominial proporcional.')
    pdf.item_lista(
        'Motor de IA: interpretação de sessões, preditividade via média '
        'móvel, precificação dinâmica ANEEL, Síndico Virtual (NLP).')
    pdf.item_lista(
        'Integrações: GoodWe API (OCPP simulado), Open Charge Map '
        '(eletropostos públicos), Google Places (localização).')
    pdf.item_lista(
        'Visualização: 5 gráficos Matplotlib, relatório PDF (fpdf), '
        'apresentação PPT (python-pptx), interface CLI interativa.')

    pdf.subtitulo('C.5 Diagrama de Fluxo de Dados')
    pdf.set_font('Courier', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '  [Cartão RFID] --> [Carregador GoodWe HCA G2]\n'
        '                           |\n'
        '                    OCPP 1.6J / MODBUS\n'
        '                           |\n'
        '                           v\n'
        '                  [Servidor EV ChargeOps]\n'
        '                    /       |       \\\n'
        '                   v        v        v\n'
        '          [Gerenciador] [Motor IA] [Faturamento]\n'
        '               |          |           |\n'
        '               v          v           v\n'
        '          [Sessões]  [Dashboard]  [Faturas]\n'
        '                        |   |\n'
        '                        v   v\n'
        '              [Gráficos] [Síndico Virtual]')
    pdf.ln(3)
    pdf.set_font('Helvetica', '', 10)

    pdf.subtitulo('C.6 Segurança')
    pdf.item_lista(
        'Autenticação por RFID com lista branca de cartões cadastrados.')
    pdf.item_lista(
        'Sessões vinculadas exclusivamente a unidades cadastradas no sistema.')
    pdf.item_lista(
        'Logs de auditoria para todas as operações (início, fim, anomalias).')
    pdf.item_lista(
        'Comunicação OCPP sobre TLS em ambiente de produção (Security Profile 2).')
    pdf.item_lista(
        'Proteção contra uso não autorizado: carregador rejeita RFID não cadastrado.')

    # ================================================================
    # ANEXO D -DECISÕES TÉCNICAS (Q&A)
    # ================================================================
    pdf.add_page()
    pdf.titulo_secao('D', 'Decisões Técnicas -Perguntas e Respostas')
    pdf.corpo(
        'Esta seção documenta as principais decisões técnicas do projeto, '
        'apresentando argumentos a favor e contra cada alternativa avaliada, '
        'com a recomendação final adotada.'
    )

    # --- Q1 ---
    pdf.ln(2)
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 8, 'Q1: Por que Python como linguagem principal?', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Prós:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Ecossistema científico robusto (NumPy, Pandas, Scikit-learn) '
        'ideal para o Motor de IA.\n'
        '- Prototipagem rápida com tipagem dinâmica -adequado para MVP.\n'
        '- Maior comunidade de desenvolvedores do mundo (Stack Overflow 2024).\n'
        '- Bibliotecas maduras para geração de PDF, PPTX e gráficos.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(230, 57, 70)
    pdf.cell(0, 6, 'Contras:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Performance inferior a linguagens compiladas (C, Go, Rust).\n'
        '- GIL (Global Interpreter Lock) limita paralelismo verdadeiro.\n'
        '- Tipagem dinâmica pode dificultar manutenção em projetos grandes.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, 'Decisão: Python -adequado para MVP e análise de dados; '
             'migração para Go/Rust se necessário em produção.', 0, 1, '', True)
    pdf.ln(3)

    # --- Q2 ---
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 8, 'Q2: OCPP aberto ou protocolo proprietário?', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Prós OCPP:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Padrão aberto mantido pela Open Charge Alliance (+300 membros).\n'
        '- Interoperabilidade com carregadores de qualquer fabricante.\n'
        '- Exigência regulatória da ANEEL para uso compartilhado.\n'
        '- Versão 2.0.1 com suporte a certificados e smart charging.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(230, 57, 70)
    pdf.cell(0, 6, 'Prós protocolo proprietário:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Performance otimizada para hardware específico.\n'
        '- Acesso a features exclusivas do fabricante.\n'
        '- Menor complexidade de implementação inicial.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, 'Decisão: OCPP 1.6J -compliance regulatório e '
             'interoperabilidade; upgrade planejado para 2.0.1.', 0, 1, '', True)
    pdf.ln(3)

    # --- Q3 ---
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 8, 'Q3: Tarifa dinâmica ou tarifa fixa?', 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(40, 40, 40)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Prós tarifa dinâmica:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Incentiva recarga fora do horário de ponta, aliviando a rede.\n'
        '- Economia real para moradores que carregam à noite.\n'
        '- Alinhada com a tarifa branca da ANEEL.\n'
        '- Permite arbitragem de custo pelo condomínio.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(230, 57, 70)
    pdf.cell(0, 6, 'Prós tarifa fixa:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Simplicidade de implementação e compreensão.\n'
        '- Previsibilidade total para o morador.\n'
        '- Sem necessidade de lógica de horário no sistema.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, 'Decisão: Tarifa dinâmica -incentivo econômico real '
             'e alinhamento com regulação ANEEL.', 0, 1, '', True)
    pdf.ln(3)

    # --- Q4 ---
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 8, 'Q4: Autenticação por RFID ou por aplicativo?', 0, 1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Prós RFID:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Operação simples: aproximar cartão e carregar.\n'
        '- Não requer smartphone, conexão de dados ou conta digital.\n'
        '- Custo baixo (R$ 5-10 por cartão ISO 14443A).\n'
        '- Hardware já incluso nos carregadores GoodWe (2 cartões).')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(230, 57, 70)
    pdf.cell(0, 6, 'Prós aplicativo:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Autenticação via QR code, NFC ou biometria.\n'
        '- Pagamento digital integrado (PIX, cartão).\n'
        '- Visualização de consumo e histórico em tempo real.\n'
        '- Reserva antecipada de carregador.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, 'Decisão: RFID como primário (base GoodWe instalada); '
             'app como evolução no roadmap.', 0, 1, '', True)
    pdf.ln(3)

    # --- Q5 ---
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 8, 'Q5: Armazenamento em memória ou banco de dados?', 0, 1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Prós em memória (dataclasses):', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Sem dependências externas (zero configuração).\n'
        '- Velocidade máxima de leitura/escrita.\n'
        '- Ideal para prototipagem e demonstração.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(230, 57, 70)
    pdf.cell(0, 6, 'Prós banco de dados:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Persistência de dados entre reinicializações.\n'
        '- Consultas SQL para relatórios complexos.\n'
        '- Backup e recuperação de dados.\n'
        '- Suporte a múltiplos usuários simultâneos.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, 'Decisão: Em memória para MVP -migração para '
             'SQLite/PostgreSQL na versão de produção.', 0, 1, '', True)
    pdf.ln(3)

    # --- Q6 ---
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 8, 'Q6: Infraestrutura on-premise ou cloud?', 0, 1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Prós on-premise:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Controle total dos dados (LGPD compliance).\n'
        '- Sem custos mensais de cloud.\n'
        '- Latência local mínima para OCPP.\n'
        '- Operação offline garantida.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(230, 57, 70)
    pdf.cell(0, 6, 'Prós cloud:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(40, 40, 40)
    pdf.multi_cell(0, 4,
        '- Escalabilidade elástica para múltiplos condomínios.\n'
        '- Backup automático e alta disponibilidade.\n'
        '- Acesso remoto ao dashboard de qualquer lugar.\n'
        '- Integração facilitada com APIs externas.')
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, 'Decisão: On-premise para MVP -migração para '
             'cloud (AWS/GCP) na versão multi-condomínio.', 0, 1, '', True)


# ============================================================================
# SLIDES ADICIONAIS DA APRESENTAÇÃO PPT (v2)
# ============================================================================

def adicionar_slides_ppt(prs, plataforma, pasta_saida):
    """Adiciona slides Sprint 01 à apresentação PPT."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    VERMELHO = RGBColor(230, 57, 70)
    AZUL_ESCURO = RGBColor(29, 53, 87)
    AZUL_MEDIO = RGBColor(69, 123, 157)
    CINZA = RGBColor(100, 100, 100)
    BRANCO = RGBColor(255, 255, 255)
    PRETO = RGBColor(40, 40, 40)
    VERDE = RGBColor(42, 157, 143)

    def add_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def titulo(slide, texto, sub=None):
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.5),
                                       Inches(11.5), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = texto
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = AZUL_ESCURO
        if sub:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(18)
            p2.font.color.rgb = CINZA
            p2.space_before = Pt(8)
        sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3), Pt(4))
        sh.fill.solid()
        sh.fill.fore_color.rgb = VERMELHO
        sh.line.fill.background()

    def caixa(slide, texto, x, y, w, h, cor=AZUL_ESCURO):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = texto
        p.font.size = Pt(14)
        p.font.color.rgb = BRANCO
        p.font.bold = True

    def card(slide, titulo_txt, desc, x, y, w, h, cor_borda=AZUL_ESCURO):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
        sh.line.color.rgb = cor_borda
        sh.line.width = Pt(2)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo_txt
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = AZUL_ESCURO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = PRETO
        p2.space_before = Pt(8)

    # =========================================================
    # SLIDE SEPARADOR: SPRINT 01
    # =========================================================
    slide = add_slide()
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_ESCURO
    bg.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1), Inches(2),
                                   Inches(11), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sprint 01"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Pesquisa e Documentação"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(5), Inches(5), Pt(3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = VERMELHO
    sh.line.fill.background()
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(5.5),
                                    Inches(11), Inches(1))
    tf2 = tx2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "Prazo: 21/06/2026"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(180, 180, 180)
    p3.alignment = PP_ALIGN.CENTER

    # =========================================================
    # SLIDE: PESQUISA DE MERCADO -BRASIL
    # =========================================================
    slide = add_slide()
    titulo(slide, "Pesquisa de Mercado -Brasil",
           "Soluções comerciais de recarga para condomínios e frotas")

    empresas_br = [
        ("Voltbras", "Maior plataforma BR\nSaaS B2B, OCPP, +10 fabricantes\nShoppings, condomínios, frotas",
         0.5, AZUL_ESCURO),
        ("WEG WEMOB", "Fabricante nacional\n7-180 kW, OCPP nativo\nIntegração inversores solares",
         3.7, AZUL_MEDIO),
        ("Zletric", "Foco em condomínios\nRateio individualizado\nIntegração administradoras",
         6.9, VERMELHO),
        ("EDP / Tupinambá", "EDP Box: residencial\nTupinambá: +500 eletropostos\nApp + pagamento digital",
         10.1, VERDE),
    ]
    for nome, desc, x, cor in empresas_br:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5),
            Inches(2.8), Inches(4.0))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = nome
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)

    # =========================================================
    # SLIDE: PESQUISA DE MERCADO -INTERNACIONAL
    # =========================================================
    slide = add_slide()
    titulo(slide, "Pesquisa de Mercado -Internacional",
           "Referências globais em infraestrutura de recarga")

    empresas_int = [
        ("ChargePoint", "EUA | 200k+ pontos\nMaior rede aberta do mundo\nCloud + analytics + pagamento",
         0.5, AZUL_ESCURO),
        ("Wallbox", "Espanha | 107 países\nQuasar 2: V2G bidirecional\nIntegração energia solar",
         3.7, AZUL_MEDIO),
        ("EVBox / Virta", "Holanda / Finlândia\nHardware + plataforma\nOCPP 2.0, roaming EU",
         6.9, RGBColor(244, 162, 97)),
        ("Tesla", "Global | 50k+ pontos\nNACS = padrão SAE J3400\n72-350 kW, Magic Dock",
         10.1, RGBColor(38, 70, 83)),
    ]
    for nome, desc, x, cor in empresas_int:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5),
            Inches(2.8), Inches(4.0))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = nome
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)

    # =========================================================
    # SLIDE: CONTEXTO REGULATÓRIO -BRASIL
    # =========================================================
    slide = add_slide()
    titulo(slide, "Contexto Regulatório -Brasil",
           "Regulamentação do setor de energia e recarga de EVs")

    regs = [
        ("ANEEL RN 1.000/2021",
         "Recarga como atividade não regulada.\n"
         "Precificação livre pelo prestador.\n"
         "Base legal do EV ChargeOps."),
        ("Lei 14.300/2022",
         "Marco da Geração Distribuída.\n"
         "Prosumidores: gerar + consumir.\n"
         "Solar em condomínios + EVs."),
        ("NBR IEC 61851 / INMETRO",
         "Requisitos de segurança para\n"
         "sistemas de recarga condutiva.\n"
         "Certificação compulsória."),
        ("PNE 2050 (EPE)",
         "4M EVs até 2030, 33M até 2050.\n"
         "Demanda: +24 TWh/ano até 2035.\n"
         "3% do consumo elétrico nacional."),
    ]
    for i, (nome, desc) in enumerate(regs):
        x = 0.5 + i * 3.15
        card(slide, nome, desc, x, 2.3, 2.9, 3.5,
             [AZUL_ESCURO, AZUL_MEDIO, VERMELHO, VERDE][i])

    # =========================================================
    # SLIDE: REGULAÇÃO INTERNACIONAL + V2G
    # =========================================================
    slide = add_slide()
    titulo(slide, "Regulação Internacional e V2G",
           "Referências globais e perspectivas de Vehicle-to-Grid")

    regs_int = [
        ("EU -AFIR 2023", "Postos a cada 60 km\n150 kW (2025), 350 kW (2030)\nPagamento por cartão", AZUL_ESCURO),
        ("EUA -NEVI", "US$ 7,5 bi investidos\n500k chargers até 2030\nOCPP + CCS obrigatório", AZUL_MEDIO),
        ("UK -Smart Charge", "Demand response obrigatório\nV2G habilitado\nCarga off-peak padrão", VERMELHO),
    ]
    for i, (nome, desc, cor) in enumerate(regs_int):
        x = 0.8 + i * 4.1
        caixa(slide, f"{nome}\n\n{desc}", x, 2.3, 3.8, 2.5, cor)

    # V2G box
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.2),
        Inches(10), Inches(1.8))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
    sh.line.color.rgb = VERDE
    sh.line.width = Pt(3)
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "V2G (Vehicle-to-Grid) -Perspectiva Futura"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AZUL_ESCURO
    p2 = tf.add_paragraph()
    p2.text = ("Bateria do EV como armazenamento distribuído. "
               "ANEEL CP 020/2024 em análise. "
               "Aplicações: peak shaving, reserva de emergência, arbitragem tarifária. "
               "EV ChargeOps: V2G no roadmap via ISO 15118.")
    p2.font.size = Pt(13)
    p2.font.color.rgb = PRETO
    p2.space_before = Pt(8)

    # =========================================================
    # SLIDE: ARQUITETURA DETALHADA
    # =========================================================
    slide = add_slide()
    titulo(slide, "Arquitetura Detalhada da Solução",
           "Componentes, protocolos e fluxo de dados")

    # 3 camadas visuais
    camadas = [
        ("CAMADA DIGITAL", "Motor de Sessões | Motor de IA | Faturamento | Integrações | Dashboard",
         AZUL_ESCURO, 2.2),
        ("CAMADA DE CONECTIVIDADE", "OCPP 1.6J (WebSocket) | MODBUS RTU/TCP | Wi-Fi + LAN + RS-485",
         AZUL_MEDIO, 3.7),
        ("CAMADA FÍSICA", "GoodWe HCA G2: 7kW | 11kW | 22kW | RFID ISO 14443A | IP65",
         RGBColor(80, 80, 80), 5.2),
    ]
    for nome, desc, cor, y in camadas:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y),
            Inches(10), Inches(1.3))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = nome
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER

    # Setas entre camadas
    for y in [3.55, 5.05]:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(y),
            Inches(0.7), Inches(0.2))
        sh.fill.solid()
        sh.fill.fore_color.rgb = VERMELHO
        sh.line.fill.background()

    # =========================================================
    # SLIDE: HARDWARE E PROTOCOLOS
    # =========================================================
    slide = add_slide()
    titulo(slide, "Hardware GoodWe HCA G2 e Protocolos",
           "Especificações técnicas do hardware de referência")

    modelos = [
        ("GW7K-HCA-20", "7 kW\nMonofásico\nAC Tipo 2\nResidencial"),
        ("GW11K-HCA-20", "11 kW\nTrifásico\nAC Tipo 2\nCondomínio"),
        ("GW22K-HCA-20", "22 kW\nTrifásico\nAC Tipo 2\nComercial"),
    ]
    for i, (modelo, specs) in enumerate(modelos):
        x = 0.8 + i * 4.1
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3),
            Inches(3.8), Inches(2.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = [AZUL_ESCURO, AZUL_MEDIO, VERMELHO][i]
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = modelo
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = specs
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)

    # Características comuns
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2),
        Inches(11.5), Inches(1.8))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
    sh.line.color.rgb = RGBColor(200, 200, 200)
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Características Comuns: "
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = AZUL_ESCURO
    run2 = p.add_run()
    run2.text = ("RFID ISO 14443A (até 10 cartões) | RS-485 + LAN + Wi-Fi + Bluetooth | "
                 "IP65 (uso externo) | Proteções: RCD tipo A, OVP, UVP, OCP, OTP | "
                 "Display LED | Firmware remoto via OCPP | OCPP 1.6J nativo")
    run2.font.size = Pt(12)
    run2.font.color.rgb = PRETO

    # =========================================================
    # SLIDE: DECISÕES TÉCNICAS (PARTE 1)
    # =========================================================
    slide = add_slide()
    titulo(slide, "Decisões Técnicas (1/2)",
           "Linguagem, protocolo e modelo de precificação")

    decisoes_1 = [
        ("Q1: Por que Python?",
         "Prós: ecossistema IA, prototipagem\n"
         "Contras: performance, GIL\n"
         "Decisão: Python para MVP",
         AZUL_ESCURO),
        ("Q2: OCPP ou proprietário?",
         "Prós OCPP: aberto, interoperável\n"
         "Prós prop.: performance, features\n"
         "Decisão: OCPP 1.6J (regulação)",
         AZUL_MEDIO),
        ("Q3: Tarifa dinâmica ou fixa?",
         "Prós din.: incentivo fora ponta\n"
         "Prós fixa: simplicidade\n"
         "Decisão: Dinâmica (ANEEL)",
         VERMELHO),
    ]
    for i, (pergunta, resp, cor) in enumerate(decisoes_1):
        x = 0.5 + i * 4.2
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3),
            Inches(3.8), Inches(4.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = pergunta
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = resp
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)

    # =========================================================
    # SLIDE: DECISÕES TÉCNICAS (PARTE 2)
    # =========================================================
    slide = add_slide()
    titulo(slide, "Decisões Técnicas (2/2)",
           "Autenticação, armazenamento e infraestrutura")

    decisoes_2 = [
        ("Q4: RFID ou App?",
         "Prós RFID: simples, sem celular\n"
         "Prós App: QR, pagamento, reserva\n"
         "Decisão: RFID + App no roadmap",
         VERDE),
        ("Q5: Memória ou Banco?",
         "Prós memória: zero config, rápido\n"
         "Prós BD: persistência, SQL, backup\n"
         "Decisão: Memória → SQLite/PG",
         RGBColor(244, 162, 97)),
        ("Q6: On-premise ou Cloud?",
         "Prós local: controle, LGPD, offline\n"
         "Prós cloud: escala, backup, acesso\n"
         "Decisão: Local → Cloud (prod)",
         RGBColor(38, 70, 83)),
    ]
    for i, (pergunta, resp, cor) in enumerate(decisoes_2):
        x = 0.5 + i * 4.2
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3),
            Inches(3.8), Inches(4.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = pergunta
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = resp
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)


# ============================================================================
# EXECUÇÃO
# ============================================================================

def main():
    pasta_saida = script_dir

    print("=" * 70)
    print("  EV ChargeOps v2 -Sprint 01: Pesquisa e Documentação")
    print("=" * 70)

    print("\n[1/4] Configurando plataforma...")
    plataforma = EVChargeOps()
    plataforma.setup_demo()

    print("\n[2/4] Gerando histórico simulado (30 dias)...")
    plataforma.gerar_historico_simulado(dias=30)

    print("\n[3/4] Gerando relatório PDF v2 (seções originais + anexos)...")
    pdf = gerar_relatorio_pdf(plataforma, pasta_saida, save=False)
    adicionar_anexos_pdf(pdf)
    caminho_pdf = os.path.join(pasta_saida, 'relatorio_ev_chargeops_v2.pdf')
    pdf.output(caminho_pdf)
    print(f"  Relatório PDF v2 salvo: {caminho_pdf}")

    print("\n[4/4] Gerando apresentação PPT v2 (slides originais + Sprint 01)...")
    gerar_graficos(plataforma, pasta_saida)
    prs = gerar_apresentacao_ppt(plataforma, pasta_saida, save=False)
    adicionar_slides_ppt(prs, plataforma, pasta_saida)
    caminho_ppt = os.path.join(pasta_saida, 'apresentacao_ev_chargeops_v2.pptx')
    prs.save(caminho_ppt)
    print(f"  Apresentação PPT v2 salva: {caminho_ppt}")

    print("\n" + "=" * 70)
    print("  v2 GERADA COM SUCESSO!")
    print(f"  - {caminho_pdf}")
    print(f"  - {caminho_ppt}")
    print("=" * 70)


if __name__ == "__main__":
    main()
