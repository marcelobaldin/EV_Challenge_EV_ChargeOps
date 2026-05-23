#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
==============================================================================
EV ChargeOps - Plataforma de Gestao Compartilhada de Recarga para Condominios
==============================================================================
Challenge FIAP + GoodWe 2026
Curso: Ciencias da Computacao (Online) - FIAP
Aluno: Marcelo Bastianello Baldin | RM568746

Requisitos atendidos (Turma Online - EV ChargeOps):
  1. Estruturacao de sessoes de recarga por usuario ou unidade habitacional
  2. Processamento de consumo individual para cobranca automatica da unidade
  3. Sistema de gerenciamento inteligente da recarga com interface para o
     usuario e gestao do condominio

Pilares do EV ChargeOps:
  - Gestao de Uso          | Sessoes por usuario/unidade
  - Gerenciamento de Recarga | Regras de uso compartilhado, horarios otimizados
  - Rateio e Faturamento   | Consumo individual (kWh) -> cobranca automatica
  - Visualizacao e Interacao | Paineis gerenciais + insights de consumo

Motor de IA:
  - Interpretacao   | Decodificacao de eventos de sessao
  - Preditividade   | Previsao de demanda e expansao de infraestrutura
  - Precificacao    | Calculo dinamico de tarifas (horario ponta/fora ponta)
  - Conversacao     | Sindico Virtual (NLP) - traduz dados em orientacoes

Integracoes:
  - GoodWe API (simulada)     | Dados operacionais dos carregadores HCA G2
  - Open Charge Map API       | Mapa publico de eletropostos
  - Google Places API         | Localizacao de pontos de recarga proximos

Hardware de referencia: GoodWe Linha HCA G2
  - GW7K-HCA-20  (7 kW)
  - GW11K-HCA-20 (11 kW)
  - GW22K-HCA-20 (22 kW)
  - RFID (2 cartoes inclusos, suporta ate 10)
  - Conectividade: RS-485 + LAN + Wi-Fi + Bluetooth

Base Regulatoria:
  - ANEEL Resolucao Normativa 1.000/2021 (precificacao livre)
  - Obrigatoriedade de interoperabilidade via OCPP
==============================================================================
"""

import json
import random
import uuid
import os
import sys
from datetime import datetime, timedelta
from collections import defaultdict
from dataclasses import dataclass, field
from typing import Optional

# ============================================================================
# 1. MODELOS DE DADOS
# ============================================================================

@dataclass
class Condominio:
    """Representa um condominio com carregadores EV compartilhados."""
    id: str = field(default_factory=lambda: str(uuid.uuid4())[:8])
    nome: str = ""
    endereco: str = ""
    capacidade_total_kw: float = 0.0
    num_carregadores: int = 0
    tarifa_kwh_base: float = 0.85  # R$/kWh base ANEEL


@dataclass
class UnidadeHabitacional:
    """Unidade (apartamento/sala) vinculada ao condominio."""
    id: str = field(default_factory=lambda: str(uuid.uuid4())[:8])
    numero: str = ""
    bloco: str = ""
    proprietario: str = ""
    condominio_id: str = ""
    rfid_cards: list = field(default_factory=list)
    saldo_devedor: float = 0.0


@dataclass
class Carregador:
    """Carregador GoodWe HCA G2 instalado no condominio."""
    id: str = field(default_factory=lambda: str(uuid.uuid4())[:8])
    modelo: str = "GW11K-HCA-20"
    potencia_kw: float = 11.0
    condominio_id: str = ""
    status: str = "disponivel"  # disponivel, em_uso, manutencao, offline
    localizacao: str = ""
    protocolo: str = "OCPP 1.6J"
    conectividade: str = "RS-485 + LAN + Wi-Fi + Bluetooth"
    sessao_atual_id: Optional[str] = None


@dataclass
class SessaoRecarga:
    """Sessao individual de recarga vinculada a uma unidade."""
    id: str = field(default_factory=lambda: str(uuid.uuid4())[:8])
    unidade_id: str = ""
    carregador_id: str = ""
    inicio: Optional[datetime] = None
    fim: Optional[datetime] = None
    energia_kwh: float = 0.0
    custo_total: float = 0.0
    tarifa_aplicada: float = 0.0
    status: str = "ativa"  # ativa, finalizada, interrompida
    tipo_tarifa: str = "fora_ponta"  # ponta, intermediaria, fora_ponta
    potencia_media_kw: float = 0.0


@dataclass
class Fatura:
    """Fatura mensal de uma unidade habitacional."""
    id: str = field(default_factory=lambda: str(uuid.uuid4())[:8])
    unidade_id: str = ""
    mes_referencia: str = ""
    sessoes: list = field(default_factory=list)
    total_kwh: float = 0.0
    total_reais: float = 0.0
    status: str = "aberta"  # aberta, fechada, paga
    data_emissao: Optional[datetime] = None
    data_vencimento: Optional[datetime] = None


# ============================================================================
# 2. SIMULADOR DE API GOODWE
# ============================================================================

class GoodWeAPISimulator:
    """
    Simula a API GoodWe para o EV Charger FIAP.
    Em producao, conectaria via OCPP/MODBUS ao hardware real.
    Retorna dados no formato JSON compativel com a API GoodWe.
    """

    @staticmethod
    def get_charger_status(carregador: Carregador) -> dict:
        """Retorna status operacional do carregador (simulado)."""
        return {
            "device_id": carregador.id,
            "model": carregador.modelo,
            "status": carregador.status,
            "power": f"{carregador.potencia_kw}kW",
            "timestamp": datetime.now().isoformat(),
            "protocol": carregador.protocolo,
            "connectivity": {
                "wifi": True,
                "lan": True,
                "rs485": True,
                "bluetooth": True
            },
            "temperature_c": round(random.uniform(25, 45), 1),
            "firmware_version": "2.1.4",
            "uptime_hours": random.randint(100, 5000)
        }

    @staticmethod
    def start_charging(carregador: Carregador, rfid_tag: str) -> dict:
        """Inicia sessao de recarga via RFID (simulado)."""
        return {
            "action": "RemoteStartTransaction",
            "device_id": carregador.id,
            "rfid_tag": rfid_tag,
            "status": "Accepted",
            "transaction_id": str(uuid.uuid4())[:8],
            "timestamp": datetime.now().isoformat(),
            "max_power_kw": carregador.potencia_kw
        }

    @staticmethod
    def stop_charging(carregador: Carregador, energia_kwh: float) -> dict:
        """Finaliza sessao de recarga (simulado)."""
        return {
            "action": "RemoteStopTransaction",
            "device_id": carregador.id,
            "status": "Accepted",
            "energy_delivered_kwh": round(energia_kwh, 2),
            "timestamp": datetime.now().isoformat()
        }

    @staticmethod
    def get_meter_values(carregador: Carregador) -> dict:
        """Retorna leituras do medidor MODBUS (simulado)."""
        potencia_atual = random.uniform(0, carregador.potencia_kw)
        return {
            "device_id": carregador.id,
            "meter_values": {
                "voltage_v": round(random.uniform(218, 232), 1),
                "current_a": round(potencia_atual * 1000 / 220, 1),
                "power_kw": round(potencia_atual, 2),
                "energy_kwh": round(random.uniform(0, 50), 2),
                "frequency_hz": round(random.uniform(59.8, 60.2), 1),
                "power_factor": round(random.uniform(0.95, 1.0), 3)
            },
            "timestamp": datetime.now().isoformat()
        }


# ============================================================================
# 3. INTEGRACOES EXTERNAS (STUBS)
# ============================================================================

class OpenChargeMapIntegration:
    """
    Integracao com Open Charge Map API.
    Retorna eletropostos publicos proximos ao condominio.
    Em producao: requests.get('https://api.openchargemap.io/v3/poi/...')
    """

    @staticmethod
    def buscar_eletropostos(latitude: float, longitude: float, raio_km: int = 5) -> list:
        """Busca eletropostos proximos (dados simulados para demonstracao)."""
        eletropostos = [
            {
                "id": "OCM-001",
                "nome": "Eletroposto Shell Recharge - Av. Paulista",
                "endereco": "Av. Paulista, 1500 - Bela Vista, Sao Paulo",
                "latitude": -23.5615,
                "longitude": -46.6559,
                "distancia_km": 1.2,
                "conectores": ["Tipo 2", "CCS2"],
                "potencia_max_kw": 22,
                "status": "Disponivel",
                "preco_kwh": 1.20
            },
            {
                "id": "OCM-002",
                "nome": "EDP Box - Shopping Ibirapuera",
                "endereco": "Av. Ibirapuera, 3103 - Moema, Sao Paulo",
                "latitude": -23.6105,
                "longitude": -46.6658,
                "distancia_km": 2.8,
                "conectores": ["Tipo 2"],
                "potencia_max_kw": 7.4,
                "status": "Disponivel",
                "preco_kwh": 0.95
            },
            {
                "id": "OCM-003",
                "nome": "Tupinamba Energia - Rua Oscar Freire",
                "endereco": "Rua Oscar Freire, 900 - Jardins, Sao Paulo",
                "latitude": -23.5655,
                "longitude": -46.6710,
                "distancia_km": 3.5,
                "conectores": ["Tipo 2", "CCS2", "CHAdeMO"],
                "potencia_max_kw": 50,
                "status": "Em uso",
                "preco_kwh": 1.50
            }
        ]
        return eletropostos


class GooglePlacesIntegration:
    """
    Integracao com Google Places API (evChargeOptions).
    Cruza infraestrutura privada com mapa publico.
    Em producao: requests.get('https://maps.googleapis.com/maps/api/place/...')
    """

    @staticmethod
    def buscar_pontos_recarga(latitude: float, longitude: float) -> list:
        """Busca pontos de recarga via Google Places (dados simulados)."""
        return [
            {
                "place_id": "GP-EV-001",
                "nome": "Estacao GoodWe - FIAP Aclimacao",
                "endereco": "R. Aclimacao, 422 - Aclimacao, Sao Paulo",
                "rating": 4.5,
                "tipo": "ev_charging_station",
                "aberto_agora": True,
                "fotos": 3
            },
            {
                "place_id": "GP-EV-002",
                "nome": "Voltbras Charging Hub",
                "endereco": "R. Augusta, 2800 - Cerqueira Cesar, Sao Paulo",
                "rating": 4.2,
                "tipo": "ev_charging_station",
                "aberto_agora": True,
                "fotos": 5
            }
        ]


# ============================================================================
# 4. MOTOR DE IA
# ============================================================================

class MotorIA:
    """
    Nucleo de Inteligencia Artificial do EV ChargeOps.

    4 funcoes principais:
    - Interpretacao:  decodifica eventos de sessao
    - Preditividade:  preve demanda e necessidade de expansao
    - Precificacao:   calcula tarifas dinamicas (ponta/fora ponta)
    - Conversacao:    Sindico Virtual (NLP simplificado)
    """

    # Horarios de ponta (ANEEL): 18h-21h dias uteis
    HORARIO_PONTA_INICIO = 18
    HORARIO_PONTA_FIM = 21
    # Horario intermediario: 17h-18h e 21h-22h
    HORARIO_INTER_RANGES = [(17, 18), (21, 22)]

    # ---- INTERPRETACAO ----

    @staticmethod
    def interpretar_sessao(sessao: SessaoRecarga) -> dict:
        """Decodifica eventos de uma sessao de recarga."""
        duracao = (sessao.fim - sessao.inicio).total_seconds() / 3600 if sessao.fim else 0
        eficiencia = (sessao.energia_kwh / (sessao.potencia_media_kw * duracao) * 100
                      if duracao > 0 and sessao.potencia_media_kw > 0 else 0)

        classificacao = "normal"
        alertas = []

        if duracao > 8:
            classificacao = "prolongada"
            alertas.append("Sessao excedeu 8 horas - possivel veiculo estacionado sem recarga ativa")
        if eficiencia < 70 and duracao > 0:
            classificacao = "baixa_eficiencia"
            alertas.append(f"Eficiencia de {eficiencia:.0f}% - verificar conexao do cabo")
        if sessao.energia_kwh > 60:
            alertas.append("Consumo elevado (>60 kWh) - veiculo com bateria grande ou recarga completa")

        return {
            "sessao_id": sessao.id,
            "duracao_horas": round(duracao, 2),
            "eficiencia_pct": round(eficiencia, 1),
            "classificacao": classificacao,
            "alertas": alertas,
            "interpretacao": f"Sessao {sessao.id}: {sessao.energia_kwh:.1f} kWh em "
                            f"{duracao:.1f}h ({classificacao})"
        }

    # ---- PREDITIVIDADE ----

    @staticmethod
    def prever_demanda(sessoes_historico: list, dias_projecao: int = 30) -> dict:
        """
        Analisa historico de sessoes para prever demanda futura.
        Usa media movel simples + deteccao de tendencia.
        """
        if not sessoes_historico:
            return {"previsao": "Dados insuficientes", "recomendacao": "Aguardar historico"}

        consumo_diario = defaultdict(float)
        sessoes_por_dia = defaultdict(int)

        for s in sessoes_historico:
            if s.inicio:
                dia = s.inicio.strftime("%Y-%m-%d")
                consumo_diario[dia] += s.energia_kwh
                sessoes_por_dia[dia] += 1

        dias = sorted(consumo_diario.keys())
        consumos = [consumo_diario[d] for d in dias]

        if len(consumos) < 2:
            media = consumos[0] if consumos else 0
            tendencia = 0
        else:
            media = sum(consumos) / len(consumos)
            # Tendencia: comparar primeira e segunda metade
            meio = len(consumos) // 2
            media_1 = sum(consumos[:meio]) / meio if meio > 0 else 0
            media_2 = sum(consumos[meio:]) / (len(consumos) - meio)
            tendencia = ((media_2 - media_1) / media_1 * 100) if media_1 > 0 else 0

        media_sessoes = (sum(sessoes_por_dia.values()) / len(sessoes_por_dia)
                        if sessoes_por_dia else 0)

        previsao_mensal = media * dias_projecao
        pico_estimado = max(consumos) * 1.15 if consumos else 0

        precisa_expansao = pico_estimado > 200  # threshold: 200 kWh/dia

        return {
            "consumo_medio_diario_kwh": round(media, 1),
            "tendencia_pct": round(tendencia, 1),
            "direcao_tendencia": "crescente" if tendencia > 5 else
                                "decrescente" if tendencia < -5 else "estavel",
            "previsao_mensal_kwh": round(previsao_mensal, 1),
            "pico_estimado_kwh_dia": round(pico_estimado, 1),
            "media_sessoes_dia": round(media_sessoes, 1),
            "necessita_expansao": precisa_expansao,
            "recomendacao": (
                "ALERTA: Considerar instalacao de carregador adicional. "
                "Pico estimado excede capacidade otima."
                if precisa_expansao else
                "Infraestrutura atual atende a demanda projetada."
            )
        }

    # ---- PRECIFICACAO ----

    @classmethod
    def calcular_tarifa(cls, horario: datetime, tarifa_base: float) -> tuple:
        """
        Calculo dinamico de tarifas baseado no horario (ANEEL).
        Retorna (tarifa_aplicada, tipo_tarifa).

        Ponta (18h-21h):         tarifa_base * 1.5
        Intermediaria (17-18h, 21-22h): tarifa_base * 1.2
        Fora ponta:              tarifa_base * 1.0
        """
        hora = horario.hour
        dia_semana = horario.weekday()  # 0=seg, 6=dom

        # Fins de semana: sempre fora ponta
        if dia_semana >= 5:
            return tarifa_base, "fora_ponta"

        # Horario de ponta
        if cls.HORARIO_PONTA_INICIO <= hora < cls.HORARIO_PONTA_FIM:
            return round(tarifa_base * 1.5, 2), "ponta"

        # Horario intermediario
        for inicio, fim in cls.HORARIO_INTER_RANGES:
            if inicio <= hora < fim:
                return round(tarifa_base * 1.2, 2), "intermediaria"

        # Fora ponta
        return tarifa_base, "fora_ponta"

    # ---- CONVERSACAO (SINDICO VIRTUAL) ----

    @staticmethod
    def sindico_virtual(pergunta: str, dados_condominio: dict) -> str:
        """
        Sindico Virtual - Agente conversacional NLP simplificado.
        Traduz dados tecnicos em respostas gerenciais simples.
        Utiliza correspondencia de palavras-chave para classificar intencoes.
        """
        pergunta_lower = pergunta.lower().strip()

        # Intencoes e respostas
        if any(p in pergunta_lower for p in ["consumo", "gasto", "quanto", "kwh"]):
            total = dados_condominio.get("consumo_total_kwh", 0)
            custo = dados_condominio.get("custo_total", 0)
            unidades = dados_condominio.get("num_unidades_ativas", 0)
            return (
                f"No periodo analisado, o condominio consumiu {total:.1f} kWh "
                f"no total, gerando um custo de R$ {custo:.2f}. "
                f"{unidades} unidades utilizaram os carregadores. "
                f"O consumo medio por unidade foi de "
                f"{total/unidades:.1f} kWh." if unidades > 0 else
                "Nenhuma unidade utilizou os carregadores neste periodo."
            )

        elif any(p in pergunta_lower for p in ["fatura", "cobran", "pagar", "boleto", "rateio"]):
            faturas = dados_condominio.get("faturas_abertas", 0)
            total_pendente = dados_condominio.get("total_pendente", 0)
            return (
                f"Existem {faturas} faturas abertas no valor total de "
                f"R$ {total_pendente:.2f}. O rateio e calculado automaticamente "
                f"com base no consumo individual de cada unidade (kWh). "
                f"A tarifa segue a precificacao dinamica da ANEEL: "
                f"fora ponta (tarifa base), intermediaria (+20%) e ponta (+50%)."
            )

        elif any(p in pergunta_lower for p in ["carregador", "disponivel", "livre", "vaga"]):
            disponiveis = dados_condominio.get("carregadores_disponiveis", 0)
            total_carreg = dados_condominio.get("total_carregadores", 0)
            return (
                f"Neste momento, {disponiveis} de {total_carreg} carregadores "
                f"estao disponiveis. "
                + ("Todos os carregadores estao livres para uso."
                   if disponiveis == total_carreg else
                   "Recomendo agendar sua recarga para horarios fora de ponta "
                   "(antes das 17h ou apos as 22h) para economia na tarifa.")
            )

        elif any(p in pergunta_lower for p in ["horario", "melhor hora", "ponta", "tarifa",
                                                 "mais barato", "economia"]):
            return (
                "Os melhores horarios para recarregar com tarifa reduzida sao:\n"
                "- Fora ponta: 22h ate 17h (tarifa base)\n"
                "- Fins de semana: tarifa base o dia todo\n"
                "- Evite: 18h-21h em dias uteis (tarifa ponta, +50%)\n"
                "Recarregar durante a madrugada gera a maior economia."
            )

        elif any(p in pergunta_lower for p in ["pico", "demanda", "previsao", "futuro"]):
            tendencia = dados_condominio.get("tendencia", "estavel")
            previsao = dados_condominio.get("previsao_mensal_kwh", 0)
            return (
                f"A tendencia de consumo esta {tendencia}. "
                f"A previsao para o proximo mes e de {previsao:.0f} kWh. "
                + ("Recomendo avaliar a instalacao de um carregador adicional."
                   if tendencia == "crescente" else
                   "A infraestrutura atual atende bem a demanda.")
            )

        elif any(p in pergunta_lower for p in ["anomalia", "problema", "erro", "falha"]):
            anomalias = dados_condominio.get("anomalias", [])
            if anomalias:
                lista = "\n".join(f"  - {a}" for a in anomalias[:5])
                return f"Foram detectadas as seguintes anomalias:\n{lista}"
            return "Nenhuma anomalia detectada no periodo. Operacao normal."

        elif any(p in pergunta_lower for p in ["ajuda", "help", "o que voce faz", "funcoes"]):
            return (
                "Sou o Sindico Virtual do EV ChargeOps. Posso ajudar com:\n"
                "- Consumo e gastos do condominio\n"
                "- Faturas e cobrancas por unidade\n"
                "- Disponibilidade dos carregadores\n"
                "- Melhores horarios para recarga (economia)\n"
                "- Previsao de demanda e tendencias\n"
                "- Anomalias e problemas detectados\n"
                "Pergunte-me qualquer coisa sobre a recarga do condominio!"
            )

        else:
            return (
                "Entendi sua pergunta, mas nao consegui classificar a intencao. "
                "Tente perguntar sobre: consumo, faturas, carregadores disponiveis, "
                "melhores horarios, previsao de demanda ou anomalias."
            )

    # ---- DETECCAO DE ANOMALIAS ----

    @staticmethod
    def detectar_anomalias(sessoes: list) -> list:
        """Detecta anomalias nas sessoes de recarga."""
        anomalias = []
        if not sessoes:
            return anomalias

        consumos = [s.energia_kwh for s in sessoes if s.energia_kwh > 0]
        if len(consumos) < 3:
            return anomalias

        media = sum(consumos) / len(consumos)
        variancia = sum((x - media) ** 2 for x in consumos) / len(consumos)
        desvio = variancia ** 0.5

        for s in sessoes:
            if s.energia_kwh > 0 and abs(s.energia_kwh - media) > 2 * desvio:
                anomalias.append(
                    f"Sessao {s.id}: consumo de {s.energia_kwh:.1f} kWh "
                    f"desvia significativamente da media ({media:.1f} kWh)"
                )

            if s.inicio and s.fim:
                duracao_h = (s.fim - s.inicio).total_seconds() / 3600
                if duracao_h > 12:
                    anomalias.append(
                        f"Sessao {s.id}: duracao de {duracao_h:.1f}h - "
                        f"possivel veiculo abandonado no carregador"
                    )
                if duracao_h > 0 and s.energia_kwh / duracao_h < 1.0:
                    anomalias.append(
                        f"Sessao {s.id}: taxa de carga muito baixa "
                        f"({s.energia_kwh/duracao_h:.1f} kWh/h) - verificar equipamento"
                    )

        return anomalias


# ============================================================================
# 5. GERENCIADOR DE SESSOES (GESTAO DE USO)
# ============================================================================

class GerenciadorSessoes:
    """Gerencia sessoes de recarga por usuario/unidade habitacional."""

    def __init__(self):
        self.sessoes: list = []
        self.sessoes_ativas: dict = {}  # carregador_id -> sessao
        self.goodwe_api = GoodWeAPISimulator()

    def iniciar_sessao(self, unidade: UnidadeHabitacional,
                       carregador: Carregador,
                       tarifa_base: float) -> SessaoRecarga:
        """Inicia uma nova sessao de recarga."""
        if carregador.status != "disponivel":
            raise ValueError(f"Carregador {carregador.id} nao esta disponivel "
                           f"(status: {carregador.status})")

        agora = datetime.now()
        tarifa, tipo_tarifa = MotorIA.calcular_tarifa(agora, tarifa_base)

        sessao = SessaoRecarga(
            unidade_id=unidade.id,
            carregador_id=carregador.id,
            inicio=agora,
            tarifa_aplicada=tarifa,
            tipo_tarifa=tipo_tarifa,
            potencia_media_kw=carregador.potencia_kw
        )

        # Simular chamada OCPP via GoodWe API
        rfid = unidade.rfid_cards[0] if unidade.rfid_cards else "RFID-DEFAULT"
        self.goodwe_api.start_charging(carregador, rfid)

        carregador.status = "em_uso"
        carregador.sessao_atual_id = sessao.id
        self.sessoes_ativas[carregador.id] = sessao
        self.sessoes.append(sessao)

        return sessao

    def finalizar_sessao(self, carregador: Carregador,
                         energia_kwh: float = None) -> SessaoRecarga:
        """Finaliza a sessao ativa no carregador."""
        if carregador.id not in self.sessoes_ativas:
            raise ValueError(f"Nenhuma sessao ativa no carregador {carregador.id}")

        sessao = self.sessoes_ativas[carregador.id]
        sessao.fim = datetime.now()
        sessao.status = "finalizada"

        if energia_kwh is None:
            duracao_h = (sessao.fim - sessao.inicio).total_seconds() / 3600
            fator = random.uniform(0.7, 0.95)
            energia_kwh = round(carregador.potencia_kw * duracao_h * fator, 2)

        sessao.energia_kwh = energia_kwh
        sessao.custo_total = round(energia_kwh * sessao.tarifa_aplicada, 2)

        # Simular parada via GoodWe API
        self.goodwe_api.stop_charging(carregador, energia_kwh)

        carregador.status = "disponivel"
        carregador.sessao_atual_id = None
        del self.sessoes_ativas[carregador.id]

        return sessao


# ============================================================================
# 6. MOTOR DE RATEIO E FATURAMENTO
# ============================================================================

class MotorFaturamento:
    """Processamento de consumo individual para cobranca automatica."""

    def __init__(self):
        self.faturas: list = []

    def gerar_fatura_mensal(self, unidade: UnidadeHabitacional,
                            sessoes: list, mes_ref: str) -> Fatura:
        """Gera fatura mensal para uma unidade com base nas sessoes."""
        sessoes_unidade = [s for s in sessoes
                          if s.unidade_id == unidade.id
                          and s.status == "finalizada"
                          and s.inicio and s.inicio.strftime("%Y-%m") == mes_ref]

        total_kwh = sum(s.energia_kwh for s in sessoes_unidade)
        total_reais = sum(s.custo_total for s in sessoes_unidade)

        hoje = datetime.now()
        fatura = Fatura(
            unidade_id=unidade.id,
            mes_referencia=mes_ref,
            sessoes=[s.id for s in sessoes_unidade],
            total_kwh=round(total_kwh, 2),
            total_reais=round(total_reais, 2),
            data_emissao=hoje,
            data_vencimento=hoje + timedelta(days=30),
            status="fechada"
        )

        unidade.saldo_devedor += fatura.total_reais
        self.faturas.append(fatura)
        return fatura

    def relatorio_rateio(self, unidades: list, sessoes: list,
                         mes_ref: str) -> dict:
        """Gera relatorio de rateio do condominio."""
        rateio = {}
        total_cond = 0

        for u in unidades:
            sessoes_u = [s for s in sessoes
                        if s.unidade_id == u.id
                        and s.status == "finalizada"
                        and s.inicio and s.inicio.strftime("%Y-%m") == mes_ref]
            kwh = sum(s.energia_kwh for s in sessoes_u)
            custo = sum(s.custo_total for s in sessoes_u)
            num = len(sessoes_u)

            rateio[u.numero] = {
                "bloco": u.bloco,
                "proprietario": u.proprietario,
                "sessoes": num,
                "kwh": round(kwh, 2),
                "custo": round(custo, 2)
            }
            total_cond += custo

        return {
            "mes_referencia": mes_ref,
            "total_condominio_reais": round(total_cond, 2),
            "rateio_por_unidade": rateio
        }


# ============================================================================
# 7. PLATAFORMA EV CHARGEOPS (ORQUESTRADOR PRINCIPAL)
# ============================================================================

class EVChargeOps:
    """
    Plataforma principal EV ChargeOps.
    Orquestra todos os modulos: sessoes, faturamento, IA, integracoes.
    """

    def __init__(self):
        self.condominio: Optional[Condominio] = None
        self.unidades: list = []
        self.carregadores: list = []
        self.gerenciador = GerenciadorSessoes()
        self.faturamento = MotorFaturamento()
        self.motor_ia = MotorIA()
        self.ocm = OpenChargeMapIntegration()
        self.google_places = GooglePlacesIntegration()

    # ---- SETUP INICIAL ----

    def setup_demo(self):
        """Configura ambiente de demonstracao com dados simulados."""
        # Criar condominio
        self.condominio = Condominio(
            nome="Condominio Residencial Parque Verde",
            endereco="Rua das Palmeiras, 500 - Vila Mariana, Sao Paulo - SP",
            capacidade_total_kw=44.0,
            num_carregadores=4
        )

        # Criar carregadores GoodWe HCA G2
        modelos = [
            ("GW7K-HCA-20", 7.0, "Garagem B1 - Vaga 01"),
            ("GW11K-HCA-20", 11.0, "Garagem B1 - Vaga 02"),
            ("GW11K-HCA-20", 11.0, "Garagem B2 - Vaga 01"),
            ("GW22K-HCA-20", 22.0, "Garagem B2 - Vaga 02")
        ]
        for modelo, pot, loc in modelos:
            c = Carregador(
                modelo=modelo,
                potencia_kw=pot,
                condominio_id=self.condominio.id,
                localizacao=loc
            )
            self.carregadores.append(c)

        # Criar unidades habitacionais
        moradores = [
            ("101", "A", "Ana Silva", ["RFID-A101-01"]),
            ("102", "A", "Bruno Costa", ["RFID-A102-01", "RFID-A102-02"]),
            ("201", "A", "Carla Mendes", ["RFID-A201-01"]),
            ("202", "A", "Daniel Oliveira", ["RFID-A202-01"]),
            ("301", "B", "Elena Souza", ["RFID-B301-01"]),
            ("302", "B", "Felipe Santos", ["RFID-B302-01"]),
            ("401", "B", "Gabriela Lima", ["RFID-B401-01"]),
            ("402", "B", "Henrique Rocha", ["RFID-B402-01"]),
        ]
        for num, bloco, prop, rfids in moradores:
            u = UnidadeHabitacional(
                numero=num,
                bloco=bloco,
                proprietario=prop,
                condominio_id=self.condominio.id,
                rfid_cards=rfids
            )
            self.unidades.append(u)

        print(f"  Condominio: {self.condominio.nome}")
        print(f"  Carregadores: {len(self.carregadores)} "
              f"(capacidade total: {self.condominio.capacidade_total_kw} kW)")
        print(f"  Unidades: {len(self.unidades)} apartamentos")

    def gerar_historico_simulado(self, dias: int = 30):
        """Gera historico de sessoes simuladas para demonstracao."""
        agora = datetime.now()
        sessoes_geradas = 0

        for dia in range(dias, 0, -1):
            data_base = agora - timedelta(days=dia)
            # 3 a 8 sessoes por dia
            n_sessoes = random.randint(3, 8)

            for _ in range(n_sessoes):
                unidade = random.choice(self.unidades)
                carregador = random.choice(self.carregadores)

                hora_inicio = random.randint(0, 23)
                minuto_inicio = random.randint(0, 59)
                inicio = data_base.replace(hour=hora_inicio, minute=minuto_inicio,
                                          second=0, microsecond=0)
                duracao_h = random.uniform(0.5, 6.0)
                fim = inicio + timedelta(hours=duracao_h)

                tarifa, tipo = self.motor_ia.calcular_tarifa(
                    inicio, self.condominio.tarifa_kwh_base)
                fator = random.uniform(0.65, 0.95)
                energia = round(carregador.potencia_kw * duracao_h * fator, 2)
                custo = round(energia * tarifa, 2)

                sessao = SessaoRecarga(
                    unidade_id=unidade.id,
                    carregador_id=carregador.id,
                    inicio=inicio,
                    fim=fim,
                    energia_kwh=energia,
                    custo_total=custo,
                    tarifa_aplicada=tarifa,
                    status="finalizada",
                    tipo_tarifa=tipo,
                    potencia_media_kw=carregador.potencia_kw
                )
                self.gerenciador.sessoes.append(sessao)
                sessoes_geradas += 1

        print(f"  {sessoes_geradas} sessoes historicas geradas ({dias} dias)")

    # ---- CONSULTAS E RELATORIOS ----

    def status_carregadores(self) -> list:
        """Retorna status de todos os carregadores."""
        status_list = []
        api = GoodWeAPISimulator()
        for c in self.carregadores:
            s = api.get_charger_status(c)
            s["localizacao"] = c.localizacao
            status_list.append(s)
        return status_list

    def dashboard_condominio(self) -> dict:
        """Gera dados para o painel gerencial do sindico."""
        sessoes = self.gerenciador.sessoes
        sessoes_finalizadas = [s for s in sessoes if s.status == "finalizada"]

        total_kwh = sum(s.energia_kwh for s in sessoes_finalizadas)
        total_custo = sum(s.custo_total for s in sessoes_finalizadas)
        total_sessoes = len(sessoes_finalizadas)

        # Consumo por unidade
        consumo_unidade = defaultdict(float)
        for s in sessoes_finalizadas:
            consumo_unidade[s.unidade_id] += s.energia_kwh

        # Unidades ativas (que usaram carregador)
        unidades_ativas = len(set(s.unidade_id for s in sessoes_finalizadas))

        # Distribuicao por tipo de tarifa
        dist_tarifa = defaultdict(int)
        for s in sessoes_finalizadas:
            dist_tarifa[s.tipo_tarifa] += 1

        # Carregadores disponiveis agora
        disponiveis = sum(1 for c in self.carregadores if c.status == "disponivel")

        # Horarios mais utilizados
        horas = defaultdict(int)
        for s in sessoes_finalizadas:
            if s.inicio:
                horas[s.inicio.hour] += 1

        hora_pico = max(horas, key=horas.get) if horas else 0

        return {
            "condominio": self.condominio.nome,
            "total_sessoes": total_sessoes,
            "total_kwh": round(total_kwh, 1),
            "total_custo": round(total_custo, 2),
            "media_kwh_sessao": round(total_kwh / total_sessoes, 1) if total_sessoes > 0 else 0,
            "unidades_ativas": unidades_ativas,
            "carregadores_disponiveis": disponiveis,
            "total_carregadores": len(self.carregadores),
            "distribuicao_tarifa": dict(dist_tarifa),
            "hora_pico": hora_pico,
            "consumo_por_unidade": dict(consumo_unidade),
        }

    def executar_analise_ia(self) -> dict:
        """Executa analise completa do Motor de IA."""
        sessoes = self.gerenciador.sessoes
        sessoes_fin = [s for s in sessoes if s.status == "finalizada"]

        # Preditividade
        previsao = self.motor_ia.prever_demanda(sessoes_fin)

        # Anomalias
        anomalias = self.motor_ia.detectar_anomalias(sessoes_fin)

        # Interpretar ultimas 5 sessoes
        interpretacoes = []
        for s in sessoes_fin[-5:]:
            interpretacoes.append(self.motor_ia.interpretar_sessao(s))

        # Dados para Sindico Virtual
        dashboard = self.dashboard_condominio()
        dados_sindico = {
            "consumo_total_kwh": dashboard["total_kwh"],
            "custo_total": dashboard["total_custo"],
            "num_unidades_ativas": dashboard["unidades_ativas"],
            "carregadores_disponiveis": dashboard["carregadores_disponiveis"],
            "total_carregadores": dashboard["total_carregadores"],
            "faturas_abertas": len(self.faturamento.faturas),
            "total_pendente": sum(f.total_reais for f in self.faturamento.faturas
                                 if f.status != "paga"),
            "tendencia": previsao.get("direcao_tendencia", "estavel"),
            "previsao_mensal_kwh": previsao.get("previsao_mensal_kwh", 0),
            "anomalias": anomalias[:5]
        }

        return {
            "previsao_demanda": previsao,
            "anomalias": anomalias,
            "interpretacoes": interpretacoes,
            "dados_sindico": dados_sindico
        }

    # ---- BUSCA ELETROPOSTOS EXTERNOS ----

    def buscar_eletropostos_proximos(self) -> dict:
        """Busca eletropostos proximos via APIs externas."""
        lat, lon = -23.5868, -46.6503  # Vila Mariana, SP
        ocm_results = self.ocm.buscar_eletropostos(lat, lon)
        gp_results = self.google_places.buscar_pontos_recarga(lat, lon)
        return {
            "open_charge_map": ocm_results,
            "google_places": gp_results
        }


# ============================================================================
# 8. VISUALIZACAO E DASHBOARD (GRAFICOS MATPLOTLIB)
# ============================================================================

def gerar_graficos(plataforma: EVChargeOps, pasta_saida: str):
    """Gera graficos de analise para o dashboard e relatorio."""
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates

    sessoes = [s for s in plataforma.gerenciador.sessoes if s.status == "finalizada"]
    if not sessoes:
        print("  Sem sessoes para gerar graficos.")
        return

    os.makedirs(pasta_saida, exist_ok=True)

    # --- GRAFICO 1: Consumo diario (kWh) ---
    consumo_diario = defaultdict(float)
    for s in sessoes:
        if s.inicio:
            dia = s.inicio.date()
            consumo_diario[dia] += s.energia_kwh

    dias = sorted(consumo_diario.keys())
    valores = [consumo_diario[d] for d in dias]

    fig, ax = plt.subplots(figsize=(12, 5))
    ax.bar(dias, valores, color='#E63946', alpha=0.85, width=0.8)
    ax.set_xlabel('Data', fontsize=11)
    ax.set_ylabel('Consumo (kWh)', fontsize=11)
    ax.set_title('Consumo Diário de Energia - EV ChargeOps', fontsize=14, fontweight='bold')
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m'))
    ax.xaxis.set_major_locator(mdates.DayLocator(interval=max(1, len(dias)//15)))
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig(os.path.join(pasta_saida, 'consumo_diario.png'), dpi=150)
    plt.close()

    # --- GRAFICO 2: Consumo por unidade ---
    consumo_unid = defaultdict(float)
    unid_map = {u.id: f"{u.bloco}-{u.numero}" for u in plataforma.unidades}
    for s in sessoes:
        label = unid_map.get(s.unidade_id, s.unidade_id[:6])
        consumo_unid[label] += s.energia_kwh

    labels = sorted(consumo_unid.keys())
    vals = [consumo_unid[l] for l in labels]

    fig, ax = plt.subplots(figsize=(10, 5))
    cores = ['#E63946', '#457B9D', '#1D3557', '#F4A261',
             '#2A9D8F', '#264653', '#E9C46A', '#F77F00']
    ax.barh(labels, vals, color=cores[:len(labels)])
    ax.set_xlabel('Consumo Total (kWh)', fontsize=11)
    ax.set_ylabel('Unidade', fontsize=11)
    ax.set_title('Consumo por Unidade Habitacional', fontsize=14, fontweight='bold')
    for i, v in enumerate(vals):
        ax.text(v + 1, i, f'{v:.0f} kWh', va='center', fontsize=9)
    plt.tight_layout()
    plt.savefig(os.path.join(pasta_saida, 'consumo_por_unidade.png'), dpi=150)
    plt.close()

    # --- GRAFICO 3: Distribuicao por horario ---
    horas = defaultdict(int)
    for s in sessoes:
        if s.inicio:
            horas[s.inicio.hour] += 1

    h_labels = list(range(24))
    h_vals = [horas.get(h, 0) for h in h_labels]

    fig, ax = plt.subplots(figsize=(12, 5))
    cores_h = ['#E63946' if 18 <= h < 21 else '#F4A261' if h in [17, 21]
               else '#457B9D' for h in h_labels]
    ax.bar(h_labels, h_vals, color=cores_h, alpha=0.85)
    ax.set_xlabel('Hora do Dia', fontsize=11)
    ax.set_ylabel('Número de Sessões', fontsize=11)
    ax.set_title('Distribuição de Sessões por Horário', fontsize=14, fontweight='bold')
    ax.set_xticks(h_labels)
    ax.set_xticklabels([f'{h:02d}h' for h in h_labels], rotation=45, fontsize=8)

    # Legenda
    from matplotlib.patches import Patch
    legend_elements = [
        Patch(facecolor='#E63946', label='Ponta (18h-21h)'),
        Patch(facecolor='#F4A261', label='Intermediária'),
        Patch(facecolor='#457B9D', label='Fora Ponta')
    ]
    ax.legend(handles=legend_elements, loc='upper left')
    plt.tight_layout()
    plt.savefig(os.path.join(pasta_saida, 'distribuicao_horario.png'), dpi=150)
    plt.close()

    # --- GRAFICO 4: Distribuicao de tarifa (Pizza) ---
    dist_tarifa = defaultdict(float)
    for s in sessoes:
        dist_tarifa[s.tipo_tarifa] += s.custo_total

    fig, ax = plt.subplots(figsize=(8, 6))
    tarifa_labels = {
        'fora_ponta': 'Fora Ponta',
        'intermediaria': 'Intermediária',
        'ponta': 'Ponta'
    }
    t_labels = [tarifa_labels.get(k, k) for k in dist_tarifa.keys()]
    t_vals = list(dist_tarifa.values())
    t_cores = ['#457B9D', '#F4A261', '#E63946']

    ax.pie(t_vals, labels=t_labels, colors=t_cores[:len(t_labels)],
           autopct='%1.1f%%', startangle=90, textprops={'fontsize': 11})
    ax.set_title('Distribuição de Custos por Tipo de Tarifa',
                fontsize=14, fontweight='bold')
    plt.tight_layout()
    plt.savefig(os.path.join(pasta_saida, 'distribuicao_tarifa.png'), dpi=150)
    plt.close()

    # --- GRAFICO 5: Evolucao acumulada de custo ---
    custo_acum = {}
    custo_total = 0
    for dia in dias:
        sessoes_dia = [s for s in sessoes if s.inicio and s.inicio.date() == dia]
        custo_total += sum(s.custo_total for s in sessoes_dia)
        custo_acum[dia] = custo_total

    fig, ax = plt.subplots(figsize=(12, 5))
    ax.fill_between(list(custo_acum.keys()), list(custo_acum.values()),
                    alpha=0.3, color='#E63946')
    ax.plot(list(custo_acum.keys()), list(custo_acum.values()),
            color='#E63946', linewidth=2)
    ax.set_xlabel('Data', fontsize=11)
    ax.set_ylabel('Custo Acumulado (R$)', fontsize=11)
    ax.set_title('Evolução Acumulada de Custos', fontsize=14, fontweight='bold')
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m'))
    ax.xaxis.set_major_locator(mdates.DayLocator(interval=max(1, len(dias)//15)))
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig(os.path.join(pasta_saida, 'custo_acumulado.png'), dpi=150)
    plt.close()

    print(f"  5 graficos gerados em: {pasta_saida}/")


# ============================================================================
# 9. GERADOR DE RELATORIO PDF
# ============================================================================

def gerar_relatorio_pdf(plataforma: EVChargeOps, pasta_saida: str, save=True):
    """Gera relatorio tecnico em PDF do EV ChargeOps."""
    from fpdf import FPDF

    class RelatorioPDF(FPDF):
        def header(self):
            self.set_font('Helvetica', 'B', 10)
            self.set_text_color(180, 0, 0)
            self.cell(0, 8, 'EV ChargeOps | FIAP + GoodWe 2026', 0, 0, 'L')
            self.cell(0, 8, 'Marcelo B. Baldin - RM568746', 0, 1, 'R')
            self.set_draw_color(180, 0, 0)
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(4)

        def footer(self):
            self.set_y(-15)
            self.set_font('Helvetica', 'I', 8)
            self.set_text_color(128, 128, 128)
            self.cell(0, 10, f'Página {self.page_no()}/{{nb}}', 0, 0, 'C')

        def titulo_secao(self, num, titulo):
            self.set_font('Helvetica', 'B', 14)
            self.set_text_color(29, 53, 87)
            self.cell(0, 10, f'{num}. {titulo}', 0, 1)
            self.set_draw_color(230, 57, 70)
            self.line(10, self.get_y(), 80, self.get_y())
            self.ln(4)

        def subtitulo(self, texto):
            self.set_font('Helvetica', 'B', 11)
            self.set_text_color(69, 123, 157)
            self.cell(0, 8, texto, 0, 1)
            self.ln(2)

        def corpo(self, texto):
            self.set_font('Helvetica', '', 10)
            self.set_text_color(40, 40, 40)
            self.multi_cell(0, 5, texto)
            self.ln(2)

        def item_lista(self, texto):
            self.set_font('Helvetica', '', 10)
            self.set_text_color(40, 40, 40)
            self.multi_cell(0, 5, f'  - {texto}')
            self.ln(1)

    pdf = RelatorioPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)

    # --- CAPA ---
    pdf.add_page()
    pdf.ln(30)
    pdf.set_font('Helvetica', 'B', 28)
    pdf.set_text_color(29, 53, 87)
    pdf.cell(0, 15, 'EV ChargeOps', 0, 1, 'C')
    pdf.set_font('Helvetica', '', 14)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 10, 'Plataforma de Gestão Compartilhada de Recarga', 0, 1, 'C')
    pdf.cell(0, 8, 'para Condomínios e Setor Corporativo', 0, 1, 'C')
    pdf.ln(10)
    pdf.set_draw_color(230, 57, 70)
    pdf.line(60, pdf.get_y(), 150, pdf.get_y())
    pdf.ln(10)
    pdf.set_font('Helvetica', '', 12)
    pdf.set_text_color(60, 60, 60)
    pdf.cell(0, 8, 'Challenge FIAP + GoodWe 2026', 0, 1, 'C')
    pdf.cell(0, 8, 'Ciências da Computação - 1º Ano Online', 0, 1, 'C')
    pdf.ln(5)
    pdf.set_font('Helvetica', 'B', 12)
    pdf.cell(0, 8, 'Marcelo Bastianello Baldin', 0, 1, 'C')
    pdf.set_font('Helvetica', '', 11)
    pdf.cell(0, 8, 'RM568746', 0, 1, 'C')
    pdf.ln(10)
    pdf.set_font('Helvetica', 'I', 10)
    pdf.set_text_color(128, 128, 128)
    pdf.cell(0, 8, datetime.now().strftime('Abril de 2026'), 0, 1, 'C')

    # --- 1. INTRODUCAO ---
    pdf.add_page()
    pdf.titulo_secao(1, 'Introdução')
    pdf.corpo(
        'O EV ChargeOps é uma plataforma de software desenvolvida para resolver '
        'o problema central da gestão compartilhada de recarga de veículos elétricos '
        'em condomínios residenciais e corporativos. Com o crescimento acelerado da '
        'frota de veículos elétricos no Brasil, surge a necessidade de regras claras, '
        'rateio justo e análise de consumo contínua em infraestruturas compartilhadas.'
    )
    pdf.corpo(
        'A solução transforma o carregador de um simples fornecedor de energia (kWh) '
        'em um hub de inteligência, capaz de gerar dados estruturados e inteligência '
        'acionável a partir de cada sessão de recarga. O sistema é construído sobre '
        'o ecossistema GoodWe (Linha HCA G2) e integra APIs de localização e um '
        'agente conversacional (Síndico Virtual) baseado em NLP.'
    )
    pdf.subtitulo('Requisitos Atendidos (Turma Online)')
    pdf.item_lista('Estruturação de sessões de recarga por usuário ou unidade habitacional')
    pdf.item_lista('Processamento de consumo individual para cobrança automática da unidade')
    pdf.item_lista('Sistema de gerenciamento inteligente com interface para usuário e gestão')

    # --- 2. ARQUITETURA ---
    pdf.add_page()
    pdf.titulo_secao(2, 'Arquitetura Funcional')
    pdf.corpo(
        'O EV ChargeOps utiliza uma arquitetura híbrida de três camadas, conforme '
        'especificado no playbook do Challenge:'
    )
    pdf.subtitulo('Camada Física')
    pdf.corpo(
        'Hardware GoodWe Linha HCA G2 (modelos GW7K, GW11K, GW22K). '
        'Comunicação via protocolos industriais MODBUS (RS-485) e OCPP 1.6J. '
        'Autenticação por RFID (até 10 cartões por carregador). '
        'Conectividade: RS-485 + LAN + Wi-Fi + Bluetooth integrados.'
    )
    pdf.subtitulo('Camada de Conectividade')
    pdf.corpo(
        'Rede local (LAN/Wi-Fi) conectando carregadores ao servidor da aplicação. '
        'Comunicação bidirecional em tempo real via protocolo OCPP.'
    )
    pdf.subtitulo('Camada Digital')
    pdf.corpo(
        'Software EV ChargeOps em Python, composto por módulos:\n'
        '- Gerenciador de Sessões (Gestão de Uso)\n'
        '- Motor de Faturamento (Rateio e Cobrança)\n'
        '- Motor de IA (Interpretação, Preditividade, Precificação, Conversação)\n'
        '- Integrações (GoodWe API, Open Charge Map, Google Places)\n'
        '- Dashboard de Visualização (gráficos e relatórios)'
    )
    pdf.subtitulo('Fluxo de Dados (Data Flow)')
    pdf.corpo(
        '1. Morador aproxima cartão RFID no carregador GoodWe\n'
        '2. Sistema identifica unidade habitacional vinculada ao RFID\n'
        '3. Sessão de recarga é iniciada via OCPP (RemoteStartTransaction)\n'
        '4. Medidor MODBUS registra consumo em tempo real (kWh, V, A, W)\n'
        '5. Ao finalizar, Motor de IA calcula tarifa dinâmica (ponta/fora ponta)\n'
        '6. Consumo e custo são registrados na conta da unidade\n'
        '7. Fatura mensal é gerada automaticamente com rateio individual\n'
        '8. Síndico Virtual disponibiliza insights via linguagem natural'
    )

    # --- 3. PILARES DO CHARGEOPS ---
    pdf.add_page()
    pdf.titulo_secao(3, 'Os 4 Pilares do EV ChargeOps')

    pdf.subtitulo('3.1 Gestão de Uso')
    pdf.corpo(
        'Estruturação de sessões de recarga por usuário ou unidade habitacional. '
        'Cada sessão registra: unidade, carregador, horário de início/fim, '
        'energia consumida (kWh), tarifa aplicada e custo total. '
        'Autenticação via RFID vinculado à unidade. '
        'Implementado na classe GerenciadorSessoes.'
    )

    pdf.subtitulo('3.2 Gerenciamento de Recarga')
    pdf.corpo(
        'Definição de regras para uso compartilhado e sugestão de horários otimizados. '
        'O Motor de IA analisa padrões de uso e sugere horários fora de ponta '
        '(antes das 17h ou após as 22h) para economia. '
        'O sistema gerencia a fila de espera e prioridade de carregadores. '
        'Integra com OCPP para controle remoto de início/parada de sessões.'
    )

    pdf.subtitulo('3.3 Rateio e Faturamento')
    pdf.corpo(
        'Processamento de consumo individual (kWh) para repasse automático na '
        'taxa condominial. Cada unidade recebe fatura mensal detalhada com: '
        'número de sessões, kWh consumidos, tarifa aplicada por sessão, custo total. '
        'Detecção de anomalias por desvio estatístico (>2 sigma). '
        'Base regulatória: ANEEL Resolução Normativa 1.000/2021 (precificação livre).'
    )

    pdf.subtitulo('3.4 Visualização e Interação')
    pdf.corpo(
        'Painéis gerenciais para síndicos com gráficos de consumo diário, '
        'consumo por unidade, distribuição horária e evolução de custos. '
        'Insights de consumo para moradores via Síndico Virtual (NLP). '
        'Gráficos gerados com Matplotlib em 5 visualizações diferentes.'
    )

    # --- 4. MOTOR DE IA ---
    pdf.add_page()
    pdf.titulo_secao(4, 'Motor de IA')
    pdf.corpo(
        'A IA não é uma feature extra; é a camada que dá sentido aos dados brutos '
        'gerados pela rede de carregadores. O Motor de IA do EV ChargeOps opera '
        'em quatro dimensões:'
    )

    pdf.subtitulo('4.1 Interpretação')
    pdf.corpo(
        'Decodificação de eventos de sessão: classifica sessões como normal, '
        'prolongada ou de baixa eficiência. Analisa duração, consumo e taxa de '
        'carga para gerar alertas automáticos. Exemplo: sessão >8h gera alerta '
        'de veículo estacionado sem recarga ativa.'
    )

    pdf.subtitulo('4.2 Preditividade')
    pdf.corpo(
        'Previsão de demanda usando média móvel e detecção de tendência. '
        'Compara primeira e segunda metade do histórico para identificar '
        'crescimento ou declínio. Projeta consumo mensal e pico diário estimado. '
        'Recomenda expansão de infraestrutura quando pico >200 kWh/dia.'
    )

    pdf.subtitulo('4.3 Precificação')
    pdf.corpo(
        'Cálculo dinâmico de tarifas baseado no horário conforme ANEEL:\n'
        '- Fora ponta (22h-17h + fins de semana): tarifa base (R$ 0,85/kWh)\n'
        '- Intermediária (17h-18h e 21h-22h): +20% (R$ 1,02/kWh)\n'
        '- Ponta (18h-21h dias úteis): +50% (R$ 1,28/kWh)\n'
        'Incentiva recarga em horários de menor demanda da rede elétrica.'
    )

    pdf.subtitulo('4.4 Conversação - Síndico Virtual')
    pdf.corpo(
        'Agente interativo conversacional baseado em NLP simplificado. '
        'Traduz dados técnicos (picos de ociosidade, faturamento parcial) em '
        'respostas gerenciais simples. Responde sobre: consumo, faturas, '
        'disponibilidade de carregadores, melhores horários, previsão de demanda '
        'e anomalias detectadas. Classificação de intenções por palavras-chave.'
    )

    # --- 5. INTEGRACOES ---
    pdf.add_page()
    pdf.titulo_secao(5, 'Integrações')

    pdf.subtitulo('5.1 GoodWe API (Simulada)')
    pdf.corpo(
        'Simulação da API GoodWe para os carregadores EV Charger FIAP. '
        'Implementa operações OCPP: RemoteStartTransaction, RemoteStopTransaction, '
        'MeterValues. Retorna dados em formato JSON: status, potência, temperatura, '
        'firmware, tensão, corrente, fator de potência. '
        'Em produção, conectaria diretamente ao hardware via OCPP 1.6J.'
    )

    pdf.subtitulo('5.2 Open Charge Map API')
    pdf.corpo(
        'Integração com o mapa público de eletropostos (openchargemap.io). '
        'Busca pontos de recarga próximos ao condomínio por coordenadas GPS. '
        'Retorna: nome, endereço, conectores, potência, preço e status. '
        'Permite cruzamento da infraestrutura privada com o mapa público '
        'para sugerir recarga de oportunidade quando carregadores do '
        'condomínio estão ocupados.'
    )

    pdf.subtitulo('5.3 Google Places API (evChargeOptions)')
    pdf.corpo(
        'Integração com Google Places para localização de pontos de recarga. '
        'Complementa o Open Charge Map com dados de avaliação (rating), '
        'fotos e horário de funcionamento. Permite ao morador encontrar '
        'alternativas de recarga na região.'
    )

    # --- 6. RESULTADOS SIMULADOS ---
    pdf.add_page()
    pdf.titulo_secao(6, 'Resultados da Simulação')

    dashboard = plataforma.dashboard_condominio()
    analise = plataforma.executar_analise_ia()
    previsao = analise["previsao_demanda"]

    pdf.subtitulo('Dashboard do Condomínio')
    pdf.corpo(f'Condomínio: {dashboard["condominio"]}')
    pdf.corpo(f'Total de sessões: {dashboard["total_sessoes"]}')
    pdf.corpo(f'Consumo total: {dashboard["total_kwh"]:.1f} kWh')
    pdf.corpo(f'Custo total: R$ {dashboard["total_custo"]:.2f}')
    pdf.corpo(f'Média por sessão: {dashboard["media_kwh_sessao"]:.1f} kWh')
    pdf.corpo(f'Unidades ativas: {dashboard["unidades_ativas"]} de {len(plataforma.unidades)}')
    pdf.corpo(f'Hora de pico: {dashboard["hora_pico"]:02d}h')

    pdf.subtitulo('Previsão de Demanda (IA)')
    pdf.corpo(f'Consumo médio diário: {previsao["consumo_medio_diario_kwh"]:.1f} kWh')
    pdf.corpo(f'Tendência: {previsao["direcao_tendencia"]} ({previsao["tendencia_pct"]:+.1f}%)')
    pdf.corpo(f'Previsão mensal: {previsao["previsao_mensal_kwh"]:.0f} kWh')
    pdf.corpo(f'Pico estimado: {previsao["pico_estimado_kwh_dia"]:.0f} kWh/dia')
    pdf.corpo(f'Recomendação: {previsao["recomendacao"]}')

    if analise["anomalias"]:
        pdf.subtitulo('Anomalias Detectadas')
        for a in analise["anomalias"][:5]:
            pdf.item_lista(a)

    # Graficos
    graficos = ['consumo_diario.png', 'consumo_por_unidade.png',
                'distribuicao_horario.png', 'distribuicao_tarifa.png',
                'custo_acumulado.png']

    for g in graficos:
        caminho = os.path.join(pasta_saida, g)
        if os.path.exists(caminho):
            pdf.add_page()
            pdf.image(caminho, x=10, y=30, w=190)

    # --- 7. BASE REGULATORIA ---
    pdf.add_page()
    pdf.titulo_secao(7, 'Base Regulatória')
    pdf.subtitulo('Exploração Comercial Livre')
    pdf.corpo(
        'Resolução Normativa ANEEL n. 1.000/2021: a recarga pública e comercial '
        'pode operar com preços livremente negociados, o que torna legal e '
        'obrigatório o desenvolvimento de plataformas de tarifação dinâmica. '
        'O EV ChargeOps implementa precificação flexível baseada em horário.'
    )
    pdf.subtitulo('Obrigatoriedade de Interoperabilidade')
    pdf.corpo(
        'ANEEL determina comunicação prévia e padrões abertos. Equipamentos não '
        'exclusivos para uso privado devem usar protocolos abertos (OCPP) para '
        'controle remoto. O sistema utiliza OCPP 1.6J como protocolo padrão, '
        'garantindo interoperabilidade com carregadores de diferentes fabricantes.'
    )

    # --- 8. CONCLUSAO ---
    pdf.add_page()
    pdf.titulo_secao(8, 'Conclusão')
    pdf.corpo(
        'O EV ChargeOps demonstra como transformar a infraestrutura de recarga '
        'condominial de um simples hardware isolado em uma plataforma inteligente '
        'de gestão compartilhada. A solução atende integralmente aos três '
        'requisitos do Challenge para turmas online:'
    )
    pdf.corpo(
        '1) Sessões estruturadas por unidade habitacional com autenticação RFID;\n'
        '2) Rateio automático com faturamento individual baseado em consumo real;\n'
        '3) Gerenciamento inteligente com Motor de IA, dashboard gerencial e '
        'Síndico Virtual conversacional.'
    )
    pdf.corpo(
        'O Motor de IA não é um acessório, mas a infraestrutura lógica central que '
        'interpreta dados brutos em inteligência acionável. As integrações com '
        'APIs externas (Open Charge Map, Google Places) expandem o ecossistema '
        'além das fronteiras do condomínio.'
    )
    pdf.corpo(
        'A base regulatória da ANEEL valida o modelo de negócio, permitindo '
        'precificação dinâmica e exigindo interoperabilidade via OCPP - ambos '
        'implementados nesta solução.'
    )

    # Salvar
    if save:
        caminho_pdf = os.path.join(pasta_saida, 'relatorio_ev_chargeops.pdf')
        pdf.output(caminho_pdf)
        print(f"  Relatorio PDF salvo: {caminho_pdf}")
    return pdf


# ============================================================================
# 10. GERADOR DE APRESENTACAO PPT
# ============================================================================

def gerar_apresentacao_ppt(plataforma: EVChargeOps, pasta_saida: str, save=True):
    """Gera apresentacao PowerPoint do EV ChargeOps."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cores do tema
    VERMELHO = RGBColor(230, 57, 70)
    AZUL_ESCURO = RGBColor(29, 53, 87)
    AZUL_MEDIO = RGBColor(69, 123, 157)
    CINZA = RGBColor(100, 100, 100)
    BRANCO = RGBColor(255, 255, 255)
    PRETO = RGBColor(40, 40, 40)

    def add_slide_limpo():
        """Adiciona slide em branco."""
        layout = prs.slide_layouts[6]  # Blank
        return prs.slides.add_slide(layout)

    def add_titulo_slide(slide, titulo, subtitulo=None, cor_titulo=AZUL_ESCURO):
        """Adiciona titulo e subtitulo ao slide."""
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5),
                                          Inches(11.5), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = cor_titulo

        if subtitulo:
            p2 = tf.add_paragraph()
            p2.text = subtitulo
            p2.font.size = Pt(18)
            p2.font.color.rgb = CINZA
            p2.space_before = Pt(8)

        # Linha vermelha decorativa
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3), Pt(4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = VERMELHO
        shape.line.fill.background()

    def add_corpo_texto(slide, textos, y_start=2.2):
        """Adiciona lista de textos ao slide."""
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y_start),
                                          Inches(11.5), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, (titulo_item, desc) in enumerate(textos):
            if i > 0:
                p_space = tf.add_paragraph()
                p_space.space_before = Pt(6)

            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            run = p.add_run()
            run.text = titulo_item
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = AZUL_ESCURO

            if desc:
                p2 = tf.add_paragraph()
                p2.text = desc
                p2.font.size = Pt(16)
                p2.font.color.rgb = PRETO
                p2.space_before = Pt(2)

    def add_box_destaque(slide, texto, x, y, w, h, cor_fundo=AZUL_ESCURO):
        """Adiciona caixa de destaque com texto."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor_fundo
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = texto
        p.font.size = Pt(14)
        p.font.color.rgb = BRANCO
        p.font.bold = True

    # =====================================================================
    # SLIDE 1: CAPA
    # =====================================================================
    slide = add_slide_limpo()

    # Fundo
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_ESCURO
    bg.line.fill.background()

    # Titulo
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5),
                                      Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EV ChargeOps"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Plataforma de Gestão Compartilhada de Recarga"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER

    # Linha vermelha
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(4), Inches(5), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = VERMELHO
    shape.line.fill.background()

    # Info
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5),
                                       Inches(11), Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for texto in ["Challenge FIAP + GoodWe 2026",
                   "Ciências da Computação - 1º Ano Online",
                   "",
                   "Marcelo Bastianello Baldin | RM568746"]:
        p = tf2.add_paragraph()
        p.text = texto
        p.font.size = Pt(18) if texto else Pt(10)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER

    # =====================================================================
    # SLIDE 2: O PROBLEMA
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "O Problema Central")

    textos = [
        ("Crescimento de EVs em condomínios",
         "O aumento acelerado de veículos elétricos exige regras claras, "
         "rateio justo e análise de consumo contínua."),
        ("Sem gestão = conflito",
         "Quem paga pela energia? Como dividir o uso dos carregadores? "
         "Como evitar sobrecarga na rede elétrica?"),
        ("Solução: EV ChargeOps",
         "Transforma cada sessão de recarga em dados estruturados e "
         "inteligência acionável para gestão do condomínio.")
    ]
    add_corpo_texto(slide, textos)

    # =====================================================================
    # SLIDE 3: REQUISITOS ATENDIDOS
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Requisitos do Challenge (Online)",
                    "Transformar solução residencial em coletiva")

    add_box_destaque(slide, "1. Estruturação de sessões de recarga\n"
                    "por usuário ou unidade habitacional",
                    0.8, 2.5, 3.6, 2.0)
    add_box_destaque(slide, "2. Processamento de consumo individual\n"
                    "para cobrança automática da unidade",
                    4.8, 2.5, 3.6, 2.0, AZUL_MEDIO)
    add_box_destaque(slide, "3. Gerenciamento inteligente com\n"
                    "interface usuário + gestão condomínio",
                    8.8, 2.5, 3.6, 2.0, VERMELHO)

    # =====================================================================
    # SLIDE 4: ARQUITETURA
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Arquitetura Híbrida de 3 Camadas")

    add_box_destaque(slide, "CAMADA DIGITAL\nRegras de Controle | Motor de IA | APIs",
                    2, 2.3, 9, 1.3, AZUL_ESCURO)
    add_box_destaque(slide, "CAMADA DE CONECTIVIDADE\nRede / Internet (LAN + Wi-Fi + OCPP)",
                    2, 3.8, 9, 1.3, AZUL_MEDIO)
    add_box_destaque(slide, "CAMADA FÍSICA\nEV Charger GoodWe HCA G2 | MODBUS | RFID",
                    2, 5.3, 9, 1.3, RGBColor(80, 80, 80))

    # =====================================================================
    # SLIDE 5: 4 PILARES
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Os 4 Pilares do EV ChargeOps")

    pilares = [
        ("Gestão de Uso", "Sessões por usuário/unidade\nAutenticação RFID", 0.5, AZUL_ESCURO),
        ("Gerenciamento de Recarga", "Regras compartilhadas\nHorários otimizados", 3.7, AZUL_MEDIO),
        ("Rateio e Faturamento", "Consumo individual (kWh)\nCobrança automática", 6.9, VERMELHO),
        ("Visualização e Interação", "Painéis gerenciais\nInsights para moradores", 10.1, RGBColor(42, 157, 143))
    ]

    for titulo, desc, x, cor in pilares:
        add_box_destaque(slide, f"{titulo}\n\n{desc}", x, 2.5, 2.8, 3.5, cor)

    # =====================================================================
    # SLIDE 6: MOTOR DE IA
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Motor de IA",
                    "A IA como infraestrutura lógica - não uma feature extra")

    ia_items = [
        ("Interpretação", "Decodificação de eventos de sessão\nem tempo real",
         1, AZUL_ESCURO),
        ("Preditividade", "Previsão de demanda e\nnecessidade de expansão",
         4.3, AZUL_MEDIO),
        ("Precificação", "Tarifa dinâmica baseada\nno horário (ANEEL)",
         7.6, VERMELHO),
        ("Conversação", "Síndico Virtual (NLP)\nDados -> orientações simples",
         10.9, RGBColor(42, 157, 143))
    ]

    for titulo, desc, x, cor in ia_items:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5),
            Inches(2.8), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = titulo
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    # =====================================================================
    # SLIDE 7: INTEGRACOES
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Integrações do Ecossistema")

    integ = [
        ("GoodWe API (Simulada)", "Dados operacionais dos carregadores HCA G2.\n"
         "OCPP 1.6J: Start/Stop Transaction, MeterValues.\n"
         'JSON: {"status":"active","power":"11kW"}'),
        ("Open Charge Map API", "Mapa público de eletropostos.\n"
         "Busca pontos próximos por GPS.\n"
         "Cruzamento infraestrutura privada x pública."),
        ("Google Places API", "Localização de pontos de recarga.\n"
         "Rating, fotos, horário de funcionamento.\n"
         "Alternativas quando carregadores ocupados."),
        ("Síndico Virtual (NLP)", "Agente conversacional inteligente.\n"
         "Traduz dados técnicos em linguagem simples.\n"
         "Consumo, faturas, horários, anomalias.")
    ]

    for i, (titulo, desc) in enumerate(integ):
        x = 0.5 + i * 3.15
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3),
            Inches(2.9), Inches(4.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = AZUL_ESCURO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = PRETO
        p2.space_before = Pt(10)

    # =====================================================================
    # SLIDE 8: PRECIFICACAO DINAMICA
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Precificação Dinâmica (ANEEL)",
                    "Tarifas calculadas automaticamente pelo Motor de IA")

    tarifas = [
        ("Fora Ponta", "22h - 17h + Fins de semana", "R$ 0,85/kWh",
         "Tarifa base", RGBColor(42, 157, 143)),
        ("Intermediária", "17h-18h e 21h-22h", "R$ 1,02/kWh",
         "+20%", RGBColor(244, 162, 97)),
        ("Ponta", "18h-21h (dias úteis)", "R$ 1,28/kWh",
         "+50%", VERMELHO)
    ]

    for i, (nome, horario, preco, perc, cor) in enumerate(tarifas):
        x = 1.5 + i * 3.5
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5),
            Inches(3.2), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        for j, (txt, sz, bold) in enumerate([
            (nome, 24, True), (horario, 14, False),
            (preco, 28, True), (perc, 16, False)
        ]):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(sz)
            run.font.bold = bold
            run.font.color.rgb = BRANCO
            p.space_before = Pt(8 if j > 0 else 0)

    # =====================================================================
    # SLIDE 9: RESULTADOS - GRAFICOS
    # =====================================================================
    graficos_slides = [
        ('consumo_diario.png', 'Consumo Diário de Energia (kWh)'),
        ('consumo_por_unidade.png', 'Consumo por Unidade Habitacional'),
        ('distribuicao_horario.png', 'Distribuição de Sessões por Horário'),
        ('distribuicao_tarifa.png', 'Distribuição de Custos por Tarifa'),
        ('custo_acumulado.png', 'Evolução Acumulada de Custos'),
    ]

    for arquivo, titulo in graficos_slides:
        caminho = os.path.join(pasta_saida, arquivo)
        if os.path.exists(caminho):
            slide = add_slide_limpo()
            add_titulo_slide(slide, titulo, "Dados simulados - 30 dias")
            slide.shapes.add_picture(caminho, Inches(1.5), Inches(2.2),
                                     Inches(10), Inches(4.5))

    # =====================================================================
    # SLIDE 14: RESULTADOS NUMERICOS
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Resultados da Simulação",
                    "30 dias de operação simulada")

    dashboard = plataforma.dashboard_condominio()
    analise = plataforma.executar_analise_ia()
    previsao = analise["previsao_demanda"]

    metricas = [
        (f'{dashboard["total_sessoes"]}', "Sessões de Recarga"),
        (f'{dashboard["total_kwh"]:.0f} kWh', "Energia Total"),
        (f'R$ {dashboard["total_custo"]:.0f}', "Custo Total"),
        (f'{dashboard["media_kwh_sessao"]:.1f} kWh', "Média/Sessão"),
        (f'{dashboard["unidades_ativas"]}/{len(plataforma.unidades)}', "Unidades Ativas"),
        (f'{previsao["direcao_tendencia"].upper()}', "Tendência"),
    ]

    for i, (valor, label) in enumerate(metricas):
        x = 0.8 + (i % 3) * 4
        y = 2.5 + (i // 3) * 2.5
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(3.5), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [AZUL_ESCURO, AZUL_MEDIO, VERMELHO,
                                      RGBColor(42, 157, 143), RGBColor(244, 162, 97),
                                      RGBColor(38, 70, 83)][i]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = valor
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = BRANCO
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.CENTER

    # =====================================================================
    # SLIDE 15: BASE REGULATORIA
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Base Regulatória")

    reg = [
        ("Exploração Comercial Livre",
         "ANEEL Resolução Normativa 1.000/2021\n\n"
         "Recarga pública e comercial pode operar com\n"
         "preços livremente negociados. Viabiliza o\n"
         "modelo de tarifação dinâmica do EV ChargeOps."),
        ("Interoperabilidade Obrigatória",
         "ANEEL - Padrões Abertos\n\n"
         "Equipamentos de uso compartilhado devem\n"
         "usar protocolos abertos (OCPP) para\n"
         "controle remoto. Justifica a arquitetura OCPP.")
    ]

    for i, (titulo, desc) in enumerate(reg):
        x = 1.5 + i * 5.5
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5),
            Inches(5), Inches(4.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = AZUL_ESCURO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = PRETO
        p2.space_before = Pt(10)

    # =====================================================================
    # SLIDE 16: USO DE IA NO PROJETO
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Uso de IA no Contexto do Software",
                    "A IA estrutura a proposta - não é um penduricalho de interface")

    textos_ia = [
        ("Motor de IA integrado ao núcleo",
         "A IA opera em 4 dimensões (Interpretação, Preditividade, Precificação, "
         "Conversação) como camada lógica central, não como feature isolada."),
        ("Detecção de anomalias por desvio estatístico",
         "Identifica sessões com consumo fora do padrão (>2 desvios) e alerta "
         "sobre veículos abandonados ou equipamento com falha."),
        ("Síndico Virtual com NLP",
         "Agente conversacional que traduz dados brutos em orientações gerenciais "
         "simples. Classificação de intenções por palavras-chave."),
        ("Previsão de demanda e expansão",
         "Média móvel + tendência para projetar consumo futuro e recomendar "
         "instalação de carregadores adicionais antes da sobrecarga.")
    ]
    add_corpo_texto(slide, textos_ia)

    # =====================================================================
    # SLIDE 17: ENTREGAVEIS
    # =====================================================================
    slide = add_slide_limpo()
    add_titulo_slide(slide, "Entregáveis do Projeto")

    entregaveis = [
        ("Arquitetura Funcional",
         "Lógica clara de entradas e saídas de dados (Data Flow).\n"
         "3 camadas: Física + Conectividade + Digital."),
        ("IA Integrada",
         "Motor lógico com 4 funções: interpretar, prever,\n"
         "precificar e conversar. Síndico Virtual NLP."),
        ("Aderência ao Contexto",
         "Solução sob medida para condomínios utilizando\n"
         "o ecossistema GoodWe/FIAP (HCA G2 + OCPP)."),
        ("Visão de Produto Real",
         "Resolve problema regulatório (ANEEL) e operacional\n"
         "(rateio justo) existente hoje em condomínios.")
    ]

    for i, (titulo, desc) in enumerate(entregaveis):
        y = 2.3 + i * 1.25
        # Check mark
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(34, 139, 34)
        shape.line.fill.background()
        tf_check = shape.text_frame
        tf_check.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_check = tf_check.paragraphs[0].add_run()
        run_check.text = "OK"  # Checkmark
        run_check.font.size = Pt(12)
        run_check.font.bold = True
        run_check.font.color.rgb = BRANCO

        # Texto
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y - 0.1),
                                          Inches(10), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo + ": "
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = AZUL_ESCURO
        run2 = p.add_run()
        run2.text = desc.replace('\n', ' ')
        run2.font.size = Pt(15)
        run2.font.color.rgb = PRETO

    # =====================================================================
    # SLIDE 18: ENCERRAMENTO
    # =====================================================================
    slide = add_slide_limpo()

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_ESCURO
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2),
                                      Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "EV ChargeOps"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Gestão inteligente de recarga para condomínios"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(5), Inches(5), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = VERMELHO
    shape.line.fill.background()

    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5),
                                       Inches(11), Inches(1.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.add_paragraph()
    p3.text = "FIAP + GoodWe | Challenge 2026"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(180, 180, 180)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf3.add_paragraph()
    p4.text = "Marcelo Bastianello Baldin | RM568746"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(150, 150, 150)
    p4.alignment = PP_ALIGN.CENTER

    # Salvar
    if save:
        caminho_ppt = os.path.join(pasta_saida, 'apresentacao_ev_chargeops.pptx')
        prs.save(caminho_ppt)
        print(f"  Apresentacao PPT salva: {caminho_ppt}")
    return prs


# ============================================================================
# 11. INTERFACE PRINCIPAL (MENU CLI)
# ============================================================================

def limpar_tela():
    os.system('cls' if os.name == 'nt' else 'clear')


def banner():
    print("=" * 70)
    print("  EV ChargeOps - Plataforma de Gestao Compartilhada de Recarga")
    print("  Challenge FIAP + GoodWe 2026 | Marcelo B. Baldin - RM568746")
    print("=" * 70)


def menu_principal(plataforma: EVChargeOps, pasta_saida: str):
    """Menu interativo principal."""
    while True:
        print("\n" + "-" * 50)
        print("  MENU PRINCIPAL")
        print("-" * 50)
        print("  1. Status dos Carregadores")
        print("  2. Dashboard do Condominio")
        print("  3. Iniciar Sessao de Recarga")
        print("  4. Finalizar Sessao de Recarga")
        print("  5. Relatorio de Rateio Mensal")
        print("  6. Analise do Motor de IA")
        print("  7. Sindico Virtual (Chat)")
        print("  8. Buscar Eletropostos Proximos")
        print("  9. Gerar Graficos")
        print(" 10. Gerar Relatorio PDF")
        print(" 11. Gerar Apresentacao PPT")
        print(" 12. Gerar Tudo (Graficos + PDF + PPT)")
        print("  0. Sair")
        print("-" * 50)

        opcao = input("  Escolha: ").strip()

        if opcao == "1":
            print("\n--- STATUS DOS CARREGADORES ---")
            for s in plataforma.status_carregadores():
                print(f"  [{s['status'].upper():^12}] {s['model']} - {s['localizacao']}")
                print(f"              Potencia: {s['power']} | Temp: {s['temperature_c']}C")

        elif opcao == "2":
            print("\n--- DASHBOARD DO CONDOMINIO ---")
            d = plataforma.dashboard_condominio()
            print(f"  Condominio: {d['condominio']}")
            print(f"  Sessoes totais: {d['total_sessoes']}")
            print(f"  Consumo total: {d['total_kwh']:.1f} kWh")
            print(f"  Custo total: R$ {d['total_custo']:.2f}")
            print(f"  Media/sessao: {d['media_kwh_sessao']:.1f} kWh")
            print(f"  Unidades ativas: {d['unidades_ativas']}/{len(plataforma.unidades)}")
            print(f"  Carregadores disponiveis: "
                  f"{d['carregadores_disponiveis']}/{d['total_carregadores']}")
            print(f"  Hora de pico: {d['hora_pico']:02d}h")
            print(f"  Distribuicao tarifaria: {d['distribuicao_tarifa']}")

        elif opcao == "3":
            print("\n--- INICIAR SESSAO ---")
            print("  Unidades disponiveis:")
            for i, u in enumerate(plataforma.unidades):
                print(f"    {i+1}. {u.bloco}-{u.numero} ({u.proprietario})")
            try:
                idx_u = int(input("  Numero da unidade: ")) - 1
                unidade = plataforma.unidades[idx_u]
            except (ValueError, IndexError):
                print("  Opcao invalida.")
                continue

            print("  Carregadores disponiveis:")
            disp = [(i, c) for i, c in enumerate(plataforma.carregadores)
                    if c.status == "disponivel"]
            if not disp:
                print("  Nenhum carregador disponivel no momento.")
                continue
            for i, c in disp:
                print(f"    {i+1}. {c.modelo} ({c.potencia_kw}kW) - {c.localizacao}")
            try:
                idx_c = int(input("  Numero do carregador: ")) - 1
                carregador = plataforma.carregadores[idx_c]
            except (ValueError, IndexError):
                print("  Opcao invalida.")
                continue

            try:
                sessao = plataforma.gerenciador.iniciar_sessao(
                    unidade, carregador, plataforma.condominio.tarifa_kwh_base)
                print(f"  Sessao {sessao.id} iniciada!")
                print(f"  Tarifa: R$ {sessao.tarifa_aplicada}/kWh ({sessao.tipo_tarifa})")
            except ValueError as e:
                print(f"  Erro: {e}")

        elif opcao == "4":
            print("\n--- FINALIZAR SESSAO ---")
            ativas = [(c, s) for c_id, s in plataforma.gerenciador.sessoes_ativas.items()
                      for c in plataforma.carregadores if c.id == c_id]
            if not ativas:
                print("  Nenhuma sessao ativa.")
                continue
            for i, (c, s) in enumerate(ativas):
                unid = next((u for u in plataforma.unidades if u.id == s.unidade_id), None)
                print(f"    {i+1}. {c.modelo} - Unidade {unid.bloco}-{unid.numero if unid else '?'}")
            try:
                idx = int(input("  Numero: ")) - 1
                c, s = ativas[idx]
                sessao = plataforma.gerenciador.finalizar_sessao(c)
                print(f"  Sessao {sessao.id} finalizada!")
                print(f"  Energia: {sessao.energia_kwh:.2f} kWh")
                print(f"  Custo: R$ {sessao.custo_total:.2f}")
            except (ValueError, IndexError):
                print("  Opcao invalida.")

        elif opcao == "5":
            print("\n--- RATEIO MENSAL ---")
            mes = input("  Mes de referencia (YYYY-MM) [Enter=atual]: ").strip()
            if not mes:
                mes = datetime.now().strftime("%Y-%m")
            rateio = plataforma.faturamento.relatorio_rateio(
                plataforma.unidades, plataforma.gerenciador.sessoes, mes)
            print(f"\n  Mes: {rateio['mes_referencia']}")
            print(f"  Total condominio: R$ {rateio['total_condominio_reais']:.2f}")
            print(f"\n  {'Unidade':<10} {'Bloco':<6} {'Proprietario':<20} "
                  f"{'Sessoes':>8} {'kWh':>10} {'Custo (R$)':>12}")
            print("  " + "-" * 70)
            for num, dados in sorted(rateio['rateio_por_unidade'].items()):
                print(f"  {num:<10} {dados['bloco']:<6} {dados['proprietario']:<20} "
                      f"{dados['sessoes']:>8} {dados['kwh']:>10.2f} "
                      f"{dados['custo']:>12.2f}")

        elif opcao == "6":
            print("\n--- ANALISE DO MOTOR DE IA ---")
            analise = plataforma.executar_analise_ia()

            print("\n  [PREVISAO DE DEMANDA]")
            prev = analise["previsao_demanda"]
            print(f"  Consumo medio diario: {prev['consumo_medio_diario_kwh']:.1f} kWh")
            print(f"  Tendencia: {prev['direcao_tendencia']} ({prev['tendencia_pct']:+.1f}%)")
            print(f"  Previsao mensal: {prev['previsao_mensal_kwh']:.0f} kWh")
            print(f"  Recomendacao: {prev['recomendacao']}")

            print("\n  [ANOMALIAS]")
            if analise["anomalias"]:
                for a in analise["anomalias"][:5]:
                    print(f"  ! {a}")
            else:
                print("  Nenhuma anomalia detectada.")

            print("\n  [ULTIMAS INTERPRETACOES]")
            for interp in analise["interpretacoes"]:
                print(f"  > {interp['interpretacao']}")

        elif opcao == "7":
            print("\n--- SINDICO VIRTUAL ---")
            print("  (Digite 'sair' para voltar ao menu)")
            analise = plataforma.executar_analise_ia()
            dados = analise["dados_sindico"]
            while True:
                pergunta = input("\n  Voce: ").strip()
                if pergunta.lower() in ["sair", "exit", "voltar"]:
                    break
                resposta = MotorIA.sindico_virtual(pergunta, dados)
                print(f"\n  Sindico Virtual: {resposta}")

        elif opcao == "8":
            print("\n--- ELETROPOSTOS PROXIMOS ---")
            resultado = plataforma.buscar_eletropostos_proximos()
            print("\n  [Open Charge Map]")
            for ep in resultado["open_charge_map"]:
                print(f"  > {ep['nome']} ({ep['distancia_km']}km)")
                print(f"    {ep['endereco']}")
                print(f"    Potencia: {ep['potencia_max_kw']}kW | "
                      f"Preco: R$ {ep['preco_kwh']}/kWh | Status: {ep['status']}")

            print("\n  [Google Places]")
            for gp in resultado["google_places"]:
                print(f"  > {gp['nome']} (Rating: {gp['rating']})")
                print(f"    {gp['endereco']}")

        elif opcao == "9":
            print("\n  Gerando graficos...")
            gerar_graficos(plataforma, pasta_saida)

        elif opcao == "10":
            print("\n  Gerando relatorio PDF...")
            gerar_relatorio_pdf(plataforma, pasta_saida)

        elif opcao == "11":
            print("\n  Gerando apresentacao PPT...")
            gerar_apresentacao_ppt(plataforma, pasta_saida)

        elif opcao == "12":
            print("\n  Gerando todos os entregaveis...")
            print("\n  [1/3] Graficos:")
            gerar_graficos(plataforma, pasta_saida)
            print("\n  [2/3] Relatorio PDF:")
            gerar_relatorio_pdf(plataforma, pasta_saida)
            print("\n  [3/3] Apresentacao PPT:")
            gerar_apresentacao_ppt(plataforma, pasta_saida)
            print("\n  Todos os entregaveis gerados com sucesso!")

        elif opcao == "0":
            print("\n  Obrigado por usar o EV ChargeOps!")
            break

        else:
            print("  Opcao invalida. Tente novamente.")


# ============================================================================
# 12. EXECUCAO AUTOMATICA (GERAR TUDO SEM MENU)
# ============================================================================

def executar_automatico(pasta_saida: str):
    """Executa geracao completa automatica (sem interacao do usuario)."""
    banner()

    print("\n[1/5] Configurando ambiente de demonstracao...")
    plataforma = EVChargeOps()
    plataforma.setup_demo()

    print("\n[2/5] Gerando historico simulado (30 dias)...")
    plataforma.gerar_historico_simulado(dias=30)

    print("\n[3/5] Gerando graficos de analise...")
    gerar_graficos(plataforma, pasta_saida)

    print("\n[4/5] Gerando relatorio PDF...")
    gerar_relatorio_pdf(plataforma, pasta_saida)

    print("\n[5/5] Gerando apresentacao PowerPoint...")
    gerar_apresentacao_ppt(plataforma, pasta_saida)

    # Exibir resumo
    print("\n" + "=" * 70)
    print("  EXECUCAO CONCLUIDA COM SUCESSO!")
    print("=" * 70)
    dashboard = plataforma.dashboard_condominio()
    print(f"  Condominio: {dashboard['condominio']}")
    print(f"  Sessoes simuladas: {dashboard['total_sessoes']}")
    print(f"  Consumo total: {dashboard['total_kwh']:.1f} kWh")
    print(f"  Custo total: R$ {dashboard['total_custo']:.2f}")
    print(f"\n  Arquivos gerados em: {pasta_saida}/")
    print(f"  - 5 graficos PNG")
    print(f"  - relatorio_ev_chargeops.pdf")
    print(f"  - apresentacao_ev_chargeops.pptx")

    return plataforma


# ============================================================================
# MAIN
# ============================================================================

if __name__ == "__main__":
    # Pasta de saida = mesma pasta do script
    PASTA_SAIDA = os.path.dirname(os.path.abspath(__file__))

    if len(sys.argv) > 1 and sys.argv[1] == "--auto":
        # Modo automatico: gera tudo sem menu
        executar_automatico(PASTA_SAIDA)
    else:
        # Modo interativo: menu CLI
        banner()
        print("\n  Inicializando plataforma...")
        plataforma = EVChargeOps()
        plataforma.setup_demo()
        print("\n  Gerando historico simulado...")
        plataforma.gerar_historico_simulado(dias=30)
        print("\n  Plataforma pronta!")
        menu_principal(plataforma, PASTA_SAIDA)
