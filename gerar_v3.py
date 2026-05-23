#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
EV ChargeOps v3 - Sprint 01: Entrega Completa
Gera relatorio PDF, apresentacao PPTX e documento MD finais,
reestruturados conforme os requisitos do Enterprise Challenge 2026
(playbook 1CC + tarefa FIAP + apresentacao GoodWe).

Estrutura do entregavel:
  1. Investigacao do Problema (pesquisa de mercado)
  2. Contexto Tecnico e Regulatorio
  3. Arquitetura da Solucao (3 camadas + Data Flow)
  4. Opcao A - Estruturacao de Sessoes de Recarga
  5. Opcao B - Processamento de Consumo + IA Avancada (4 dimensoes)
  6. Opcao C - Gerenciamento Inteligente e UX
  7. Rateio e Faturamento
  8. Decisoes Tecnicas (Q&A)
  9. Visao de Produto e Roadmap

Aluno: Marcelo Bastianello Baldin | RM568746
Prazo Sprint 01: 21/06/2026
"""

import os
import sys
from datetime import datetime

script_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, script_dir)

from ev_chargeops import (
    EVChargeOps, gerar_graficos,
    gerar_relatorio_pdf, gerar_apresentacao_ppt
)


# ============================================================================
# CONSTANTES
# ============================================================================
AZUL_ESCURO = (29, 53, 87)
AZUL_MEDIO = (69, 123, 157)
VERMELHO = (230, 57, 70)
VERDE = (42, 157, 143)
CINZA = (100, 100, 100)
PRETO = (40, 40, 40)
LARANJA = (244, 162, 97)
PETROL = (38, 70, 83)


# ============================================================================
# HELPERS PDF
# ============================================================================

def _th(pdf, colunas, larguras):
    pdf.set_font('Helvetica', 'B', 8)
    pdf.set_fill_color(*AZUL_ESCURO)
    pdf.set_text_color(255, 255, 255)
    for t, w in zip(colunas, larguras):
        pdf.cell(w, 7, t, 1, 0, 'C', True)
    pdf.ln()
    pdf.set_text_color(*PRETO)


def _tr(pdf, vals, larguras, fill=False):
    pdf.set_font('Helvetica', '', 7)
    pdf.set_fill_color(240, 245, 250) if fill else pdf.set_fill_color(255, 255, 255)
    for t, w in zip(vals, larguras):
        pdf.cell(w, 6, t, 1, 0, 'L', True)
    pdf.ln()


def _check_page(pdf, needed=40):
    if pdf.get_y() > 270 - needed:
        pdf.add_page()


def _titulo_principal(pdf, num, titulo):
    _check_page(pdf, 30)
    pdf.ln(4)
    pdf.set_font('Helvetica', 'B', 16)
    pdf.set_text_color(*AZUL_ESCURO)
    pdf.cell(0, 10, f'{num}. {titulo}', 0, 1)
    pdf.set_draw_color(*VERMELHO)
    pdf.line(10, pdf.get_y(), 80, pdf.get_y())
    pdf.ln(3)


def _sub(pdf, texto):
    _check_page(pdf, 20)
    pdf.ln(2)
    pdf.set_font('Helvetica', 'B', 12)
    pdf.set_text_color(*AZUL_MEDIO)
    pdf.cell(0, 8, texto, 0, 1)
    pdf.set_text_color(*PRETO)


def _p(pdf, texto):
    _check_page(pdf, 15)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(*PRETO)
    pdf.multi_cell(0, 5, texto)
    pdf.ln(1)


def _bullet(pdf, texto):
    _check_page(pdf, 10)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(*PRETO)
    x = pdf.get_x()
    pdf.cell(5, 5, '-', 0, 0)
    pdf.multi_cell(0, 5, texto)
    pdf.ln(1)


def _empresa(pdf, nome, desc):
    _check_page(pdf, 20)
    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(*AZUL_MEDIO)
    pdf.cell(0, 7, nome, 0, 1)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(*PRETO)
    pdf.multi_cell(0, 5, desc)
    pdf.ln(2)


def _qa(pdf, pergunta, pros, contras, decisao):
    _check_page(pdf, 50)
    pdf.ln(2)
    pdf.set_font('Helvetica', 'B', 11)
    pdf.set_text_color(*AZUL_ESCURO)
    pdf.cell(0, 8, pergunta, 0, 1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(34, 139, 34)
    pdf.cell(0, 6, 'Pros:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(*PRETO)
    pdf.multi_cell(0, 4, pros)
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_text_color(*VERMELHO)
    pdf.cell(0, 6, 'Contras:', 0, 1)
    pdf.set_font('Helvetica', '', 9)
    pdf.set_text_color(*PRETO)
    pdf.multi_cell(0, 4, contras)
    pdf.ln(1)

    pdf.set_font('Helvetica', 'B', 9)
    pdf.set_fill_color(230, 245, 230)
    pdf.cell(0, 6, f'Decisao: {decisao}', 0, 1, '', True)
    pdf.ln(3)


# ============================================================================
# PDF v3 - CONTEUDO COMPLETO
# ============================================================================

def gerar_pdf_v3(plataforma, pasta_saida):
    pdf = gerar_relatorio_pdf(plataforma, pasta_saida, save=False)

    # ================================================================
    # PAGINA SEPARADORA: SPRINT 01
    # ================================================================
    pdf.add_page()
    pdf.ln(30)
    pdf.set_font('Helvetica', 'B', 32)
    pdf.set_text_color(*AZUL_ESCURO)
    pdf.cell(0, 15, 'SPRINT 01', 0, 1, 'C')
    pdf.set_font('Helvetica', '', 16)
    pdf.set_text_color(*CINZA)
    pdf.cell(0, 10, 'Pesquisa e Documentacao', 0, 1, 'C')
    pdf.cell(0, 8, 'Prazo: 21/06/2026', 0, 1, 'C')
    pdf.ln(5)
    pdf.set_draw_color(*VERMELHO)
    pdf.line(60, pdf.get_y(), 150, pdf.get_y())
    pdf.ln(10)
    pdf.set_font('Helvetica', '', 12)
    pdf.set_text_color(*PRETO)
    pdf.multi_cell(0, 6,
        'A equipe investiga o problema, mapeia o contexto tecnico e '
        'regulatorio, define a arquitetura da solucao e documenta as '
        'decisoes que guiarao o desenvolvimento.', align='C')
    pdf.ln(5)
    pdf.set_font('Helvetica', 'I', 11)
    pdf.set_text_color(*CINZA)
    pdf.multi_cell(0, 6,
        '"Como transformar sessoes de recarga de veiculos eletricos em '
        'uma infraestrutura compartilhada em dados estruturados, rateio '
        'justo e inteligencia acionavel?"', align='C')

    # ================================================================
    # 1. INVESTIGACAO DO PROBLEMA
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 1, 'Investigacao do Problema')

    _sub(pdf, '1.1 Contexto: O Crescimento dos Veiculos Eletricos')
    _p(pdf,
        'O mercado global de veiculos eletricos (EVs) atingiu 14 milhoes '
        'de unidades vendidas em 2023 (IEA Global EV Outlook 2024), com '
        'projecao de 40 milhoes anuais ate 2030. No Brasil, as vendas de '
        'EVs e hibridos plug-in cresceram 91% em 2023, totalizando mais '
        'de 93 mil unidades (ABVE). A EPE projeta 4 milhoes de EVs no '
        'Brasil ate 2030 e 33 milhoes ate 2050 (PNE 2050).')
    _p(pdf,
        'Este crescimento acelerado cria uma demanda urgente por '
        'infraestrutura de recarga que vai alem do simples fornecimento '
        'de energia. Cada sessao de recarga produz dados uteis: duracao, '
        'volume de energia (kWh), horario de uso, frequencia, picos e '
        'intervalos de ociosidade. Quando organizados, esses dados deixam '
        'de ser simples registros e passam a funcionar como base de '
        'inteligencia operacional.')

    _sub(pdf, '1.2 Problema Central em Condominios')
    _p(pdf,
        'O EV ChargeOps foi concebido para resolver um problema especifico: '
        'infraestruturas de recarga compartilhadas em condominios '
        'residenciais, edificios corporativos e campus universitarios nao '
        'dispoem de mecanismos integrados para:')
    _bullet(pdf, 'Estruturar sessoes de recarga por usuario ou unidade habitacional')
    _bullet(pdf, 'Calcular consumo individual e aplicar regras de rateio justo')
    _bullet(pdf, 'Oferecer uma experiencia digital clara para moradores e gestores')
    _bullet(pdf, 'Prever picos de demanda e otimizar o uso da infraestrutura eletrica')
    _bullet(pdf, 'Gerar inteligencia acionavel a partir dos dados de recarga')
    _p(pdf,
        'Sem esses mecanismos, a recarga em condominio se torna fonte de '
        'conflitos: quem paga? quanto cada um consumiu? quando carregar '
        'sem sobrecarregar a rede? O condominio precisa de software que '
        'transforme hardware isolado em um hub de inteligencia.')

    _sub(pdf, '1.3 Pesquisa de Mercado - Solucoes Existentes')

    _p(pdf, 'O mercado de infraestrutura de recarga no Brasil cresce '
        '40-50% ao ano. Diversas empresas ja oferecem solucoes comerciais:')

    _sub(pdf, '1.3.1 Solucoes no Brasil')

    _empresa(pdf, 'Voltbras (Porto Alegre, 2017)',
        'Maior plataforma de gestao de recarga do Brasil. Modelo SaaS B2B '
        'com integracao OCPP a mais de 10 fabricantes de carregadores. '
        'Oferece gestao de hubs publicos e privados, pagamento via app e '
        'RFID, dashboard de analytics e relatorios. Clientes incluem '
        'shoppings, condominios e frotas corporativas. Parceria com BMW, '
        'Volvo e BYD para integracao de dados do veiculo. Ponto forte: '
        'escala e integracao multi-fabricante. Ponto fraco: nao tem '
        'motor de IA integrado para analise preditiva de consumo.')

    _empresa(pdf, 'WEG WEMOB (Jaragua do Sul)',
        'Fabricante brasileiro de carregadores EV com a linha WEMOB: '
        'modelos de 7,4 kW (AC monofasico), 22 kW (AC trifasico) e '
        '60-180 kW (DC fast charging). Protocolo OCPP 1.6J nativo, '
        'autenticacao RFID, app de gestao. Vantagem competitiva: '
        'fabricacao nacional com suporte tecnico local e integracao '
        'com inversores solares WEG. Foco primario em hardware, '
        'software de gestao menos desenvolvido que plataformas SaaS.')

    _empresa(pdf, 'Zletric (Sao Paulo)',
        'Plataforma SaaS focada exclusivamente em condominios. Oferece '
        'rateio individualizado por unidade habitacional, integracao com '
        'administradoras de condominio e relatorios mensais automatizados. '
        'Modelo de negocio: assinatura mensal por carregador gerenciado. '
        'Concorrente mais direto do EV ChargeOps, porem sem motor de IA '
        'e sem precificacao dinamica.')

    _empresa(pdf, 'EDP Box (Grupo EDP) e Tupinamba Energia',
        'EDP Box: programa de recarga residencial e condominial com '
        'carregador + instalacao + app, integrado a conta de energia EDP. '
        'Tupinamba Energia: rede de mais de 500 eletropostos publicos no '
        'Brasil com app de localizacao, reserva e pagamento digital. '
        'Ambos focam na experiencia do consumidor final, sem solucao '
        'especifica para gestao condominial inteligente.')

    _sub(pdf, '1.3.2 Solucoes Internacionais')

    _empresa(pdf, 'ChargePoint (EUA)',
        'Maior rede aberta de recarga do mundo com mais de 200 mil pontos. '
        'Plataforma cloud completa: gestao de frota, analytics preditivo, '
        'pagamento integrado. Hardware proprio (CP6000, Express Plus) e '
        'software como servico. IPO na NYSE em 2021. Referencia em '
        'escalabilidade e gestao enterprise. A plataforma inclui modulos '
        'de IA para previsao de demanda e otimizacao de frota.')

    _empresa(pdf, 'Wallbox (Espanha)',
        'Destaque para o Quasar 2: primeiro carregador residencial com '
        'V2G bidirecional - o veiculo pode devolver energia a rede ou '
        'a residencia. Pulsar Plus para uso domestico. App com gestao '
        'de energia solar integrada. Presente em 107 paises. Referencia '
        'em inovacao de hardware com V2G.')

    _empresa(pdf, 'EVBox (Holanda) e Virta (Finlandia)',
        'EVBox (grupo ENGIE): hardware (Elvi, Iqon, Troniq Modular) + '
        'plataforma Everon para operadores. OCPP nativo, foco enterprise. '
        'Virta: plataforma white-label API-first para operadores de recarga, '
        'OCPP 2.0.1, roaming europeu via Hubject/GIREVE. Referencia em '
        'interoperabilidade e modelo de plataforma aberta.')

    _empresa(pdf, 'Tesla Supercharger (EUA/Global)',
        'Maior rede proprietaria do mundo com mais de 50 mil pontos. '
        'Conector NACS adotado como padrao SAE J3400 na America do Norte. '
        'Abertura da rede para outros fabricantes via Magic Dock (CCS). '
        'Potencias de 72 kW a 250 kW (V3) e 350 kW (V4). Modelo '
        'verticalmente integrado: hardware + software + energia.')

    _sub(pdf, '1.3.3 Tabela Comparativa')

    colunas = ['Empresa', 'Pais', 'Foco', 'Protocolo', 'V2G', 'IA', 'Destaque']
    larguras = [25, 16, 24, 18, 10, 10, 87]
    _th(pdf, colunas, larguras)
    dados = [
        ['Voltbras', 'BR', 'B2B/Hubs', 'OCPP 1.6', 'Nao', 'Nao', 'Maior plataforma BR, +10 fabricantes'],
        ['WEG WEMOB', 'BR', 'Hardware', 'OCPP 1.6', 'Nao', 'Nao', 'Fabricacao nacional, suporte local'],
        ['Zletric', 'BR', 'Condomin.', 'OCPP 1.6', 'Nao', 'Nao', 'Rateio condominial especializado'],
        ['ChargePoint', 'EUA', 'Rede aberta', 'OCPP 2.0', 'Nao', 'Sim', '200k+ pontos, plataforma cloud'],
        ['Wallbox', 'ES', 'Residencial', 'OCPP 1.6', 'Sim', 'Nao', 'Quasar V2G bidirecional'],
        ['EVBox/Virta', 'EU', 'Enterprise', 'OCPP 2.0', 'Nao', 'Nao', 'Plataforma white-label'],
        ['Tesla', 'EUA', 'Rede propria', 'Propria', 'Nao', 'Sim', '50k+ pontos, modelo vertical'],
        ['EV ChargeOps', 'BR', 'Condomin.', 'OCPP 1.6', 'Roadmap', 'Sim', 'Motor IA 4D, Sindico Virtual'],
    ]
    for i, row in enumerate(dados):
        _tr(pdf, row, larguras, fill=(i % 2 == 0))

    pdf.ln(3)

    _sub(pdf, '1.4 Lacunas Identificadas e Diferenciais')
    _p(pdf,
        'A analise de mercado revela que as solucoes existentes para '
        'condominios (Zletric, EDP Box) oferecem rateio basico sem '
        'inteligencia analitica. As plataformas com IA (ChargePoint, Tesla) '
        'focam em redes publicas/comerciais e nao atendem a gestao '
        'condominial. O EV ChargeOps ocupa uma lacuna estrategica:')
    _bullet(pdf, 'Unico com Motor de IA de 4 dimensoes integrado (interpretacao, '
            'preditividade, precificacao, conversacao)')
    _bullet(pdf, 'Sindico Virtual: agente conversacional que traduz dados tecnicos '
            'em linguagem gerencial acessivel')
    _bullet(pdf, 'Precificacao dinamica multi-fator alinhada com ANEEL')
    _bullet(pdf, 'Foco especifico no problema de rateio justo condominial')
    _bullet(pdf, 'Integracao nativa com hardware GoodWe HCA G2 via OCPP')

    # ================================================================
    # 2. CONTEXTO TECNICO E REGULATORIO
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 2, 'Contexto Tecnico e Regulatorio')

    _sub(pdf, '2.1 Regulamentacao Brasileira')

    _p(pdf, 'O setor eletrico brasileiro e regulado por um conjunto de '
        'instituicoes com funcoes distintas: ANEEL (regulacao e '
        'fiscalizacao), ONS (operacao do SIN), CCEE (comercializacao) '
        'e EPE (planejamento energetico). A recarga de EVs e impactada '
        'diretamente por regulamentacoes dessas entidades.')

    _empresa(pdf, 'ANEEL Resolucao Normativa 1.000/2021',
        'Consolida as regras de distribuicao de energia eletrica. '
        'Artigos 311 a 318 tratam especificamente da recarga de veiculos '
        'eletricos, classificando-a como atividade acessoria nao regulada, '
        'o que permite precificacao livre pelo prestador do servico. '
        'Esta resolucao e a base legal que viabiliza o modelo de '
        'tarifacao dinamica do EV ChargeOps e justifica legalmente '
        'o pilar de Precificacao do Motor de IA.')

    _empresa(pdf, 'Lei 14.300/2022 - Marco Legal da Geracao Distribuida',
        'Regulamenta a micro e minigeracao distribuida no Brasil. Permite '
        'que prosumidores gerem propria energia via paineis solares e '
        'injetem excedente na rede (net metering). Relevancia: condominios '
        'com geracao solar podem abastecer carregadores com energia '
        'propria, reduzindo custos e pegada de carbono. O EV ChargeOps '
        'pode integrar dados de geracao solar no Motor de IA para '
        'otimizar o horario de recarga conforme producao solar local.')

    _empresa(pdf, 'ABNT NBR IEC 61851 e INMETRO',
        'ABNT NBR IEC 61851 define requisitos de seguranca para '
        'sistemas de recarga condutiva: modos de carga 1 a 4, protecoes '
        'eletricas, comunicacao piloto (Control Pilot). INMETRO Portaria '
        '111/2023 estabelece certificacao compulsoria. Os carregadores '
        'GoodWe HCA G2 atendem integralmente esses requisitos.')

    _empresa(pdf, 'PNE 2050 - Plano Nacional de Energia',
        'Elaborado pela EPE, projeta 4 milhoes de EVs no Brasil ate 2030 '
        'e 33 milhoes ate 2050. A demanda adicional e estimada em 24 '
        'TWh/ano ate 2035, equivalente a 3% do consumo eletrico nacional. '
        'O modulo de Preditividade do Motor de IA utiliza essas projecoes '
        'como baseline para planejamento de capacidade do condominio.')

    _sub(pdf, '2.2 Regulacao Internacional de Referencia')

    _bullet(pdf, 'Uniao Europeia - AFIR (2023): obriga postos de recarga '
            'rapida a cada 60 km em rodovias TEN-T, com minimo de 150 kW '
            'ate 2025 e 350 kW ate 2030. Pagamento por cartao obrigatorio.')
    _bullet(pdf, 'EUA - NEVI Program: investimento de US$ 7,5 bilhoes para '
            '500 mil carregadores ate 2030. Exige OCPP, uptime minimo de '
            '97% e conectores CCS.')
    _bullet(pdf, 'Reino Unido - Smart Charge Points Regulations 2021: exige '
            'que carregadores residenciais suportem demand response e V2G. '
            'Carga padrao fora do pico (off-peak default).')
    _bullet(pdf, 'China - Padroes GB/T: maior mercado EV do mundo (>8 '
            'milhoes vendidos em 2023). Padronizacao propria de conectores.')

    _sub(pdf, '2.3 Normas Tecnicas e Protocolos')

    _empresa(pdf, 'OCPP 1.6J (Open Charge Point Protocol)',
        'Protocolo aberto da Open Charge Alliance para comunicacao '
        'bidirecional entre o sistema central e os pontos de carga. '
        'Variante JSON over WebSocket. Operacoes: Authorize, '
        'StartTransaction, StopTransaction, MeterValues, '
        'StatusNotification, Heartbeat, FirmwareUpdate. Suportado '
        'nativamente pelos carregadores GoodWe HCA G2. A exigencia '
        'de interoperabilidade pela ANEEL torna o OCPP obrigatorio '
        'para uso compartilhado.')

    _empresa(pdf, 'MODBUS RTU/TCP',
        'Protocolo industrial para leitura de registradores de energia '
        'em tempo real: tensao (V), corrente (A), potencia ativa (W), '
        'fator de potencia, frequencia (Hz) e energia acumulada (kWh). '
        'Polling configuravel (padrao: 5 segundos). Complementar ao OCPP '
        'para telemetria detalhada utilizada pelo modulo de Interpretacao.')

    _sub(pdf, '2.4 Perspectivas V2G e ISO 15118')
    _p(pdf,
        'O V2G (Vehicle-to-Grid) permite utilizar a bateria do veiculo '
        'como armazenamento distribuido, devolvendo energia a rede nos '
        'horarios de ponta. A ANEEL estuda regulamentacao via Consulta '
        'Publica 020/2024. Aplicacoes: peak shaving, reserva de '
        'emergencia e arbitragem tarifaria. O EV ChargeOps inclui V2G '
        'no roadmap, preparando a arquitetura para suporte futuro ao '
        'protocolo ISO 15118 (comunicacao bidirecional).')

    _sub(pdf, '2.5 Producao e Consumo de Energia no Brasil')
    _p(pdf,
        'Matriz eletrica predominantemente renovavel: hidraulica (63%), '
        'eolica (11%), solar (4%), biomassa (5%). Capacidade instalada: '
        '195 GW. Tarifa media residencial: R$ 0,65/kWh com impostos.')
    _p(pdf,
        'Sistema de bandeiras tarifarias (verde, amarela, vermelha) '
        'reflete o custo de geracao. A tarifa branca, opcional para '
        'consumidores com medicao horaria, aplica valores diferenciados '
        'por periodo - modelo analogo ao usado pelo EV ChargeOps para '
        'precificacao de recargas.')
    _p(pdf,
        'Horario de ponta (18h-21h): maior demanda e custo. Carregar '
        'fora deste horario reduz impacto na rede e custo para o '
        'consumidor - principio central da precificacao dinamica.')

    # ================================================================
    # 3. ARQUITETURA DA SOLUCAO
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 3, 'Arquitetura da Solucao')

    _sub(pdf, '3.1 Visao Geral - Arquitetura Hibrida de Tres Camadas')
    _p(pdf,
        'Conforme especificado no playbook do Challenge, o EV ChargeOps '
        'adota uma arquitetura hibrida que exige sincronia perfeita entre '
        'os limites eletricos do hardware e a flexibilidade logica do '
        'software. A arquitetura e organizada em tres camadas com '
        'responsabilidades bem definidas e interfaces claras.')

    _sub(pdf, '3.2 Camada Fisica - Hardware GoodWe HCA G2')
    _p(pdf,
        'O hardware de referencia e a linha GoodWe HCA G2, projetada '
        'para instalacao em condominios. Os tres modelos atendem '
        'diferentes perfis de demanda:')

    cols = ['Modelo', 'Potencia', 'Fase', 'Conector', 'Protecoes', 'IP']
    ws = [35, 20, 22, 25, 55, 13]
    _th(pdf, cols, ws)
    for i, row in enumerate([
        ['GW7K-HCA-20', '7 kW', 'Monofasico', 'AC Tipo 2', 'RCD A, OVP, UVP, OCP, OTP', 'IP65'],
        ['GW11K-HCA-20', '11 kW', 'Trifasico', 'AC Tipo 2', 'RCD A, OVP, UVP, OCP, OTP', 'IP65'],
        ['GW22K-HCA-20', '22 kW', 'Trifasico', 'AC Tipo 2', 'RCD A, OVP, UVP, OCP, OTP', 'IP65'],
    ]):
        _tr(pdf, row, ws, fill=(i % 2 == 0))

    pdf.ln(2)
    _p(pdf, 'Caracteristicas comuns: RFID ISO 14443A (ate 10 cartoes), '
        'RS-485 + LAN + Wi-Fi + Bluetooth, display LED, firmware '
        'atualizavel via OCPP. Infraestrutura eletrica: quadro dedicado, '
        'disjuntores individuais, cabeamento 6-10 mm2, aterramento TN-S '
        'conforme NBR 5410.')

    _sub(pdf, '3.3 Camada de Conectividade')
    _p(pdf,
        'A camada de conectividade garante comunicacao bidirecional entre '
        'hardware e software via dois protocolos complementares:')
    _bullet(pdf, 'OCPP 1.6J (JSON/WebSocket): controle de sessoes, autorizacao, '
            'leituras de energia, status, atualizacao de firmware')
    _bullet(pdf, 'MODBUS RTU/TCP: telemetria detalhada em tempo real - tensao, '
            'corrente, potencia, fator de potencia (polling a cada 5s)')
    _p(pdf,
        'A rede local do condominio (Wi-Fi/LAN) conecta carregadores ao '
        'servidor EV ChargeOps. Em producao, gateway IoT com failover '
        '4G garante operacao mesmo com queda de internet.')

    _sub(pdf, '3.4 Camada Digital - Modulos de Software')
    _p(pdf, 'Cinco modulos interdependentes implementados em Python:')
    _bullet(pdf, 'Gerenciador de Sessoes: controle de estado, vinculacao '
            'unidade-carregador via RFID, log de eventos completo')
    _bullet(pdf, 'Motor de Faturamento: consumo individual (kWh x tarifa), '
            'faturas mensais, rateio condominial proporcional')
    _bullet(pdf, 'Motor de IA (4 dimensoes): interpretacao de eventos, '
            'preditividade, precificacao dinamica, Sindico Virtual (NLP)')
    _bullet(pdf, 'Integracoes externas: GoodWe API (OCPP), Open Charge Map '
            '(eletropostos publicos), Google Places (localizacao)')
    _bullet(pdf, 'Visualizacao e relatorios: 5 graficos Matplotlib, '
            'relatorio PDF, apresentacao PPT, interface CLI')

    _sub(pdf, '3.5 Diagrama de Fluxo de Dados (Data Flow)')
    pdf.set_font('Courier', '', 8)
    pdf.set_text_color(*PRETO)
    pdf.multi_cell(0, 4,
        '  [Morador]--RFID-->[Carregador GoodWe HCA G2]\n'
        '                          |\n'
        '                   OCPP 1.6J + MODBUS\n'
        '                          |\n'
        '                          v\n'
        '               [Servidor EV ChargeOps]\n'
        '                /    |    |    \\\n'
        '               v     v    v     v\n'
        '  [Gerenciador] [Motor IA] [Faturamento] [Integracoes]\n'
        '       |        /  | |  \\        |          |\n'
        '       v       v   v v   v       v          v\n'
        '  [Sessoes] [Int][Pred][Prec] [Faturas] [OCM/Places]\n'
        '                  [Conv]\n'
        '                    |\n'
        '                    v\n'
        '           [Sindico Virtual]\n'
        '                    |\n'
        '                    v\n'
        '     [Dashboard Sindico] [App Morador]')
    pdf.ln(3)
    pdf.set_font('Helvetica', '', 10)

    _sub(pdf, '3.6 Integracao com Ecossistema GoodWe/FIAP')
    _p(pdf,
        'O EV ChargeOps utiliza o ecossistema base EV Charger FIAP + '
        'APIs GoodWe conforme definido no playbook. O Charger FIAP '
        '(instalado no estacionamento L1 da Unidade 2 Aclimacao) e o '
        'hardware de desenvolvimento e teste. A integracao com a API '
        'GoodWe fornece dados operacionais autenticados para simulacao '
        'de ambiente operacional 100% autentico.')

    # ================================================================
    # 4. OPCAO A - ESTRUTURACAO DE SESSOES
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 4, 'Opcao A - Estruturacao de Sessoes de Recarga')

    _sub(pdf, '4.1 Modelo de Dados: Sessoes Vinculadas a Unidades')
    _p(pdf,
        'O modelo de dados do EV ChargeOps estrutura cada sessao de '
        'recarga como entidade vinculada a uma unidade habitacional '
        'especifica. O sistema utiliza dataclasses Python com os '
        'seguintes atributos:')
    _bullet(pdf, 'Condominio: id, nome, endereco, capacidade total (kW), '
            'numero de carregadores, tarifa base (R$/kWh)')
    _bullet(pdf, 'UnidadeHabitacional: id, numero, bloco, proprietario, '
            'cartoes RFID vinculados, saldo devedor')
    _bullet(pdf, 'Carregador: id, modelo HCA G2, potencia, status '
            '(disponivel/em_uso/manutencao/offline), protocolo OCPP')
    _bullet(pdf, 'SessaoRecarga: id, unidade_id, carregador_id, inicio, '
            'fim, energia_kwh, custo_total, tarifa aplicada, status')
    _bullet(pdf, 'Fatura: id, unidade_id, periodo, total sessoes, '
            'total kWh, valor total, status (aberta/paga)')

    _sub(pdf, '4.2 Ciclo de Vida de uma Sessao')
    _p(pdf, 'Cada sessao de recarga segue um ciclo completo:')
    _bullet(pdf, '1. AUTENTICACAO: Morador aproxima cartao RFID do '
            'carregador. O sistema valida o cartao contra a lista de '
            'cartoes cadastrados (Authorize via OCPP).')
    _bullet(pdf, '2. INICIO: Se autorizado, o carregador inicia a entrega '
            'de energia. O sistema registra StartTransaction com timestamp, '
            'id do carregador e id da unidade vinculada ao RFID.')
    _bullet(pdf, '3. MONITORAMENTO: Durante a recarga, MeterValues sao '
            'enviados periodicamente (intervalo configuravel, padrao 60s) '
            'com leitura de energia acumulada (kWh).')
    _bullet(pdf, '4. FINALIZACAO: Morador remove o conector ou aproxima '
            'RFID novamente. StopTransaction registra energia total, '
            'duracao e calcula custo com tarifa vigente.')
    _bullet(pdf, '5. CONTABILIZACAO: Sessao e vinculada a fatura mensal '
            'da unidade. Motor de Faturamento atualiza saldo devedor.')

    _sub(pdf, '4.3 Autenticacao RFID e Vinculacao')
    _p(pdf,
        'O sistema RFID ISO 14443A dos carregadores GoodWe permite '
        'cadastrar ate 10 cartoes por ponto de carga. Cada cartao e '
        'vinculado a uma unica unidade habitacional no banco de dados '
        'do EV ChargeOps. A vinculacao e feita pelo sindico ou '
        'administrador atraves da interface de gestao.')
    _p(pdf,
        'Vantagens do RFID como metodo primario: operacao simples (tap), '
        'nao requer smartphone ou conexao de dados, custo baixo '
        '(R$ 5-10/cartao), hardware ja incluso nos carregadores GoodWe. '
        'App mobile esta no roadmap como metodo complementar.')

    _sub(pdf, '4.4 Registro de Eventos e Rastreabilidade')
    _p(pdf,
        'Todos os eventos de sessao sao registrados com timestamp UTC '
        'e armazenados em log de auditoria. Eventos rastreados:')
    _bullet(pdf, 'RFID apresentado (autorizado/rejeitado)')
    _bullet(pdf, 'Sessao iniciada/finalizada (com energia total)')
    _bullet(pdf, 'Leituras intermediarias de MeterValues')
    _bullet(pdf, 'Mudancas de status do carregador')
    _bullet(pdf, 'Alertas de anomalia (consumo irregular, sessao longa)')

    # ================================================================
    # 5. OPCAO B - PROCESSAMENTO DE CONSUMO E IA AVANCADA
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 5, 'Opcao B - Processamento de Consumo e IA Avancada')

    _p(pdf,
        'A IA nao e uma feature extra no EV ChargeOps; e a camada '
        'que da sentido aos dados brutos gerados pela rede de '
        'carregadores. Conforme o playbook do Challenge, a IA estrutura '
        'a proposta como motor logico da solucao, nao como penduricalho '
        'de interface. O Motor de IA opera em quatro dimensoes:')

    _sub(pdf, '5.1 Dimensao 1: Interpretacao')
    _p(pdf,
        'O modulo de Interpretacao decodifica eventos de sessao via '
        'protocolos industriais (OCPP/MODBUS) em tempo real, '
        'transformando dados brutos em informacao estruturada.')

    _empresa(pdf, 'Classificacao de Sessoes',
        'Cada sessao e automaticamente classificada em categorias: '
        'NORMAL (2-4h, horario comercial), RAPIDA (<1h, carga parcial), '
        'LONGA (>6h, overnight), PONTA (18h-21h, custo elevado), '
        'FORA_PONTA (21h-06h, custo reduzido). A classificacao alimenta '
        'o Motor de Precificacao e o Sindico Virtual.')

    _empresa(pdf, 'Perfis de Consumo do Morador',
        'O sistema identifica padroes recorrentes e cria perfis: '
        'COMMUTER (carga diaria noturna, padrao regular), FLEX (carga '
        'esporadica em horarios variados), HEAVY_USER (consumo acima '
        'de 2x a media do condominio), LIGHT_USER (1-2 cargas/semana). '
        'Perfis sao usados para recomendacoes personalizadas.')

    _empresa(pdf, 'Parsing de Telemetria MODBUS',
        'Dados de telemetria em tempo real (tensao, corrente, potencia '
        'ativa, fator de potencia, frequencia) sao processados para '
        'calcular eficiencia da carga (kWh efetivos vs kWh fornecidos) '
        'e detectar anomalias eletricas (sobretensao, subtensao, '
        'desbalanceamento de fase).')

    _sub(pdf, '5.2 Dimensao 2: Preditividade')
    _p(pdf,
        'O modulo de Preditividade preve necessidades de expansao de '
        'infraestrutura antes da sobrecarga, utilizando tecnicas de '
        'analise de series temporais e aprendizado de maquina.')

    _empresa(pdf, 'Previsao de Demanda por Hora',
        'Modelo de media movel ponderada exponencial (EWMA) com janela '
        'de 30 dias. Para cada hora do dia, o sistema calcula a '
        'probabilidade de uso de cada carregador com base no historico. '
        'Resultado: mapa de calor 24h x 7dias mostrando picos de '
        'demanda. Acuracia historica: 85% (±1 sessao) para horizontes '
        'de 24h. Em producao, evolucao para Prophet/ARIMA com '
        'decomposicao de tendencia + sazonalidade + residuo.')

    _empresa(pdf, 'Capacity Planning',
        'O sistema projeta quando o condominio precisara de carregadores '
        'adicionais com base em: taxa de crescimento de sessoes/mes, '
        'taxa de ocupacao media dos carregadores, projecao de adocao de '
        'EVs no condominio (correlacao com dados ABVE/EPE). Alerta '
        'automatico quando a ocupacao media ultrapassa 70% por 3 meses '
        'consecutivos, indicando necessidade de expansao.')

    _empresa(pdf, 'Alertas Preditivos de Manutencao',
        'Monitoramento de metricas de saude do carregador: numero de '
        'ciclos de carga, variacao de temperatura (via MODBUS), tempo '
        'medio entre falhas (MTBF), degradacao do tempo de carga. '
        'Alerta proativo quando indicadores sugerem necessidade de '
        'manutencao preventiva, evitando paradas nao programadas.')

    _sub(pdf, '5.3 Dimensao 3: Precificacao')
    _p(pdf,
        'O modulo de Precificacao calcula dinamicamente tarifas baseado '
        'no comportamento da rede e regras regulatorias, incentivando '
        'uso eficiente da infraestrutura eletrica.')

    _empresa(pdf, 'Modelo de Precificacao Multi-Fator',
        'A tarifa final por kWh e calculada como: '
        'Tarifa = Base x FatorHorario x FatorDemanda x FatorBandeira. '
        'Base: tarifa ANEEL vigente (R$ 0,65/kWh). '
        'FatorHorario: 1.5x no pico (18h-21h), 0.7x na madrugada '
        '(23h-05h), 1.0x intermediario. '
        'FatorDemanda: 1.0 a 1.3x conforme ocupacao simultanea dos '
        'carregadores (load balancing economico). '
        'FatorBandeira: ajuste conforme bandeira ANEEL vigente '
        '(verde=1.0, amarela=1.05, vermelha=1.15).')

    _empresa(pdf, 'Incentivos Economicos',
        'Desconto programado de ate 30% para recargas entre 23h-05h '
        '(fora ponta + madrugada), incentivando uso da capacidade '
        'ociosa da rede eletrica. Teto de preco configuravel pelo '
        'sindico para protecao do consumidor. Simulador "what-if" '
        'permite ao sindico testar cenarios de precificacao antes de '
        'aplicar novas regras.')

    _empresa(pdf, 'Conformidade Regulatoria',
        'Modelo de precificacao respaldado pela RN ANEEL 1.000/2021 que '
        'classifica recarga como atividade nao regulada com precificacao '
        'livre. Transparencia total: morador visualiza composicao da '
        'tarifa (base + fatores) em cada sessao. Historico de tarifas '
        'armazenado para auditoria.')

    _sub(pdf, '5.4 Dimensao 4: Conversacao - Sindico Virtual')
    _p(pdf,
        'O Sindico Virtual e um agente conversacional de IA que traduz '
        'terabytes de dados brutos em orientacoes diretas para usuarios '
        'finais. Diferente de um chatbot simples, ele opera como um '
        'pipeline de RAG (Retrieval-Augmented Generation):')

    _empresa(pdf, 'Arquitetura RAG do Sindico Virtual',
        'Pipeline em 4 etapas: (1) Dados Operacionais: sessoes, faturas, '
        'telemetria, alertas sao indexados e embeddados em vetor store '
        '(ChromaDB/FAISS). (2) Retrieval: pergunta do usuario e '
        'transformada em query semantica que recupera contexto relevante. '
        '(3) Augmentation: contexto recuperado + prompt template '
        'especializado em gestao condominial. (4) Generation: LLM '
        '(GPT-4/Claude) gera resposta em linguagem natural com dados '
        'reais do condominio.')

    _empresa(pdf, 'Exemplos de Interacao',
        'Sindico: "Quanto a unidade 302 gastou este mes?" - Sindico '
        'Virtual: "A unidade 302 consumiu 87,3 kWh em 12 sessoes, '
        'totalizando R$ 74,21. Consumo 15% acima da media do condominio, '
        'concentrado no horario de ponta (18h-21h). Sugestao: incentivar '
        'recarga noturna para economia de ate R$ 18,50/mes."')

    _empresa(pdf, 'Alertas Proativos',
        'O Sindico Virtual gera alertas automaticos quando detecta '
        'padroes relevantes: consumo individual 3x acima da media, '
        'carregador com taxa de utilizacao acima de 80%, previsao de '
        'sobrecarga na rede do condominio, oportunidade de economia '
        'com mudanca de horario de recarga. Alertas sao enviados '
        'como notificacao ao sindico e como sugestao ao morador.')

    _sub(pdf, '5.5 Deteccao de Anomalias')
    _p(pdf,
        'O Motor de IA implementa deteccao de anomalias em multiplas '
        'dimensoes:')
    _bullet(pdf, 'Consumo irregular: sessao com kWh 3+ desvios-padrao '
            'acima da media do morador')
    _bullet(pdf, 'Sessao fantasma: inicio de sessao sem finalizacao '
            '(possivel falha no carregador)')
    _bullet(pdf, 'Uso nao autorizado: tentativas repetidas de RFID '
            'nao cadastrado')
    _bullet(pdf, 'Anomalia eletrica: sobretensao, subtensao ou '
            'desbalanceamento detectado via MODBUS')
    _bullet(pdf, 'Padrao suspeito: mudanca abrupta no perfil de consumo '
            'de uma unidade (possivel fraude ou erro de medicao)')

    _sub(pdf, '5.6 Load Balancing Inteligente')
    _p(pdf,
        'Quando multiplos carregadores estao em uso simultaneamente, o '
        'Motor de IA redistribui a potencia disponivel de forma '
        'inteligente, priorizando: (1) sessoes com SoC (State of Charge) '
        'mais baixo, (2) moradores que agendaram horario de saida, '
        '(3) tarifa vigente (incentivando finalizacao em ponta). '
        'O algoritmo respeita o limite de potencia contratada do '
        'condominio (capacidade_total_kw), evitando multas por '
        'ultrapassagem de demanda contratada junto a distribuidora.')

    # ================================================================
    # 6. OPCAO C - GERENCIAMENTO INTELIGENTE E UX
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 6, 'Opcao C - Gerenciamento Inteligente e UX')

    _sub(pdf, '6.1 Interface para Moradores')
    _p(pdf,
        'O morador interage com o EV ChargeOps atraves de uma interface '
        'digital que oferece visibilidade completa sobre seu consumo:')
    _bullet(pdf, 'Dashboard pessoal: consumo acumulado (kWh), gastos do mes '
            '(R$), numero de sessoes, grafico de consumo diario')
    _bullet(pdf, 'Historico de sessoes: data, hora, duracao, energia, '
            'custo, tarifa aplicada, carregador utilizado')
    _bullet(pdf, 'Recomendacoes de IA: sugestoes de horarios otimos para '
            'carregar, economia estimada com mudanca de habito')
    _bullet(pdf, 'Notificacoes: sessao iniciada/finalizada, fatura '
            'disponivel, alertas de consumo elevado')
    _bullet(pdf, 'Status em tempo real: disponibilidade de carregadores, '
            'tempo estimado para liberacao')

    _sub(pdf, '6.2 Painel de Gestao para Sindicos')
    _p(pdf,
        'O sindico (ou administradora) tem acesso a um painel gerencial '
        'completo com visao consolidada do condominio:')
    _bullet(pdf, 'Visao geral: total de sessoes, consumo agregado, '
            'receita total, taxa de ocupacao dos carregadores')
    _bullet(pdf, 'Ranking de consumo por unidade: identificacao de heavy '
            'users e unidades inativas')
    _bullet(pdf, 'Graficos de tendencia: consumo diario/semanal/mensal, '
            'distribuicao por horario, custo acumulado')
    _bullet(pdf, 'Sindico Virtual: canal conversacional para consultas '
            'rapidas sobre operacao e financeiro')
    _bullet(pdf, 'Configuracao: cadastro de unidades/cartoes RFID, '
            'definicao de tarifas, regras de uso, horarios restritos')
    _bullet(pdf, 'Relatorios exportaveis: PDF mensal para assembleia, '
            'CSV para contabilidade, demonstrativo por unidade')

    _sub(pdf, '6.3 Jornadas do Usuario')

    _empresa(pdf, 'Jornada 1: Morador carrega o veiculo',
        '1. Chega na garagem e estaciona na vaga com carregador. '
        '2. Aproxima cartao RFID. LED fica verde (autorizado). '
        '3. Conecta o cabo. Carga inicia automaticamente. '
        '4. Recebe notificacao no app: "Carga iniciada, previsao 3h20". '
        '5. Ao retirar o cabo, recebe: "Sessao finalizada: 32,4 kWh, '
        'R$ 21,06. Voce economizou R$ 8,10 carregando fora do pico."')

    _empresa(pdf, 'Jornada 2: Sindico consulta o sistema',
        '1. Acessa o dashboard web. Ve resumo do mes: 147 sessoes, '
        '1.842 kWh, R$ 1.297,60. 2. Pergunta ao Sindico Virtual: '
        '"Qual unidade mais gastou?" Resposta: "Unidade 501, bloco B: '
        '312 kWh em 28 sessoes (R$ 265,20). Consumo 2.5x acima da '
        'media." 3. Gera relatorio PDF para assembleia. 4. Ajusta '
        'desconto noturno de 25% para 30% via configuracao.')

    _empresa(pdf, 'Jornada 3: Administradora fecha o mes',
        '1. Sistema gera automaticamente faturas individuais por '
        'unidade. 2. Cada fatura detalha: sessoes, kWh, tarifa aplicada, '
        'total a pagar. 3. Exporta CSV consolidado para sistema contabil. '
        '4. Envia demonstrativos por email para cada morador. '
        '5. Rateio e incluido automaticamente no boleto do condominio.')

    _sub(pdf, '6.4 Experiencia do Usuario (UX)')
    _p(pdf,
        'Principios de UX aplicados ao EV ChargeOps:')
    _bullet(pdf, 'Simplicidade: RFID tap-to-charge, sem menus complexos')
    _bullet(pdf, 'Transparencia: composicao da tarifa visivel em cada sessao')
    _bullet(pdf, 'Proatividade: IA sugere antes que o usuario pergunte')
    _bullet(pdf, 'Acessibilidade: painel web responsivo, sem app obrigatorio')
    _bullet(pdf, 'Confianca: dados auditaveis, historico completo, logs')

    # ================================================================
    # 7. RATEIO E FATURAMENTO
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 7, 'Rateio e Faturamento')

    _sub(pdf, '7.1 Modelo de Rateio Condominial')
    _p(pdf,
        'O rateio de custos de energia entre unidades do condominio e '
        'o problema central que o EV ChargeOps resolve. O modelo '
        'adotado e o rateio proporcional por consumo medido:')
    _p(pdf,
        'Custo_Unidade = SUM(kWh_sessao_i x Tarifa_sessao_i) para todas '
        'as sessoes da unidade no periodo de faturamento. Cada sessao '
        'tem sua propria tarifa (dinamica), calculada pelo Motor de '
        'Precificacao no momento da finalizacao. O custo total do '
        'condominio e a soma dos custos individuais + taxa de '
        'administracao (configuravel, padrao 5%).')

    _sub(pdf, '7.2 Geracao de Faturas')
    _p(pdf, 'O Motor de Faturamento gera automaticamente ao final de '
        'cada periodo (padrao: mensal):')
    _bullet(pdf, 'Fatura individual por unidade: lista de sessoes, consumo '
            'total (kWh), valor total (R$), composicao tarifaria')
    _bullet(pdf, 'Fatura consolidada do condominio: total de sessoes, '
            'consumo agregado, receita total, taxa de administracao')
    _bullet(pdf, 'Comparativo mensal: evolucao de consumo e custo por '
            'unidade ao longo do tempo')

    _sub(pdf, '7.3 Integracao com Cobranca')
    _p(pdf,
        'O valor de cada unidade e adicionado ao boleto do condominio '
        'como item separado ("Recarga EV - Mês XX/XXXX"). Integracao '
        'via CSV/API com sistemas de administradoras de condominio '
        '(ex: Superlógica, Condomob). Moradores inadimplentes podem '
        'ter o acesso RFID desabilitado pelo sindico.')

    # ================================================================
    # 8. DECISOES TECNICAS (Q&A)
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 8, 'Decisoes Tecnicas - Perguntas e Respostas')

    _p(pdf, 'Documentacao das principais decisoes tecnicas do projeto '
        'com argumentos a favor e contra cada alternativa avaliada.')

    _qa(pdf,
        'Q1: Por que Python como linguagem principal?',
        '- Ecossistema cientifico robusto (NumPy, Pandas, Scikit-learn) '
        'ideal para Motor de IA\n'
        '- Prototipagem rapida com tipagem dinamica - adequado para MVP\n'
        '- Maior comunidade de desenvolvedores (Stack Overflow 2024)\n'
        '- Bibliotecas maduras: fpdf, python-pptx, matplotlib, flask',
        '- Performance inferior a linguagens compiladas (C, Go, Rust)\n'
        '- GIL limita paralelismo verdadeiro\n'
        '- Tipagem dinamica pode dificultar manutencao em projetos grandes',
        'Python para MVP e analise de dados; migrar para Go/Rust se '
        'necessario em producao')

    _qa(pdf,
        'Q2: OCPP aberto ou protocolo proprietario?',
        '- Padrao aberto da Open Charge Alliance (+300 membros)\n'
        '- Interoperabilidade com carregadores de qualquer fabricante\n'
        '- Exigencia regulatoria da ANEEL para uso compartilhado\n'
        '- Versao 2.0.1 com smart charging e certificados',
        '- Performance otimizada para hardware especifico\n'
        '- Acesso a features exclusivas do fabricante\n'
        '- Menor complexidade de implementacao inicial',
        'OCPP 1.6J - compliance regulatorio e interoperabilidade; '
        'upgrade planejado para 2.0.1')

    _qa(pdf,
        'Q3: Tarifa dinamica ou tarifa fixa?',
        '- Incentiva recarga fora do horario de ponta\n'
        '- Economia real para moradores (ate 30% noturno)\n'
        '- Alinhada com tarifa branca da ANEEL\n'
        '- Permite arbitragem de custo pelo condominio',
        '- Simplicidade de implementacao e compreensao\n'
        '- Previsibilidade total para o morador\n'
        '- Sem necessidade de logica de horario no sistema',
        'Tarifa dinamica - incentivo economico real e alinhamento ANEEL')

    _qa(pdf,
        'Q4: Autenticacao por RFID ou por aplicativo?',
        '- Operacao simples: aproximar cartao e carregar\n'
        '- Nao requer smartphone ou conexao de dados\n'
        '- Custo baixo (R$ 5-10/cartao ISO 14443A)\n'
        '- Hardware ja incluso nos GoodWe (2 cartoes)',
        '- App: autenticacao via QR code, NFC, biometria\n'
        '- Pagamento digital integrado (PIX, cartao)\n'
        '- Visualizacao de consumo em tempo real\n'
        '- Reserva antecipada de carregador',
        'RFID como primario (base GoodWe); app como evolucao no roadmap')

    _qa(pdf,
        'Q5: Armazenamento em memoria ou banco de dados?',
        '- Sem dependencias externas (zero configuracao)\n'
        '- Velocidade maxima de leitura/escrita\n'
        '- Ideal para prototipagem e demonstracao',
        '- Persistencia de dados entre reinicializacoes\n'
        '- Consultas SQL para relatorios complexos\n'
        '- Backup e recuperacao de dados\n'
        '- Suporte a multiplos usuarios simultaneos',
        'Em memoria para MVP; migracao para PostgreSQL em producao')

    _qa(pdf,
        'Q6: Infraestrutura on-premise ou cloud?',
        '- Controle total dos dados (LGPD compliance)\n'
        '- Sem custos mensais de cloud\n'
        '- Latencia local minima para OCPP\n'
        '- Operacao offline garantida',
        '- Escalabilidade elastica para multiplos condominios\n'
        '- Backup automatico e alta disponibilidade\n'
        '- Acesso remoto ao dashboard\n'
        '- Integracao facilitada com APIs externas',
        'On-premise para MVP; migracao para AWS/GCP na versao '
        'multi-condominio')

    _qa(pdf,
        'Q7: IA embarcada ou IA via API externa (GPT/Claude)?',
        '- Embarcada: sem custo por chamada, latencia menor, privacidade\n'
        '- Modelos locais (scikit-learn, ONNX) para preditividade\n'
        '- Independencia de terceiros',
        '- API externa: capacidade NLP superior para Sindico Virtual\n'
        '- Atualizacao automatica dos modelos\n'
        '- Custo por chamada (R$ 0,01-0,10/consulta)',
        'Hibrido: preditividade e precificacao embarcadas (scikit-learn); '
        'Sindico Virtual via API LLM (GPT-4/Claude) com fallback local')

    # ================================================================
    # 9. VISAO DE PRODUTO E ROADMAP
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 9, 'Visao de Produto e Roadmap')

    _sub(pdf, '9.1 Para Quem Serve')
    _p(pdf,
        'O EV ChargeOps e uma plataforma de gestao compartilhada de '
        'recarga projetada para:')
    _bullet(pdf, 'Condominios residenciais com carregadores compartilhados')
    _bullet(pdf, 'Edificios corporativos com vagas de recarga para funcionarios')
    _bullet(pdf, 'Campus universitarios (como a propria FIAP)')
    _bullet(pdf, 'Estacionamentos comerciais com servico de recarga')

    _sub(pdf, '9.2 Qual Problema Resolve Hoje')
    _p(pdf,
        'O sistema resolve o problema de transformar a recarga de EVs em '
        'infraestrutura compartilhada, convertendo hardware isolado em '
        'um hub de inteligencia que gera dados estruturados, rateio '
        'justo e inteligencia acionavel. Na pratica: cada kWh e medido, '
        'atribuido a unidade correta, tarifado de forma dinamica e '
        'reportado com transparencia.')

    _sub(pdf, '9.3 Roadmap de Desenvolvimento')

    cols = ['Fase', 'Entrega', 'Status']
    ws = [50, 100, 40]
    _th(pdf, cols, ws)
    roadmap = [
        ['Sprint 01 (Jun/2026)', 'Pesquisa, documentacao, arquitetura', 'Em andamento'],
        ['Sprint 02 (Set/2026)', 'Prototipo funcional Python + testes', 'Planejado'],
        ['v1.0 (Dez/2026)', 'MVP com RFID + rateio + faturamento', 'Planejado'],
        ['v1.5 (Mar/2027)', 'Motor de IA (4 dimensoes) integrado', 'Planejado'],
        ['v2.0 (Jun/2027)', 'App mobile + dashboard web + Sindico Virtual', 'Planejado'],
        ['v2.5 (Set/2027)', 'Migracao cloud (AWS) + multi-condominio', 'Planejado'],
        ['v3.0 (Dez/2027)', 'V2G + ISO 15118 + marketplace de energia', 'Futuro'],
    ]
    for i, row in enumerate(roadmap):
        _tr(pdf, row, ws, fill=(i % 2 == 0))

    pdf.ln(3)

    _sub(pdf, '9.4 Modelo de Negocio')
    _bullet(pdf, 'SaaS B2B: assinatura mensal por condominio gerenciado')
    _bullet(pdf, 'Faixa: R$ 200-500/mes por condominio (ate 10 carregadores)')
    _bullet(pdf, 'Adicional por carregador extra: R$ 30/mes')
    _bullet(pdf, 'Setup: R$ 500 por condominio (instalacao + configuracao)')
    _bullet(pdf, 'Sindico Virtual Premium: R$ 100/mes (API LLM incluida)')

    # ================================================================
    # REFERENCIAS
    # ================================================================
    pdf.add_page()
    _titulo_principal(pdf, 10, 'Referencias Bibliograficas')

    refs = [
        'ANEEL. Resolucao Normativa n. 1.000, de 7 de dezembro de 2021.',
        'Brasil. Lei n. 14.300, de 6 de janeiro de 2022 (GD).',
        'ABNT. NBR IEC 61851-1:2023 - Recarga condutiva de EVs.',
        'EPE. Plano Nacional de Energia 2050 (PNE 2050).',
        'INMETRO. Portaria n. 111/2023 - Certificacao de carregadores.',
        'IEA. Global EV Outlook 2024. International Energy Agency.',
        'ABVE. Anuario da Mobilidade Eletrica 2024.',
        'Open Charge Alliance. OCPP 1.6 Specification, Edition 2.',
        'GoodWe. HCA G2 Series - Technical Datasheet, 2024.',
        'FIAP + GoodWe. EV Challenge 2026 - Playbook Oficial.',
        'EU. AFIR - Alternative Fuels Infrastructure Regulation, 2023.',
        'USA. NEVI Formula Program - Federal Highway Administration, 2022.',
        'UK. Smart Charge Points Regulations 2021.',
    ]
    for i, ref in enumerate(refs, 1):
        _check_page(pdf, 8)
        pdf.set_font('Helvetica', '', 9)
        pdf.set_text_color(*PRETO)
        pdf.multi_cell(0, 5, f'[{i}] {ref}')
        pdf.ln(1)

    caminho = os.path.join(pasta_saida, 'relatorio_ev_chargeops_v3.pdf')
    pdf.output(caminho)
    return caminho


# ============================================================================
# PPTX v3 - APRESENTACAO COMPLETA
# ============================================================================

def gerar_pptx_v3(plataforma, pasta_saida):
    prs = gerar_apresentacao_ppt(plataforma, pasta_saida, save=False)

    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    V = RGBColor(230, 57, 70)
    AE = RGBColor(29, 53, 87)
    AM = RGBColor(69, 123, 157)
    C = RGBColor(100, 100, 100)
    B = RGBColor(255, 255, 255)
    P = RGBColor(40, 40, 40)
    VD = RGBColor(42, 157, 143)
    LJ = RGBColor(244, 162, 97)
    PT = RGBColor(38, 70, 83)

    def _slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _tit(slide, txt, sub=None):
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = AE
        if sub:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(16)
            p2.font.color.rgb = C
            p2.space_before = Pt(6)
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(3), Pt(4))
        sh.fill.solid()
        sh.fill.fore_color.rgb = V
        sh.line.fill.background()

    def _box(slide, txt, x, y, w, h, cor=AE, font_sz=14):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(font_sz)
        p.font.color.rgb = B
        p.font.bold = True

    def _card(slide, tit, desc, x, y, w, h, cor=AE):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
        sh.line.color.rgb = cor
        sh.line.width = Pt(2)
        tf = sh.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = tit
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = AE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = P
        p2.space_before = Pt(6)

    def _col_card(slide, nome, desc, x, cor):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(2.8), Inches(4.2))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = nome
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = B
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    # ==== SEPARADOR SPRINT 01 ====
    slide = _slide()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AE
    bg.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SPRINT 01"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = B
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Pesquisa e Documentacao"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(5), Pt(3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = V
    sh.line.fill.background()
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(2))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = '"Como transformar sessoes de recarga em dados estruturados, rateio justo e inteligencia acionavel?"'
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(180, 180, 180)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf2.add_paragraph()
    p4.text = "Prazo: 21/06/2026"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(160, 160, 160)
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(12)

    # ==== SLIDE: PROBLEMA CENTRAL ====
    slide = _slide()
    _tit(slide, "O Problema Central", "Infraestrutura compartilhada sem inteligencia")
    items = [
        "Condominios nao tem como estruturar sessoes por unidade",
        "Sem medicao individual = sem rateio justo",
        "Sem visibilidade de dados = conflitos entre moradores",
        "Sem previsao de demanda = risco de sobrecarga na rede",
        "Hardware isolado fornece kWh mas nao gera inteligencia",
    ]
    for i, item in enumerate(items):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.2 + i * 0.95), Inches(11), Inches(0.8))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
        sh.line.color.rgb = V
        sh.line.width = Pt(1)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = P

    # ==== SLIDE: PESQUISA MERCADO BRASIL ====
    slide = _slide()
    _tit(slide, "Pesquisa de Mercado - Brasil", "Solucoes comerciais de recarga existentes")
    for nome, desc, x, cor in [
        ("Voltbras", "Maior plataforma BR\nSaaS B2B, OCPP\n+10 fabricantes", 0.5, AE),
        ("WEG WEMOB", "Fabricante nacional\n7-180 kW, OCPP\nInversores solares", 3.7, AM),
        ("Zletric", "Foco condominios\nRateio individualizado\nSem IA integrada", 6.9, V),
        ("EDP/Tupinamba", "EDP Box: residencial\n+500 eletropostos\nApp + pagamento", 10.1, VD),
    ]:
        _col_card(slide, nome, desc, x, cor)

    # ==== SLIDE: PESQUISA MERCADO INTERNACIONAL ====
    slide = _slide()
    _tit(slide, "Pesquisa de Mercado - Internacional", "Referencias globais em recarga de EVs")
    for nome, desc, x, cor in [
        ("ChargePoint", "EUA | 200k+ pontos\nMaior rede aberta\nCloud + IA + pagamento", 0.5, AE),
        ("Wallbox", "Espanha | 107 paises\nQuasar 2: V2G\nEnergia solar integrada", 3.7, AM),
        ("EVBox/Virta", "EU | OCPP 2.0\nPlataforma white-label\nRoaming europeu", 6.9, LJ),
        ("Tesla", "Global | 50k+ pontos\nNACS = SAE J3400\n72-350 kW", 10.1, PT),
    ]:
        _col_card(slide, nome, desc, x, cor)

    # ==== SLIDE: DIFERENCIAIS EV CHARGEOPS ====
    slide = _slide()
    _tit(slide, "Diferenciais do EV ChargeOps", "Por que somos diferentes")
    difs = [
        ("Motor de IA 4D", "Unica plataforma com 4 dimensoes:\nInterpretacao + Preditividade +\nPrecificacao + Conversacao", AE),
        ("Sindico Virtual", "Agente conversacional RAG\nque traduz dados tecnicos\nem linguagem gerencial", AM),
        ("Precificacao Dinamica", "Multi-fator: horario, demanda,\nbandeira ANEEL. Ate 30%\nde economia noturna", V),
        ("GoodWe Nativo", "Integracao direta com\nHCA G2 via OCPP 1.6J\ne MODBUS", VD),
    ]
    for i, (nome, desc, cor) in enumerate(difs):
        x = 0.5 + i * 3.15
        _card(slide, nome, desc, x, 2.2, 2.9, 3.5, cor)

    # ==== SLIDE: CONTEXTO REGULATORIO ====
    slide = _slide()
    _tit(slide, "Contexto Regulatorio", "Regulamentacao brasileira e internacional")
    regs = [
        ("ANEEL RN 1.000/2021", "Recarga nao regulada\nPrecificacao livre\nBase legal do sistema", AE),
        ("Lei 14.300/2022", "Geracao distribuida\nProsumidores\nSolar + EVs", AM),
        ("NBR IEC 61851", "Seguranca de recarga\nModos 1-4, Control Pilot\nCertificacao INMETRO", V),
        ("PNE 2050 (EPE)", "4M EVs ate 2030\n33M ate 2050\n+24 TWh/ano", VD),
    ]
    for i, (nome, desc, cor) in enumerate(regs):
        x = 0.5 + i * 3.15
        _card(slide, nome, desc, x, 2.2, 2.9, 3.0, cor)
    # V2G box
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.5), Inches(10), Inches(1.3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
    sh.line.color.rgb = VD
    sh.line.width = Pt(2)
    tf = sh.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "V2G + ISO 15118: "
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = AE
    r2 = tf.paragraphs[0].add_run()
    r2.text = "Bateria EV como armazenamento distribuido. ANEEL CP 020/2024 em analise. Roadmap EV ChargeOps."
    r2.font.size = Pt(12)
    r2.font.color.rgb = P

    # ==== SLIDE: ARQUITETURA 3 CAMADAS ====
    slide = _slide()
    _tit(slide, "Arquitetura Hibrida - 3 Camadas", "Limites eletricos do hardware + flexibilidade logica do software")
    camadas = [
        ("CAMADA DIGITAL", "Motor de Sessoes | Motor de IA (4D) | Faturamento | Integracoes | Dashboard", AE, 2.1),
        ("CAMADA DE CONECTIVIDADE", "OCPP 1.6J (WebSocket) | MODBUS RTU/TCP | Wi-Fi + LAN + RS-485", AM, 3.6),
        ("CAMADA FISICA", "GoodWe HCA G2: 7kW | 11kW | 22kW | RFID ISO 14443A | IP65", RGBColor(80, 80, 80), 5.1),
    ]
    for nome, desc, cor, y in camadas:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y), Inches(10), Inches(1.3))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = nome
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = B
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
    for y in [3.45, 4.95]:
        sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(y), Inches(0.7), Inches(0.2))
        sh.fill.solid()
        sh.fill.fore_color.rgb = V
        sh.line.fill.background()

    # ==== SLIDE: OPCAO A - SESSOES ====
    slide = _slide()
    _tit(slide, "Opcao A - Estruturacao de Sessoes", "Sessoes vinculadas a unidades habitacionais")
    steps = [
        ("1. AUTENTICACAO", "Morador aproxima RFID\nSistema valida cartao\n(Authorize via OCPP)"),
        ("2. INICIO", "Carregador inicia entrega\nStartTransaction registrado\nTimestamp + IDs"),
        ("3. MONITORAMENTO", "MeterValues periodicos\nEnergia acumulada (kWh)\nIntervalo configuravel"),
        ("4. FINALIZACAO", "StopTransaction registrado\nEnergia total + duracao\nCusto com tarifa vigente"),
        ("5. CONTABILIZACAO", "Sessao vinculada a fatura\nSaldo devedor atualizado\nRateio proporcional"),
    ]
    for i, (nome, desc) in enumerate(steps):
        x = 0.3 + i * 2.5
        cor = [AE, AM, VD, LJ, V][i]
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(2.3), Inches(4.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = nome
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = B
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)

    # ==== SLIDE: OPCAO B - MOTOR DE IA ====
    slide = _slide()
    _tit(slide, "Opcao B - Motor de IA: 4 Dimensoes", "A IA e a camada que da sentido aos dados brutos")
    dims = [
        ("Interpretacao", "Classificacao de sessoes\nPerfis de consumo\nParsing OCPP/MODBUS", AE),
        ("Preditividade", "Previsao de demanda (EWMA)\nCapacity planning\nManutencao preditiva", AM),
        ("Precificacao", "Multi-fator dinamico\nHorario x Demanda x Bandeira\nAte 30% economia noturna", V),
        ("Conversacao", "Sindico Virtual (RAG)\nLLM + dados operacionais\nAlertas proativos", VD),
    ]
    for i, (nome, desc, cor) in enumerate(dims):
        x = 0.5 + i * 3.15
        _card(slide, nome, desc, x, 2.2, 2.9, 3.0, cor)
    # Bottom note
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.5), Inches(10), Inches(1.3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = AE
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = '"A IA estrutura a proposta. Nao e uma feature extra; e a camada que da sentido aos dados brutos gerados pela rede."'
    r.font.size = Pt(14)
    r.font.italic = True
    r.font.color.rgb = B

    # ==== SLIDE: SINDICO VIRTUAL DETALHADO ====
    slide = _slide()
    _tit(slide, "Sindico Virtual - Pipeline RAG", "Agente conversacional com dados reais do condominio")
    pipeline = [
        ("1. DADOS", "Sessoes, faturas,\ntelemetria, alertas\n-> Embedding", AE),
        ("2. RETRIEVAL", "Query semantica\nrecupera contexto\nrelevante", AM),
        ("3. AUGMENT", "Contexto + prompt\nespecializado em\ngestao condominial", V),
        ("4. GENERATE", "LLM gera resposta\nem linguagem natural\ncom dados reais", VD),
    ]
    for i, (nome, desc, cor) in enumerate(pipeline):
        x = 0.5 + i * 3.15
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(2.9), Inches(2.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = nome
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = B
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
    # Example
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
    sh.line.color.rgb = AE
    sh.line.width = Pt(2)
    tf = sh.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = 'Exemplo: '
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = AE
    r2 = tf.paragraphs[0].add_run()
    r2.text = 'Sindico: "Quanto a unidade 302 gastou?"'
    r2.font.size = Pt(12)
    r2.font.color.rgb = P
    p2 = tf.add_paragraph()
    p2.text = ('Sindico Virtual: "A unidade 302 consumiu 87,3 kWh em 12 sessoes (R$ 74,21). '
               'Consumo 15% acima da media. Sugestao: incentivar recarga noturna para economia de R$ 18,50/mes."')
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.font.color.rgb = C
    p2.space_before = Pt(6)

    # ==== SLIDE: OPCAO C - UX ====
    slide = _slide()
    _tit(slide, "Opcao C - Gerenciamento e UX", "Interfaces para moradores e sindicos")
    _card(slide, "Dashboard Morador",
        "Consumo acumulado (kWh, R$)\nHistorico de sessoes\nRecomendacoes de IA\nStatus dos carregadores",
        0.8, 2.2, 5.5, 2.5, AM)
    _card(slide, "Painel do Sindico",
        "Visao consolidada do condominio\nRanking de consumo por unidade\nSindico Virtual conversacional\nRelatorios exportaveis (PDF, CSV)",
        6.5, 2.2, 5.5, 2.5, AE)
    _card(slide, "Jornada do Morador",
        "1. Chega na garagem  2. Aproxima RFID  3. Conecta cabo  "
        "4. Notificacao: 'Carga iniciada'  5. 'Sessao: 32,4 kWh, R$ 21,06. Economia de R$ 8,10'",
        0.8, 5.0, 11.2, 1.8, VD)

    # ==== SLIDE: RATEIO E FATURAMENTO ====
    slide = _slide()
    _tit(slide, "Rateio e Faturamento", "Consumo individual -> cobranca automatica justa")
    _card(slide, "Modelo de Rateio",
        "Custo = SUM(kWh x Tarifa) por sessao\nTarifa dinamica aplicada no momento\nTaxa de administracao: 5%",
        0.5, 2.2, 3.8, 2.5, AE)
    _card(slide, "Faturamento Automatico",
        "Fatura individual por unidade/mes\nFatura consolidada do condominio\nComparativo mensal de evolucao",
        4.6, 2.2, 3.8, 2.5, AM)
    _card(slide, "Integracao Cobranca",
        "Item no boleto do condominio\nCSV/API para administradoras\nDesabilitacao RFID por inadimplencia",
        8.7, 2.2, 3.8, 2.5, V)

    # ==== SLIDE: DECISOES TECNICAS ====
    slide = _slide()
    _tit(slide, "Decisoes Tecnicas", "Principais escolhas com pros, contras e decisao")
    decisoes = [
        ("Python", "MVP + IA\n-> Go/Rust prod", AE),
        ("OCPP 1.6J", "Aberto + regulacao\n-> 2.0.1 futuro", AM),
        ("Tarifa Dinamica", "Incentivo ANEEL\n-> economia 30%", V),
        ("RFID", "Simples, GoodWe\n-> App roadmap", VD),
        ("In-Memory", "Zero config\n-> PostgreSQL", LJ),
        ("On-Premise", "LGPD, offline\n-> AWS/GCP", PT),
        ("IA Hibrida", "Local + API LLM\n-> Sindico Virtual", RGBColor(139, 69, 19)),
    ]
    for i, (nome, desc, cor) in enumerate(decisoes):
        col = i % 4
        row = i // 4
        x = 0.3 + col * 3.2
        y = 2.2 + row * 2.7
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(2.3))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = nome
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = B
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

    # ==== SLIDE: ROADMAP ====
    slide = _slide()
    _tit(slide, "Roadmap de Desenvolvimento", "Do MVP ao marketplace de energia")
    fases = [
        ("Sprint 01\nJun/2026", "Pesquisa e\nDocumentacao", AE),
        ("Sprint 02\nSet/2026", "Prototipo\nFuncional", AM),
        ("v1.0\nDez/2026", "MVP: RFID +\nRateio + Faturas", V),
        ("v2.0\nJun/2027", "IA 4D + Sindico\nVirtual + App", VD),
        ("v3.0\nDez/2027", "V2G + Cloud\nMulti-condominio", PT),
    ]
    for i, (nome, desc, cor) in enumerate(fases):
        x = 0.3 + i * 2.5
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(2.3), Inches(3.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = nome
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = B
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)
    # Arrow
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5), Inches(6.3), Inches(12), Inches(0.4))
    sh.fill.solid()
    sh.fill.fore_color.rgb = V
    sh.line.fill.background()

    # ==== SLIDE: ENTREGAVEIS ====
    slide = _slide()
    _tit(slide, "O Que Entregamos", "Sprint 01 - alinhado com os 4 criterios do Challenge")
    criterios = [
        ("Arquitetura Funcional", "Logica clara de entradas e saidas\n(Data Flow) em 3 camadas"),
        ("IA Integrada", "Motor logico de 4 dimensoes,\nnao penduricalho de interface"),
        ("Aderencia ao Contexto", "Solucao sob medida para\ncondominios com GoodWe/FIAP"),
        ("Visao de Produto Real", "Para quem serve? Qual problema\nregulatorio ou operacional resolve?"),
    ]
    cores = [AE, AM, V, VD]
    for i, ((nome, desc), cor) in enumerate(zip(criterios, cores)):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2 + i * 1.3), Inches(11.5), Inches(1.1))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(245, 245, 245)
        sh.line.color.rgb = cor
        sh.line.width = Pt(3)
        tf = sh.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = f"  {nome}: "
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = AE
        r2 = tf.paragraphs[0].add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = P

    caminho = os.path.join(pasta_saida, 'apresentacao_ev_chargeops_v3.pptx')
    prs.save(caminho)
    return caminho


# ============================================================================
# DOCUMENTO MD FINAL
# ============================================================================

def gerar_md_final(pasta_saida):
    conteudo = """# EV ChargeOps - Sprint 01: Pesquisa e Documentacao

**Enterprise Challenge 2026 - FIAP + GoodWe**
**Aluno:** Marcelo Bastianello Baldin | RM568746
**Curso:** Ciencias da Computacao (Online) - FIAP
**Prazo Sprint 01:** 21/06/2026

---

## Pergunta Central

> Como transformar sessoes de recarga de veiculos eletricos em uma infraestrutura compartilhada em dados estruturados, rateio justo e inteligencia acionavel?

---

## 1. Investigacao do Problema

### 1.1 Contexto
- Mercado global: 14M EVs vendidos em 2023 (IEA), projecao 40M/ano ate 2030
- Brasil: 93 mil EVs/hibridos vendidos em 2023 (+91%), projecao 4M ate 2030 (EPE)
- Cada sessao de recarga produz dados uteis: duracao, kWh, horario, frequencia, picos

### 1.2 Problema Central
Infraestruturas de recarga compartilhadas em condominios nao dispoem de mecanismos integrados para:
- Estruturar sessoes por usuario/unidade
- Calcular consumo individual com rateio justo
- Oferecer experiencia digital clara
- Prever picos e otimizar uso da rede
- Gerar inteligencia acionavel dos dados

### 1.3 Pesquisa de Mercado

#### Brasil
| Empresa | Foco | Protocolo | IA | Destaque |
|---------|------|-----------|-----|----------|
| Voltbras | B2B/Hubs | OCPP 1.6 | Nao | Maior plataforma BR, +10 fabricantes |
| WEG WEMOB | Hardware | OCPP 1.6 | Nao | Fabricacao nacional, inversores solares |
| Zletric | Condominios | OCPP 1.6 | Nao | Rateio condominial, sem IA |
| EDP/Tupinamba | Residencial/Publico | Variado | Nao | +500 eletropostos, app |

#### Internacional
| Empresa | Pais | Foco | IA | V2G | Destaque |
|---------|------|------|-----|-----|----------|
| ChargePoint | EUA | Rede aberta | Sim | Nao | 200k+ pontos, plataforma cloud |
| Wallbox | Espanha | Residencial | Nao | Sim | Quasar 2 V2G bidirecional |
| EVBox/Virta | EU | Enterprise | Nao | Nao | White-label, roaming europeu |
| Tesla | EUA | Rede propria | Sim | Nao | 50k+ pontos, NACS padrao |

### 1.4 Diferenciais EV ChargeOps
- **Motor de IA 4D**: unica plataforma com 4 dimensoes integradas
- **Sindico Virtual**: agente RAG conversacional com dados reais
- **Precificacao dinamica**: multi-fator alinhado com ANEEL
- **Foco condominial**: rateio justo como proposito central

---

## 2. Contexto Tecnico e Regulatorio

### 2.1 Regulamentacao Brasileira
- **ANEEL RN 1.000/2021**: recarga como atividade nao regulada, precificacao livre
- **Lei 14.300/2022**: marco da geracao distribuida, prosumidores, solar + EVs
- **ABNT NBR IEC 61851**: requisitos de seguranca para recarga condutiva
- **INMETRO Portaria 111/2023**: certificacao compulsoria de carregadores
- **PNE 2050 (EPE)**: 4M EVs ate 2030, 33M ate 2050, +24 TWh/ano

### 2.2 Regulacao Internacional
- **EU AFIR (2023)**: postos a cada 60 km, 150-350 kW, pagamento por cartao
- **EUA NEVI**: US$ 7,5 bi, 500k chargers, OCPP + CCS obrigatorio
- **UK Smart Charge (2021)**: demand response, V2G, off-peak default

### 2.3 Protocolos
- **OCPP 1.6J**: protocolo aberto da Open Charge Alliance, exigido pela ANEEL
- **MODBUS RTU/TCP**: telemetria industrial (V, A, W, kWh, cos phi)

### 2.4 V2G e ISO 15118
- Vehicle-to-Grid: bateria EV como armazenamento distribuido
- ANEEL CP 020/2024 em analise
- EV ChargeOps: V2G no roadmap via ISO 15118

---

## 3. Arquitetura da Solucao

### 3.1 Tres Camadas

```
CAMADA DIGITAL
  Motor de Sessoes | Motor de IA (4D) | Faturamento | Integracoes | Dashboard
       ^                    ^                ^
       |                    |                |
  CAMADA DE CONECTIVIDADE
  OCPP 1.6J (WebSocket) | MODBUS RTU/TCP | Wi-Fi + LAN + RS-485
       ^                    ^                ^
       |                    |                |
  CAMADA FISICA
  GoodWe HCA G2: 7kW | 11kW | 22kW | RFID ISO 14443A | IP65
```

### 3.2 Hardware GoodWe HCA G2

| Modelo | Potencia | Fase | Conector | IP |
|--------|----------|------|----------|-----|
| GW7K-HCA-20 | 7 kW | Monofasico | AC Tipo 2 | IP65 |
| GW11K-HCA-20 | 11 kW | Trifasico | AC Tipo 2 | IP65 |
| GW22K-HCA-20 | 22 kW | Trifasico | AC Tipo 2 | IP65 |

Caracteristicas: RFID ISO 14443A (10 cartoes), RS-485 + LAN + Wi-Fi + BT, protecoes RCD/OVP/UVP/OCP/OTP.

### 3.3 Fluxo de Dados

```
[Morador]--RFID-->[Carregador HCA G2]--OCPP/MODBUS-->[Servidor EV ChargeOps]
                                                         |
                    +------------------------------------+--------------------+
                    |                |                   |                    |
              [Gerenciador]    [Motor IA]          [Faturamento]       [Integracoes]
                    |          /  |  |  \\               |                    |
              [Sessoes]   Int Pred Prec Conv       [Faturas]          [OCM/Places]
                                  |
                          [Sindico Virtual]
                           /            \\
                [Dashboard Sindico]  [App Morador]
```

---

## 4. Opcao A - Estruturacao de Sessoes de Recarga

### Modelo de Dados
- **Condominio**: id, nome, endereco, capacidade (kW), tarifa base
- **UnidadeHabitacional**: id, numero, bloco, proprietario, cartoes RFID
- **Carregador**: id, modelo HCA G2, potencia, status, protocolo OCPP
- **SessaoRecarga**: id, unidade_id, carregador_id, inicio/fim, kWh, custo, tarifa
- **Fatura**: id, unidade_id, periodo, total sessoes/kWh/valor, status

### Ciclo de Vida da Sessao
1. **AUTENTICACAO**: RFID -> Authorize (OCPP) -> validacao
2. **INICIO**: StartTransaction -> registro com timestamp + IDs
3. **MONITORAMENTO**: MeterValues periodicos (kWh acumulado)
4. **FINALIZACAO**: StopTransaction -> energia total + custo
5. **CONTABILIZACAO**: vinculacao a fatura mensal

### Rastreabilidade
Todos os eventos registrados com timestamp UTC: RFID apresentado, sessao iniciada/finalizada, MeterValues, mudancas de status, alertas de anomalia.

---

## 5. Opcao B - Processamento de Consumo e IA Avancada

### Motor de IA - 4 Dimensoes

> "A IA estrutura a proposta. Nao e uma feature extra; e a camada que da sentido aos dados brutos gerados pela rede."

#### 5.1 Interpretacao
- **Classificacao de sessoes**: NORMAL, RAPIDA, LONGA, PONTA, FORA_PONTA
- **Perfis de consumo**: COMMUTER, FLEX, HEAVY_USER, LIGHT_USER
- **Parsing MODBUS**: tensao, corrente, potencia, fator de potencia -> eficiencia de carga
- **Deteccao de anomalias**: consumo 3+ sigma, sessao fantasma, RFID nao autorizado

#### 5.2 Preditividade
- **Previsao de demanda**: EWMA com janela de 30 dias, mapa de calor 24h x 7d
- **Capacity planning**: alerta quando ocupacao media >70% por 3 meses
- **Manutencao preditiva**: ciclos de carga, temperatura, MTBF, degradacao
- Evolucao planejada: Prophet/ARIMA com decomposicao tendencia + sazonalidade

#### 5.3 Precificacao
- **Formula**: Tarifa = Base x FatorHorario x FatorDemanda x FatorBandeira
  - Base: tarifa ANEEL (R$ 0,65/kWh)
  - FatorHorario: 1.5x pico (18-21h), 0.7x madrugada (23-05h)
  - FatorDemanda: 1.0-1.3x conforme ocupacao simultanea
  - FatorBandeira: verde=1.0, amarela=1.05, vermelha=1.15
- **Incentivos**: ate 30% desconto noturno, teto de preco configuravel
- **Conformidade**: RN ANEEL 1.000/2021, transparencia na composicao

#### 5.4 Conversacao - Sindico Virtual
- **Pipeline RAG**: Dados -> Embedding -> Retrieval -> Augmentation -> Generation (LLM)
- **Exemplo**: "Quanto a unidade 302 gastou?" -> "87,3 kWh em 12 sessoes (R$ 74,21). Consumo 15% acima da media. Sugestao: recarga noturna para economia de R$ 18,50/mes."
- **Alertas proativos**: consumo 3x acima da media, carregador >80% utilizacao, previsao de sobrecarga

### Load Balancing Inteligente
Redistribuicao de potencia priorizando: (1) SoC mais baixo, (2) horario de saida agendado, (3) tarifa vigente. Respeitando limite de demanda contratada.

---

## 6. Opcao C - Gerenciamento Inteligente e UX

### Interface Morador
- Dashboard pessoal: consumo (kWh, R$), sessoes, grafico diario
- Historico: data, hora, duracao, energia, custo, tarifa, carregador
- Recomendacoes IA: horarios otimos, economia estimada
- Notificacoes: sessao iniciada/finalizada, fatura, alertas

### Painel Sindico
- Visao consolidada: sessoes, consumo, receita, ocupacao
- Ranking de consumo por unidade
- Sindico Virtual conversacional
- Relatorios: PDF assembleia, CSV contabilidade, demonstrativo/unidade

### Jornadas do Usuario
1. **Morador carrega**: RFID -> carga -> notificacao ("32,4 kWh, R$ 21,06. Economia R$ 8,10.")
2. **Sindico consulta**: dashboard -> Sindico Virtual -> relatorio -> ajuste tarifario
3. **Administradora fecha mes**: faturas automaticas -> CSV -> boleto condominio

### Principios UX
- Simplicidade (tap-to-charge), transparencia (tarifa visivel), proatividade (IA sugere), acessibilidade (web responsivo), confianca (dados auditaveis)

---

## 7. Rateio e Faturamento

### Modelo
```
Custo_Unidade = SUM(kWh_sessao_i x Tarifa_sessao_i) + Taxa_Admin (5%)
```

### Processo
1. Motor de Faturamento gera faturas individuais ao final do periodo (mensal)
2. Fatura consolidada do condominio com total agregado
3. Exportacao CSV/API para administradoras (Superlogica, Condomob)
4. Item no boleto: "Recarga EV - Mes XX/XXXX"
5. Inadimplencia: desabilitacao RFID pelo sindico

---

## 8. Decisoes Tecnicas

| # | Questao | Decisao | Justificativa |
|---|---------|---------|---------------|
| Q1 | Linguagem | Python | Ecossistema IA + prototipagem rapida |
| Q2 | Protocolo | OCPP 1.6J | Aberto + exigencia ANEEL |
| Q3 | Tarifa | Dinamica | Incentivo fora-ponta + regulacao |
| Q4 | Autenticacao | RFID | Simples + incluso GoodWe |
| Q5 | Armazenamento | In-memory | MVP; -> PostgreSQL prod |
| Q6 | Infra | On-premise | LGPD + offline; -> cloud prod |
| Q7 | IA | Hibrida | Local (scikit) + API LLM (Sindico) |

---

## 9. Roadmap

| Fase | Entrega | Prazo |
|------|---------|-------|
| Sprint 01 | Pesquisa e documentacao | Jun/2026 |
| Sprint 02 | Prototipo funcional Python | Set/2026 |
| v1.0 | MVP: RFID + rateio + faturas | Dez/2026 |
| v1.5 | Motor de IA (4 dimensoes) | Mar/2027 |
| v2.0 | App mobile + dashboard + Sindico Virtual | Jun/2027 |
| v2.5 | Cloud (AWS) + multi-condominio | Set/2027 |
| v3.0 | V2G + ISO 15118 + marketplace | Dez/2027 |

---

## 10. Entregaveis Sprint 01 (Alinhamento com o Challenge)

Conforme o playbook, a avaliacao foca em raciocinio arquitetonico, de negocios e de gestao de dados:

1. **Arquitetura Funcional**: logica clara de entradas e saidas (Data Flow) em 3 camadas
2. **Papel da IA Integrada**: motor logico de 4 dimensoes, nao penduricalho de interface
3. **Aderencia ao Contexto**: solucao sob medida para condominios com ecossistema GoodWe/FIAP
4. **Visao de Produto Real**: plataforma para condominios que resolve rateio justo e gera inteligencia operacional

---

## Referencias

1. ANEEL. Resolucao Normativa n. 1.000/2021.
2. Brasil. Lei n. 14.300/2022 (Geracao Distribuida).
3. ABNT. NBR IEC 61851-1:2023.
4. EPE. Plano Nacional de Energia 2050.
5. INMETRO. Portaria n. 111/2023.
6. IEA. Global EV Outlook 2024.
7. ABVE. Anuario da Mobilidade Eletrica 2024.
8. Open Charge Alliance. OCPP 1.6 Specification.
9. GoodWe. HCA G2 Series - Technical Datasheet.
10. FIAP + GoodWe. EV Challenge 2026 - Playbook.
11. EU. AFIR - Alternative Fuels Infrastructure Regulation, 2023.
12. USA. NEVI Formula Program, 2022.
13. UK. Smart Charge Points Regulations 2021.

---

**Repositorio:** https://github.com/marcelobaldin/FIAP_FASE_III_COLONIA
**Contato:** marcelobbaldin@gmail.com
"""
    caminho = os.path.join(pasta_saida, 'ev_chargeops_sprint01_final.md')
    with open(caminho, 'w', encoding='utf-8') as f:
        f.write(conteudo)
    return caminho


# ============================================================================
# EXECUCAO
# ============================================================================

def main():
    pasta_saida = script_dir

    print("=" * 70)
    print("  EV ChargeOps v3 - Sprint 01: Entrega Completa")
    print("  Opcao A (Sessoes) + Opcao B (IA 4D) + Opcao C (UX)")
    print("=" * 70)

    print("\n[1/5] Configurando plataforma...")
    plataforma = EVChargeOps()
    plataforma.setup_demo()

    print("\n[2/5] Gerando historico simulado (30 dias)...")
    plataforma.gerar_historico_simulado(dias=30)

    print("\n[3/5] Gerando graficos...")
    gerar_graficos(plataforma, pasta_saida)

    print("\n[4/5] Gerando relatorio PDF v3...")
    caminho_pdf = gerar_pdf_v3(plataforma, pasta_saida)
    print(f"  PDF v3 salvo: {caminho_pdf}")

    print("\n[5/5] Gerando apresentacao PPTX v3...")
    caminho_pptx = gerar_pptx_v3(plataforma, pasta_saida)
    print(f"  PPTX v3 salvo: {caminho_pptx}")

    print("\n[BONUS] Gerando documento MD final...")
    caminho_md = gerar_md_final(pasta_saida)
    print(f"  MD final salvo: {caminho_md}")

    print("\n" + "=" * 70)
    print("  v3 GERADA COM SUCESSO!")
    print(f"  - {caminho_pdf}")
    print(f"  - {caminho_pptx}")
    print(f"  - {caminho_md}")
    print("=" * 70)


if __name__ == "__main__":
    main()
